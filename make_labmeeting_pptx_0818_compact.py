#!/usr/bin/env python3
"""2026-08-18 랩미팅 PPT 요약본 — 핵심 진행상황 · 사용 모델 · 진행 계획 · 사진/영상만.

  python3 make_labmeeting_pptx_0818_compact.py
  → /home/kim/tx90/대화록 및 PPT/랩미팅_20260818_요약.pptx

원본(make_labmeeting_pptx_0818.py, 17장)에서 글 위주 슬라이드를 걷어낸 7장 구성.
v5/v6 RViz 영상은 자리(placeholder)만 — 발표 전 직접 삽입.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

NAVY = RGBColor(0x2F, 0x3D, 0x9E)
INK = RGBColor(0x15, 0x18, 0x1D)
GREY = RGBColor(0x62, 0x6B, 0x78)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
RED = RGBColor(0xA9, 0x33, 0x1D)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
BGSOFT = RGBColor(0xED, 0xEF, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LILAC = RGBColor(0xB9, 0xC2, 0xF0)
AMBER = RGBColor(0x9A, 0x6A, 0x00)

FONT = "맑은 고딕"
MONO = "Consolas"
SW, SH = Cm(33.867), Cm(19.05)
BASE = "/home/kim/tx90"
OUTDIR = os.path.join(BASE, "대화록 및 PPT")
WSDIR = os.path.join(BASE, "omx_workspace")
IMG13 = os.path.join(OUTDIR, "img_0813")
IMG14 = os.path.join(OUTDIR, "img_0814")

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size=15, bold=False, color=INK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide():
    return prs.slides.add_slide(BLANK)


def bar(s, title, step=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Cm(2.0))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    tf = r.text_frame
    tf.margin_left = Cm(0.9); tf.margin_top = Cm(0.26)
    p = tf.paragraphs[0]
    if step:
        run = p.add_run(); run.text = step + "  "
        _set(run, 15, True, LILAC)
    run = p.add_run(); run.text = title
    _set(run, 21, True, WHITE)
    return r


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullets(tf, items, size=15):
    first = True
    for it in items:
        lv, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lv
        p.space_after = Pt(opt.get("after", 6))
        marks = {0: "▪  ", 1: "–  ", 2: "·  "}
        run = p.add_run()
        run.text = ("" if opt.get("nomark") else marks.get(lv, "")) + txt
        _set(run, opt.get("size", size - lv), opt.get("bold", False),
             opt.get("color", INK), MONO if opt.get("mono") else FONT)


def table(s, x, y, w, rows, widths=None, size=12.5, header=True, rh=Cm(0.8)):
    shp = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, rh * len(rows))
    t = shp.table
    if widths:
        total = sum(widths)
        for i, cw in enumerate(widths):
            t.columns[i].width = int(w * cw / total)
    for ri, row in enumerate(rows):
        t.rows[ri].height = rh
        for ci, cell in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Cm(0.16); c.margin_right = Cm(0.16)
            c.margin_top = Cm(0.05); c.margin_bottom = Cm(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt = cell if isinstance(cell, str) else cell[0]
            opt = {} if isinstance(cell, str) else cell[1]
            c.fill.solid()
            c.fill.fore_color.rgb = (NAVY if (header and ri == 0)
                                     else opt.get("bg", WHITE if ri % 2 else BGSOFT))
            p = c.text_frame.paragraphs[0]
            run = p.add_run(); run.text = txt
            _set(run, opt.get("size", size),
                 opt.get("bold", header and ri == 0),
                 WHITE if (header and ri == 0) else opt.get("color", INK),
                 MONO if opt.get("mono") else FONT)
    return t


def foot(s, txt):
    tf = textbox(s, Cm(0.9), SH - Cm(0.95), SW - Cm(1.8), Cm(0.7))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, 10, False, GREY)


def flowbox(s, x, y, w, h, txt, fill=NAVY, fg=WHITE, size=12.5, bold=True):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = LINE; b.line.width = Pt(0.75)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.12); tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.08); tf.margin_bottom = Cm(0.08)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = line
        _set(run, size if i == 0 else size - 2, bold if i == 0 else False, fg)
    return b


def arrow(s, x, y, w=Cm(0.7)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Cm(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = GREY
    a.line.fill.background()
    return a


def videobox(s, x, y, w, h, label, note="발표 전 이 자리에 영상 삽입"):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = BGSOFT
    b.line.color.rgb = NAVY; b.line.width = Pt(1.5)
    try:
        from pptx.enum.line import MSO_LINE
        b.line.dash_style = MSO_LINE.DASH
    except Exception:
        pass
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "▶  " + label
    _set(run, 17, True, NAVY)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = note
    _set(run, 12, False, GREY)
    return b


def caption(s, x, y, w, txt, size=11):
    tf = textbox(s, x, y, w, Cm(0.9))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, size, False, GREY)


# ═══════════ 1. 표지 ═══════════
s = slide()
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
r.fill.solid(); r.fill.fore_color.rgb = NAVY; r.line.fill.background()
tf = textbox(s, Cm(2.2), Cm(4.6), Cm(29.5), Cm(8))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "OMX 시연의 TX2-90 전이학습"
_set(run, 36, True, WHITE)
p = tf.add_paragraph(); p.space_before = Pt(8)
run = p.add_run()
run.text = "EE pose 정책(v5)에서 joint 정책(v6)까지 — 시뮬레이션 검증 완료"
_set(run, 20, False, LILAC)
p = tf.add_paragraph(); p.space_before = Pt(30)
run = p.add_run()
run.text = "좌표 이식 residual 0.02 mm   |   v5 자세오차 1.8°   |   v6 관절오차 ~1.1°   |   실행 시 IK 호출 0회"
_set(run, 15, True, WHITE)
tf2 = textbox(s, Cm(2.2), SH - Cm(1.7), Cm(24), Cm(0.9))
run = tf2.paragraphs[0].add_run()
run.text = "랩미팅 · 2026-08-18 (화) · 이채원 · 한국생산기술연구원 현장실습"
_set(run, 13, False, LILAC)

# ═══════════ 2. 핵심 진행 상황 ═══════════
s = slide()
bar(s, "핵심 진행 상황 — 변환 파이프라인 완성 · 모델 2종 학습 · 시뮬 검증 완료", "01 진행")
rows = [
    ["날짜", "한 일", "핵심 결과"],
    ["8/12", "OMX 데이터 분석 → 4코너 실측 → Umeyama 변환 → TX90 EE 데이터셋",
     "residual 0.02 mm — 좌표 이식 완료"],
    ["8/13", "오일러 wrap 결함 수정(v5) → ACT 재학습 → RViz 재생",
     ("자세오차 평균 28.9° → 1.8°", {"bold": True})],
    ["8/14", "EE→joint 오프라인 IK 변환(v6) → ACT 재학습 → IK 없는 RViz 재생",
     ("관절오차 ~1.1° · 실행 IK 0회", {"bold": True})],
]
table(s, Cm(0.9), Cm(2.7), Cm(32.1), rows, widths=[1.6, 10.4, 6.4], size=13.5, rh=Cm(1.75))
y1 = Cm(10.6); h = Cm(2.5); w = Cm(4.85); gap = Cm(0.75)
steps = [
    ("OMX 시연\n관절 기록 162ep", GREEN),
    ("OMX EE pose\nFK 계산", GREEN),
    ("TX90 EE pose\nUmeyama 이식 (8/12)", GREEN),
    ("데이터셋 v5\nwrap 수정 (8/13)", GREEN),
    ("데이터셋 v6\n오프라인 IK (8/14)", NAVY),
    ("v7 (예정)\n실기 수집·이어학습", GREY),
]
x = Cm(0.9)
for i, (txt, c) in enumerate(steps):
    flowbox(s, x, y1, w, h, txt, fill=c, size=12)
    if i < len(steps) - 1:
        arrow(s, x + w + Cm(0.03), y1 + h / 2 - Cm(0.25), Cm(0.66))
    x += w + gap
tf = textbox(s, Cm(0.9), Cm(13.7), Cm(32), Cm(3.6))
bullets(tf, [
    (0, "159개 시연 전체가 TX90 관절 궤적으로 번역 완료 — 시뮬레이션(RViz)에서 할 수 있는 검증은 끝", {"bold": True, "size": 16}),
    (1, "남은 갭: 카메라 영상은 끝까지 OMX 장면 (시각 도메인 갭 → v7 에서 해소 예정)", {"size": 14, "color": RED}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 3. 사용 모델 ═══════════
s = slide()
bar(s, "사용 모델 — ACT (구조 동일), v5 는 EE pose 출력 · v6 는 joint 출력", "02 모델")
tf = textbox(s, Cm(0.9), Cm(2.3), Cm(32), Cm(0.8))
bullets(tf, [(0, "모델 v5 — EE pose 를 읽고 EE pose 를 낸다 (실행에 IK 필요)", {"bold": True, "size": 14.5, "color": NAVY})])
flowbox(s, Cm(0.9), Cm(3.15), Cm(6.8), Cm(2.5), "입력\n카메라 2대 480×640\n+ state 7D [x,y,z,rx,ry,rz,grip]",
        fill=BGSOFT, fg=INK, size=11.5)
flowbox(s, Cm(8.5), Cm(3.15), Cm(8.6), Cm(2.5), "ACT (공통 구조)\nResNet18×2 + Transformer\n디코더 쿼리 100개 = 청크", size=11.5)
flowbox(s, Cm(17.9), Cm(3.15), Cm(15.0), Cm(2.5),
        "출력: EE pose 청크 100×7\n→ 실행하려면 런타임 IK 번역 필요\n(느림·특이점·손목 공회전)",
        fill=RED, size=11.5)
arrow(s, Cm(7.72), Cm(4.15)); arrow(s, Cm(17.12), Cm(4.15))
tf = textbox(s, Cm(0.9), Cm(6.1), Cm(32), Cm(0.8))
bullets(tf, [(0, "모델 v6 — joint 를 읽고 joint 를 낸다 (IK 를 데이터 준비 시점으로 이동)", {"bold": True, "size": 14.5, "color": GREEN})])
flowbox(s, Cm(0.9), Cm(6.95), Cm(6.8), Cm(2.5), "입력\n카메라 2대 (완전 동일)\n+ state 7D [j1…j6, grip]",
        fill=BGSOFT, fg=INK, size=11.5)
flowbox(s, Cm(8.5), Cm(6.95), Cm(8.6), Cm(2.5), "ACT (완전 동일 구조)\n가중치만 v6 데이터로 재학습\n(100k 스텝 · 2.9시간)", size=11.5)
flowbox(s, Cm(17.9), Cm(6.95), Cm(15.0), Cm(2.5),
        "출력: joint 청크 100×7\n→ 그대로 로봇 명령 (IK 호출 0회)",
        fill=GREEN, size=11.5)
arrow(s, Cm(7.72), Cm(7.95)); arrow(s, Cm(17.12), Cm(7.95))
rows = [
    ["항목", "모델 v5 (EE)", "모델 v6 (joint)"],
    ["학습", "ACT 100k 스텝 · ~3시간 · loss 0.061", "ACT 100k 스텝 · 2.9시간 · loss 0.052"],
    ["정확도 (open-loop)", "자세 평균 1.5~2.2° · 위치 6~14 mm", ("관절 평균 0.93~1.55° · 중앙값 0.4~0.6°", {"bold": True})],
    ["그리퍼 일치율", "95~99.5%", "96.8~99.2%"],
    ["RViz 재생", "MoveIt2 IK 번역 후 재생 (52/52)", ("예측 joint 직접 발행 — IK 0회", {"bold": True, "color": GREEN})],
]
table(s, Cm(0.9), Cm(10.2), Cm(32.1), rows, widths=[4.6, 8.6, 8.6], size=12.5, rh=Cm(1.35))
foot(s, "데이터: 159 ep · 76,345 프레임 (불량 4개 제외) · /root/train_tx90_act_v5 · v6_joint")

# ═══════════ 4. 진행 계획 ═══════════
s = slide()
bar(s, "진행 계획 — Isaac Sim 은 바로 시작, Stäubli 확인은 사람 몫", "03 계획")
rows = [
    ["단계", "상태", "비고"],
    ["OMX 분석 → 좌표 이식 → v5 학습·RViz 검증", ("완료", {"color": GREEN, "bold": True}), "8/12~13"],
    ["v6 joint 데이터셋 · 모델 · 무IK 재생", ("완료", {"color": GREEN, "bold": True}), "8/14 — v5(EE)·v6(joint) 두 노선 확보"],
    ["Stäubli 계정 · velocity 확장 견적 확인", ("진행 필요", {"color": RED, "bold": True}), "반응형 closed-loop 의 상한을 결정하는 관문"],
    ["Isaac Sim: URDF 임포트 → v6 물리 재생", ("이번 주", {"color": RED, "bold": True}), "\"실제로 집히는가\" 첫 증명"],
    ["Isaac Sim: closed-loop + 실행기 완성", ("다음", {"color": AMBER}), "재계획형 → 반응형 실험 (실기 없이 완성 가능)"],
    ["실기: VAL3 재생 → 카메라 설치 → v7 전이학습", ("예정", {"color": GREY}), "외부 제어 개통 후 · 30~50 ep 수집"],
    ["실기 closed-loop 데모", ("목표", {"color": GREY}), "\"주사위를 옮기면 따라온다\""],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[9.0, 2.4, 7.6], size=12.5, rh=Cm(1.35))
tf = textbox(s, Cm(0.9), Cm(13.6), Cm(32), Cm(4.4))
bullets(tf, [
    (0, "논의드리고 싶은 것", {"bold": True, "size": 16}),
    (1, "① 작업대 재배치 승인 — 몸쪽 20 cm + 왼쪽 15 cm (159개 전수 스윕으로 특이점 회피 확인, v6 데이터에 반영됨)", {"size": 14}),
    (1, "② Stäubli CS9 계정·velocity 확장 견적 확인 (40분짜리 확인 작업)", {"size": 14}),
    (1, "③ 실기 수집(v7) 시점 조율 — 카메라 2대 · 그리퍼 조달 포함", {"size": 14}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 5. 사진 ① OMX 작업공간 분석 ═══════════
s = slide()
bar(s, "사진 ① — OMX 시연 전체 궤적 분석 (좌표 이식의 근거)", "04 사진")
s.shapes.add_picture(os.path.join(WSDIR, "OMX_작업공간 추출 위한 시각화.png"),
                     Cm(0.9), Cm(2.5), width=Cm(32.0))
caption(s, Cm(0.9), Cm(10.2), Cm(32),
        "왼쪽부터: XY 평면(빨강 = 작업 사각형·grasp/release 지점) · XZ 측면 · 높이 분포 · 그리퍼 개폐 분포")
s.shapes.add_picture(os.path.join(WSDIR, "all_episodes_summary.png"),
                     Cm(4.0), Cm(11.2), height=Cm(6.6))
caption(s, Cm(16.5), Cm(13.5), Cm(15.5),
        "162개 에피소드 전수 요약 — 깨끗 149 · 재시도 13(유지) · 불량 4(제외) → 159 ep 사용")
foot(s, "생성: extract_omx_workspace.py · summarize_all_episodes.py")

# ═══════════ 6. 사진 ② v5·v6 데이터 분포 ═══════════
s = slide()
bar(s, "사진 ② — 같은 시연, 두 가지 표현 (위 v5 EE · 아래 v6 joint)", "05 사진")
s.shapes.add_picture(os.path.join(WSDIR, "tx90_v5_workspace.png"),
                     Cm(0.9), Cm(2.35), width=Cm(32.0))
s.shapes.add_picture(os.path.join(WSDIR, "tx90_v6_workspace.png"),
                     Cm(0.9), Cm(9.85), width=Cm(32.0))
caption(s, Cm(0.9), Cm(17.35), Cm(32),
        "grasp(빨강)/release(파랑) 패턴과 그리퍼 쌍봉은 동일, 표현만 다름 — v6 는 작업대 오프셋(−20 cm, +15 cm) 반영", 10.5)
foot(s, "생성: plot_v5_v6_workspace.py")

# ═══════════ 7. 영상 — v5 · v6 RViz ═══════════
s = slide()
bar(s, "영상 — v5 vs v6 RViz 재생 (비교 포인트: 4번 손목 축 공회전)", "06 영상")
videobox(s, Cm(0.9), Cm(2.6), Cm(15.8), Cm(8.2), "v5 RViz 시뮬레이션 영상")
videobox(s, Cm(17.2), Cm(2.6), Cm(15.8), Cm(8.2), "v6 RViz 시뮬레이션 영상")
s.shapes.add_picture(os.path.join(IMG13, "replay_t1.png"), Cm(0.9), Cm(11.2), height=Cm(6.4))
s.shapes.add_picture(os.path.join(IMG14, "v6_f2.png"), Cm(17.2), Cm(11.2), height=Cm(6.4))
caption(s, Cm(0.9), Cm(17.65), Cm(15.8), "v5 정지화면 — MoveIt2 IK 번역 후 재생 (손목 공회전 관찰됨)")
caption(s, Cm(17.2), Cm(17.65), Cm(15.8), "v6 정지화면 — 무IK 재생, 작업대 오프셋 반영 위치 (손목 공회전 없음)")
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18 · 이채원")

out = os.path.join(OUTDIR, "랩미팅_20260818_요약.pptx")
prs.save(out)
print("저장:", out)
