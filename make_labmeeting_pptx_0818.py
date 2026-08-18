#!/usr/bin/env python3
"""2026-08-18 랩미팅 PPT — OMX→TX90 전이 (8/12~8/14, Umeyama 변환 이후 전체).

  python3 make_labmeeting_pptx_0818.py
  → /home/kim/tx90/대화록 및 PPT/랩미팅_20260818.pptx

구성 원칙: 슬라이드당 메시지 하나 · 진행상황/다음단계 중심 · v5→v6 전환 상세 설명.
v5/v6 RViz 시뮬레이션 영상은 자리(placeholder)만 만들어 둠 — 발표 전 직접 삽입.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

NAVY = RGBColor(0x2F, 0x3D, 0x9E)
INK = RGBColor(0x15, 0x18, 0x1D)
GREY = RGBColor(0x62, 0x6B, 0x78)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
RED = RGBColor(0xA9, 0x33, 0x1D)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
BGSOFT = RGBColor(0xED, 0xEF, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LILAC = RGBColor(0xB9, 0xC2, 0xF0)
AMBER = RGBColor(0x9A, 0x6A, 0x00)

FONT = "맑은 고딕"
MONO = "Consolas"
SW, SH = Cm(33.867), Cm(19.05)
BASE = "/home/kim/tx90"
OUTDIR = os.path.join(BASE, "대화록 및 PPT")
WSDIR = os.path.join(BASE, "omx_workspace")
IMG13 = os.path.join(OUTDIR, "img_0813")
IMG14 = os.path.join(OUTDIR, "img_0814")

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size=15, bold=False, color=INK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide():
    return prs.slides.add_slide(BLANK)


def bar(s, title, step=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Cm(2.0))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    tf = r.text_frame
    tf.margin_left = Cm(0.9); tf.margin_top = Cm(0.26)
    p = tf.paragraphs[0]
    if step:
        run = p.add_run(); run.text = step + "  "
        _set(run, 15, True, LILAC)
    run = p.add_run(); run.text = title
    _set(run, 21, True, WHITE)
    return r


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullets(tf, items, size=15):
    first = True
    for it in items:
        lv, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lv
        p.space_after = Pt(opt.get("after", 6))
        marks = {0: "▪  ", 1: "–  ", 2: "·  "}
        run = p.add_run()
        run.text = ("" if opt.get("nomark") else marks.get(lv, "")) + txt
        _set(run, opt.get("size", size - lv), opt.get("bold", False),
             opt.get("color", INK), MONO if opt.get("mono") else FONT)


def table(s, x, y, w, rows, widths=None, size=12.5, header=True, rh=Cm(0.8)):
    shp = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, rh * len(rows))
    t = shp.table
    if widths:
        total = sum(widths)
        for i, cw in enumerate(widths):
            t.columns[i].width = int(w * cw / total)
    for ri, row in enumerate(rows):
        t.rows[ri].height = rh
        for ci, cell in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Cm(0.16); c.margin_right = Cm(0.16)
            c.margin_top = Cm(0.05); c.margin_bottom = Cm(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt = cell if isinstance(cell, str) else cell[0]
            opt = {} if isinstance(cell, str) else cell[1]
            c.fill.solid()
            c.fill.fore_color.rgb = (NAVY if (header and ri == 0)
                                     else opt.get("bg", WHITE if ri % 2 else BGSOFT))
            p = c.text_frame.paragraphs[0]
            run = p.add_run(); run.text = txt
            _set(run, opt.get("size", size),
                 opt.get("bold", header and ri == 0),
                 WHITE if (header and ri == 0) else opt.get("color", INK),
                 MONO if opt.get("mono") else FONT)
    return t


def foot(s, txt):
    tf = textbox(s, Cm(0.9), SH - Cm(0.95), SW - Cm(1.8), Cm(0.7))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, 10, False, GREY)


def flowbox(s, x, y, w, h, txt, fill=NAVY, fg=WHITE, size=12.5, bold=True):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = LINE; b.line.width = Pt(0.75)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.12); tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.08); tf.margin_bottom = Cm(0.08)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = line
        _set(run, size if i == 0 else size - 2, bold if i == 0 else False, fg)
    return b


def arrow(s, x, y, w=Cm(0.7)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Cm(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = GREY
    a.line.fill.background()
    return a


def videobox(s, x, y, w, h, label, note="발표 전 이 자리에 영상 삽입"):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = BGSOFT
    b.line.color.rgb = NAVY; b.line.width = Pt(1.5)
    try:
        from pptx.enum.line import MSO_LINE
        b.line.dash_style = MSO_LINE.DASH
    except Exception:
        pass
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "▶  " + label
    _set(run, 17, True, NAVY)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = note
    _set(run, 12, False, GREY)
    return b


def caption(s, x, y, w, txt, size=11):
    tf = textbox(s, x, y, w, Cm(0.9))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, size, False, GREY)


# ═══════════ 1. 표지 ═══════════
s = slide()
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
r.fill.solid(); r.fill.fore_color.rgb = NAVY; r.line.fill.background()
tf = textbox(s, Cm(2.2), Cm(4.6), Cm(29.5), Cm(8))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "OMX 시연의 TX2-90 전이학습"
_set(run, 36, True, WHITE)
p = tf.add_paragraph(); p.space_before = Pt(8)
run = p.add_run()
run.text = "EE pose 정책(v5)에서 joint 정책(v6)까지 — 시뮬레이션 검증 완료"
_set(run, 20, False, LILAC)
p = tf.add_paragraph(); p.space_before = Pt(30)
run = p.add_run()
run.text = "좌표 이식 residual 0.02 mm   |   v5 자세오차 1.8°   |   v6 관절오차 ~1.1°   |   실행 시 IK 호출 0회"
_set(run, 15, True, WHITE)
tf2 = textbox(s, Cm(2.2), SH - Cm(1.7), Cm(24), Cm(0.9))
run = tf2.paragraphs[0].add_run()
run.text = "랩미팅 · 2026-08-18 (화) · 이채원 · 한국생산기술연구원 현장실습"
_set(run, 13, False, LILAC)

# ═══════════ 2. 3일 요약 ═══════════
s = slide()
bar(s, "3일 요약 — 변환 파이프라인 완성 · 모델 2종 학습 · 시뮬 검증 완료", "01 요약")
rows = [
    ["날짜", "한 일", "핵심 결과"],
    ["8/12", "OMX 데이터 분석·FK 검증 → 4코너 실측 → Umeyama 유사변환 → TX90 EE 데이터셋 조립",
     "residual 0.02 mm — 좌표 이식 완료"],
    ["8/13", "오일러 wrap 결함 발견·수정(v5) → ACT 100k 재학습 → RViz 재생 배선",
     ("자세오차 평균 28.9° → 1.8°", {"bold": True})],
    ["8/14", "EE→joint 오프라인 IK 변환(v6) → ACT 100k 재학습 → IK 없는 RViz 재생",
     ("관절오차 ~1.1° · 실행 IK 0회", {"bold": True})],
]
table(s, Cm(0.9), Cm(2.7), Cm(32.1), rows, widths=[1.6, 11.2, 5.6], size=13, rh=Cm(1.75))
tf = textbox(s, Cm(0.9), Cm(10.6), Cm(32), Cm(6.8))
bullets(tf, [
    (0, "현재 위치: 시뮬레이션(RViz)에서 할 수 있는 검증은 끝났다", {"bold": True, "size": 17}),
    (1, "159개 시연 전체가 TX90 관절 궤적으로 번역되어 있고, 두 정책(v5·v6) 모두 정확도 확인 완료", {"size": 14.5}),
    (0, "다음 무대 두 개 (병렬 진행)", {"bold": True, "size": 17, "after": 4}),
    (1, "Isaac Sim — 물리 엔진 위에서 \"실제로 집히는가\" + closed-loop 첫 가동 (우리가 바로 시작 가능)", {"size": 14.5}),
    (1, "Stäubli 외부 제어 — CS9 velocity 확장 계정·견적 확인 (사람이 뚫어야 하는 관문)", {"size": 14.5}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 3. 전체 계보 ═══════════
s = slide()
bar(s, "전체 그림 — 같은 159개 시연이 표현만 바꿔가며 내려온다", "02 지도")
y1 = Cm(2.8); h = Cm(2.5); w = Cm(4.85); gap = Cm(0.75)
steps = [
    ("OMX 시연\n관절 기록 162ep\n(OpenManipulator-X)", GREEN),
    ("OMX EE pose\nFK 계산\n(fk_omx.py)", GREEN),
    ("TX90 EE pose\nUmeyama + 툴 보정\n= 좌표 이식 (8/12)", GREEN),
    ("데이터셋 v5\nwrap 수정 → 모델 v5\nEE 출력 (8/13)", GREEN),
    ("데이터셋 v6\n오프라인 IK → 모델 v6\njoint 출력 (8/14)", NAVY),
    ("v7 (예정)\n실기 수집 + 이어학습\nTX90 카메라 장면", GREY),
]
x = Cm(0.9)
for i, (txt, c) in enumerate(steps):
    flowbox(s, x, y1, w, h, txt, fill=c, size=12)
    if i < len(steps) - 1:
        arrow(s, x + w + Cm(0.03), y1 + h / 2 - Cm(0.25), Cm(0.66))
    x += w + gap
tf = textbox(s, Cm(0.9), Cm(6.1), Cm(32), Cm(11.5))
bullets(tf, [
    (0, "읽는 법", {"bold": True, "size": 16}),
    (1, "시연의 '내용'(주사위 pick & place 동작 패턴)은 끝까지 같고, 숫자의 '언어'만 단계마다 바뀐다", {"size": 14.5}),
    (1, "OMX 관절각  →  OMX 손끝 위치  →  TX90 손끝 위치  →  TX90 관절각  순서의 번역", {"size": 14.5}),
    (0, "모델 구조는 내내 동일 (ACT) — 바뀐 것은 가르친 데이터뿐", {"bold": True, "size": 16, "after": 4}),
    (1, "v5 와 v6 는 같은 시연에서 나온 형제 — v6 는 \"모델 v5 의 출력\"이 아니라 원본 시연을 다시 번역한 것", {"size": 14.5}),
    (1, "그래서 v5 노선(EE·로봇 중립)도 폐기하지 않고 병행 유지 — 다른 로봇으로 옮길 때의 보험", {"size": 14.5}),
    (0, "카메라 영상만은 번역 불가 — 끝까지 OMX 장면 (시각 도메인 갭, v7 에서 해소 예정)", {"bold": True, "size": 15,
        "color": RED}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 4. 8/12 좌표 이식 ═══════════
s = slide()
bar(s, "기반 (8/12) — Umeyama 유사변환으로 OMX 좌표를 TX90 좌표로 이식", "03 기반")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(7.2))
bullets(tf, [
    (0, "대응점: 양쪽 작업 사각형의 4코너 — OMX 는 시연 궤적 분석으로, TX90 은 펜던트 실측으로 획득", {"size": 14.5}),
    (1, "OMX 쪽 사각형 241 × 361 mm (grasp/release 지점 기준) · TX90 쪽 동일 크기 테이프 사각형", {"size": 13.5}),
    (0, "결과: s = 0.9998 (≈1) · R ≈ I · t = [+0.500, −0.014, +0.306] m · residual 0.02 mm — 사실상 순수 평행이동",
     {"bold": True, "size": 15}),
    (0, "이식 과정에서 잡은 함정 2가지 (놓쳤으면 궤적 전체가 틀어졌음)", {"size": 14.5, "after": 4}),
    (1, "펜던트 World 원점 = base_link 가 아니라 +478 mm 위의 base 링크 → 보정 안 하면 궤적이 478 mm 내려감", {"size": 13.5}),
    (1, "OMX 와 TX90 의 툴 프레임 접근축이 90° 다름 → Ry(+90°) 보정 안 하면 그리퍼가 수평으로 누움", {"size": 13.5}),
])
s.shapes.add_picture(os.path.join(WSDIR, "OMX_작업공간 추출 위한 시각화.png"),
                     Cm(0.9), Cm(10.0), width=Cm(32.0))
caption(s, Cm(0.9), Cm(17.7), Cm(32),
        "OMX 시연 전체 궤적 분석 — 왼쪽부터: XY 평면(빨강 = 작업 사각형·grasp/release 지점) · XZ 측면 · 높이 분포 · 그리퍼 개폐 분포")
foot(s, "산출: compute_transform.py · transform.npy · 데이터셋 v4 (TX90 EE pose)")

# ═══════════ 5. 데이터 정제 → v5 ═══════════
s = slide()
bar(s, "데이터 정제 — 162 → 159 에피소드, 오일러 wrap 결함 수정 = 데이터셋 v5", "04 정제")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(17.6), Cm(15.2))
bullets(tf, [
    (0, "전수 검사로 불량 시연 4개 제외 → 159 ep · 76,345 프레임", {"bold": True, "size": 15}),
    (1, "ep0 카메라 좌우 뒤바뀜 · ep62 그리퍼 미작동 · ep111/115 과제 실패", {"size": 13}),
    (0, "치명 결함: 오일러 각의 ±π 경계 점프", {"bold": True, "size": 15, "after": 4}),
    (1, "그리퍼가 수직 아래를 보는 자세라 rx·rz 가 ±π 근처에 몰림 → 저장값이 프레임 사이에 ~2π 씩 점프", {"size": 13}),
    (1, "MSE 로 학습하는 ACT 가 불연속의 '중간값'을 예측 → 자세오차 87° (v4, 10k 스텝 실측)", {"size": 13}),
    (1, "수정: rx·rz 를 [0, 2π) 로 재매핑 — 데이터가 π 중심이라 새 경계에 닿는 프레임이 0개", {"size": 13}),
    (0, "효과 (open-loop, 정답 대비)", {"bold": True, "size": 15, "after": 4}),
    (1, "자세오차 평균 28.9° → 1.8° · 최대 177° → 7.6°", {"bold": True, "size": 14, "color": GREEN}),
    (0, "교훈: 각도 표현의 불연속은 회귀 학습을 조용히 망가뜨린다 — 데이터 분포를 먼저 봐야 잡힌다",
     {"size": 13.5, "color": GREY, "after": 4}),
])
s.shapes.add_picture(os.path.join(WSDIR, "all_episodes_summary.png"),
                     Cm(19.0), Cm(2.6), width=Cm(14.0))
caption(s, Cm(19.0), Cm(10.5), Cm(14.0),
        "162개 에피소드 전수 요약 (grasp/release 위치·높이·시점) — 깨끗 149 · 재시도 13(유지) · 불량 4(제외)")
foot(s, "산출: fix_euler_wrap.py → 데이터셋 v5 (tx90_act_pick_and_place_v5_ee)")

# ═══════════ 6. 모델 v5 결과 ═══════════
s = slide()
bar(s, "모델 v5 (EE pose 출력) — 정확도는 충분, 그러나 실행에 IK 가 필요", "05 v5")
rows = [
    ["항목", "결과"],
    ["학습", "ACT 100,000 스텝 · 약 3시간 · 최종 loss 0.061"],
    ["자세 오차 (open-loop)", ("평균 1.5~2.2° · 최대 7.6°  (wrap 수정 전 28.9°)", {"bold": True})],
    ["위치 오차", "평균 6~14 mm"],
    ["그리퍼 개폐 일치율", "95~99.5%"],
    ["RViz 재생", "MoveIt2 IK 로 관절 번역 후 재생 — 작업대 −10 cm 이동 시 52/52 성공"],
]
table(s, Cm(0.9), Cm(2.6), Cm(19.2), rows, widths=[5.4, 13.8], size=12.5, rh=Cm(1.25))
tf = textbox(s, Cm(0.9), Cm(11.0), Cm(19.2), Cm(6.6))
bullets(tf, [
    (0, "정책 자체는 시연을 거의 완벽히 재현", {"bold": True, "size": 15}),
    (0, "그러나 출력이 '손끝 위치'라서, 로봇이 알아들으려면 매번 IK 번역이 필요", {"bold": True, "size": 15}),
    (1, "이 번역 과정에서 문제가 드러남 → 다음 장", {"size": 14, "color": RED}),
])
videobox(s, Cm(21.0), Cm(2.6), Cm(12.0), Cm(6.8), "v5 RViz 시뮬레이션 영상")
s.shapes.add_picture(os.path.join(IMG13, "replay_t1.png"), Cm(21.0), Cm(9.9), height=Cm(7.2))
caption(s, Cm(27.8), Cm(9.9), Cm(5.2), "RViz 재생 정지화면 (TX2-90, 예측 궤적)")
foot(s, "산출: run_policy_tx90.py (predict / replay) · /root/train_tx90_act_v5")

# ═══════════ 7. v5 의 한계 ═══════════
s = slide()
bar(s, "v5 의 한계 — 실행할 때마다 IK 를 푸는 구조가 병목", "06 전환 이유")
flowbox(s, Cm(0.9), Cm(2.6), Cm(7.6), Cm(2.3), "모델 v5 출력\n손끝 위치·자세 (EE pose)", fill=BGSOFT, fg=INK, size=13)
flowbox(s, Cm(9.4), Cm(2.6), Cm(8.6), Cm(2.3), "IK 번역 (런타임)\n\"이 손끝 = 관절 몇 도?\"", fill=RED, size=13)
flowbox(s, Cm(18.9), Cm(2.6), Cm(7.6), Cm(2.3), "관절 명령\n로봇이 알아듣는 언어", fill=GREEN, size=13)
arrow(s, Cm(8.6), Cm(3.5)); arrow(s, Cm(18.1), Cm(3.5))
tf = textbox(s, Cm(0.9), Cm(5.6), Cm(32), Cm(11.8))
bullets(tf, [
    (0, "빨간 상자(런타임 IK)에서 실제로 겪은 문제 3가지", {"bold": True, "size": 16}),
    (1, "느리고 불확실 — 청크마다 0.1~0.5초, 실패 가능. 실시간(30Hz) 제어에 치명적", {"size": 14}),
    (1, "손목 특이점 — 이 작업 자세의 IK 해가 특이점(j5≈0°) 근처라 waypoint 가 끊김 (초기 재생 44/52)", {"size": 14}),
    (1, "손목 공회전 — IK 가 특이점 근처의 '헐거운' 자유도 안에서 표류 → j4/j6 가 태스크와 무관하게 헛돎 (영상에서 보임)",
     {"size": 14}),
    (0, "시드 체인·브랜치 거부 등으로 완화했지만, \"실행 중에 푼다\"는 구조 자체가 남긴 리스크", {"size": 14.5, "after": 4}),
    (1, "실기라면 공회전은 실제 모터 동작 — 마모 · 케이블 꼬임 · 관절 한계 정지", {"size": 14, "color": RED}),
    (0, "전환 아이디어: IK 를 실행 시점에서 데이터 준비 시점으로 옮기자", {"bold": True, "size": 17, "color": NAVY,
        "after": 4}),
    (1, "오프라인에서는 시간 제약이 없다 — 전역 탐색(DP)·재시도·전수 검증까지 하고, 실행 때는 IK 를 아예 안 푼다", {"size": 14.5}),
    (1, "이것이 v6 = joint 출력 모델 (박사님이 말씀하신 \"디코더에서 joint 출력\"의 구현)", {"bold": True, "size": 14.5}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 8. v5 vs v6 입출력 ═══════════
s = slide()
bar(s, "v5 vs v6 — 모델 구조는 동일, 입출력 벡터의 '의미'만 교체", "07 v6 입출력")
tf = textbox(s, Cm(0.9), Cm(2.3), Cm(32), Cm(0.8))
bullets(tf, [(0, "모델 v5 — EE pose 를 읽고 EE pose 를 낸다", {"bold": True, "size": 14.5, "color": NAVY})])
flowbox(s, Cm(0.9), Cm(3.15), Cm(6.8), Cm(2.5), "입력\n카메라 2대 480×640\n+ state 7D [x,y,z,rx,ry,rz,grip]",
        fill=BGSOFT, fg=INK, size=11.5)
flowbox(s, Cm(8.5), Cm(3.15), Cm(8.6), Cm(2.5), "ACT (공통 구조)\nResNet18×2 + Transformer\n디코더 쿼리 100개 = 청크", size=11.5)
flowbox(s, Cm(17.9), Cm(3.15), Cm(15.0), Cm(2.5),
        "출력: EE pose 청크 100×7\n예) [0.777, −0.125, 0.386, 3.38, 0.84, 3.10, 0.451]\n→ 실행하려면 IK 번역 필요",
        fill=RED, size=11.5)
arrow(s, Cm(7.72), Cm(4.15)); arrow(s, Cm(17.12), Cm(4.15))
tf = textbox(s, Cm(0.9), Cm(6.1), Cm(32), Cm(0.8))
bullets(tf, [(0, "모델 v6 — joint 를 읽고 joint 를 낸다", {"bold": True, "size": 14.5, "color": GREEN})])
flowbox(s, Cm(0.9), Cm(6.95), Cm(6.8), Cm(2.5), "입력\n카메라 2대 (완전 동일)\n+ state 7D [j1…j6, grip]",
        fill=BGSOFT, fg=INK, size=11.5)
flowbox(s, Cm(8.5), Cm(6.95), Cm(8.6), Cm(2.5), "ACT (완전 동일 구조)\n가중치만 v6 데이터로 재학습\n(100k 스텝 · 2.9시간)", size=11.5)
flowbox(s, Cm(17.9), Cm(6.95), Cm(15.0), Cm(2.5),
        "출력: joint 청크 100×7\n예) [0.1°, 35.9°, 115.1°, 42.2°, −24.1°, −37.5°, 0.451]\n→ 그대로 로봇 명령 (IK 불필요)",
        fill=GREEN, size=11.5)
arrow(s, Cm(7.72), Cm(7.95)); arrow(s, Cm(17.12), Cm(7.95))
tf = textbox(s, Cm(0.9), Cm(10.1), Cm(32), Cm(7.5))
bullets(tf, [
    (0, "두 예시 벡터는 같은 순간(에피소드 0, 집는 순간 t=227)의 실제 데이터 — 같은 사건의 두 언어", {"bold": True, "size": 15.5}),
    (1, "v5: \"손끝을 (0.777, −0.125, 0.386) m 에, 자세 (194°, 48°, 178°) 로, 그리퍼 닫아라\"", {"size": 14}),
    (1, "v6: \"관절 6개를 (0.1°, 35.9°, 115.1°, 42.2°, −24.1°, −37.5°) 로, 그리퍼 닫아라\"", {"size": 14}),
    (0, "입력의 state 도 같은 원리로 교체: 현재 EE pose → 현재 관절각 (로봇의 고유감각)", {"size": 14.5, "after": 4}),
    (1, "실기에서 로봇이 실제로 재는 것도 관절각(엔코더)뿐 — EE pose 는 언제나 FK 로 계산한 2차 값", {"size": 13.5, "color": GREY}),
    (0, "안 바뀐 것: 카메라 영상 · 그리퍼 값 · 청크 100 · 학습 손실 — \"어떤 모델이냐\"가 아니라 \"무엇을 가르쳤냐\"의 차이",
     {"bold": True, "size": 14.5}),
])
foot(s, "실측 값 출처: 각 데이터셋 episode_000000.parquet")

# ═══════════ 9. 데이터셋 변환 방법 ═══════════
s = slide()
bar(s, "데이터셋 v5 → v6 변환 — 오프라인 IK, 규약 5가지", "08 v6 변환")
rows = [
    ["#", "규약", "왜"],
    ["①", "작업대 오프셋 (x −0.20, y +0.15) m 적용", "159개 전체가 손목 특이점을 벗어나는 위치 (28종 스윕으로 탐색)"],
    ["②", "7프레임 이동평균 평활", "OMX 시연의 손떨림 지터 제거 — IK 를 흔드는 원인"],
    ["③", "앵커 IK — 6Hz 만 IK, 사이는 관절 선형 보간", "특이점을 관절 공간에서 무리 없이 통과 (30Hz 전 프레임 IK 는 실패율 30%)"],
    ["④", "state[t] = action[t−1] 로 재정의", "OMX 측정 노이즈 배제 · 연속성 자동 보장 (완전 추종 가정)"],
    ["⑤", "그리디 실패 시 DP 전역 탐색 폴백", "에피소드 전체를 보고 관절 배치(브랜치) 일관 선택 — 런타임 IK 는 불가능한 사치"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[0.8, 8.2, 10.6], size=12.5, rh=Cm(1.45))
tf = textbox(s, Cm(0.9), Cm(11.6), Cm(20.6), Cm(6.0))
bullets(tf, [
    (0, "결과: 159 / 159 에피소드 변환 성공", {"bold": True, "size": 16, "color": GREEN}),
    (1, "그리디 체인 158 + DP 폴백 1 (ep143)", {"size": 13.5}),
    (0, "오프라인이라 가능했던 것", {"bold": True, "size": 15, "after": 4}),
    (1, "시간 제약 없음 → 오프셋 28종 × 159 에피소드 전수 스윕, 전역 탐색, 4종 검증까지", {"size": 13.5}),
    (1, "런타임 IK 는 미래를 모르니 그 자리 최선(그리디)밖에 못 함 — 이것이 v6 방식의 본질적 이점", {"size": 13.5, "color": GREY}),
    (0, "영상·그리퍼·타임스탬프는 그대로 — 숫자 6개의 언어만 교체", {"bold": True, "size": 14.5, "after": 4}),
])
s.shapes.add_picture(os.path.join(IMG14, "cam1_grasp.png"), Cm(22.2), Cm(11.6), width=Cm(5.3))
s.shapes.add_picture(os.path.join(IMG14, "cam2_grasp.png"), Cm(27.7), Cm(11.6), width=Cm(5.3))
caption(s, Cm(22.2), Cm(15.7), Cm(10.8), "t=227 의 두 카메라 — v5·v6 에서 픽셀 단위로 동일 (OMX 장면 = 시각 도메인 갭의 실체)")
foot(s, "산출: convert_v6.py · 규약 전체는 v6 meta/conversion_meta.json 에 기록")

# ═══════════ 10. 검증 4종 + 작업대 갱신 ═══════════
s = slide()
bar(s, "데이터셋 v6 검증 4종 통과 — 작업대 위치는 (−20 cm, +15 cm) 로 갱신", "09 검증")
rows = [
    ["검증", "방법", "결과", "판정"],
    ["변환 완전성", "159 에피소드 앵커 체인 (+DP 폴백)", "159/159", ("통과", {"color": GREEN, "bold": True})],
    ["연속성", "30Hz 프레임 간 max|Δjoint| 전수", "최대 16.2° · 평균 4.9° (임계 30°)", ("통과", {"color": GREEN, "bold": True})],
    ["FK 역검증", "관절 → 손끝 왕복 계산 비교", "위치 ≤5.1 mm · 자세 ≤1.0°", ("통과", {"color": GREEN, "bold": True})],
    ["다봉성 (kNN)", "같은 입력 → 다른 관절 출력 쌍 탐색", "4,770 샘플 중 0건 — 단봉 확인", ("통과", {"color": GREEN, "bold": True})],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[3.4, 7.0, 6.4, 1.6], size=12.5, rh=Cm(1.35))
b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.9), Cm(10.1), Cm(32.1), Cm(5.2))
b.fill.solid(); b.fill.fore_color.rgb = BGSOFT
b.line.color.rgb = RED; b.line.width = Pt(1.5)
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.3); tf.margin_right = Cm(0.5)
bullets(tf, [
    (0, "⚠ 결정 갱신 — 작업대(테이프 사각형) 위치: 로봇 기준 몸쪽 20 cm + 왼쪽 15 cm 이동", {"bold": True, "size": 16, "color": RED, "nomark": True}),
    (1, "어제의 −10 cm 는 예측 궤적 1개 기준이었음 → 전체 159개 기준으로 28종 스윕: (−0.10, 0) 은 131개만 통과, (−0.20, +0.15) 에서 159개 전부 통과",
     {"size": 13.5}),
    (1, "±2 cm 여유 있는 고원 지대라 실측 오차에 강건", {"size": 13.5}),
    (1, "실기 적용 시 테이프 사각형을 이 위치로 옮겨야 함 — v6 데이터에 이미 구워진 값", {"bold": True, "size": 13.5}),
])
tf = textbox(s, Cm(0.9), Cm(15.8), Cm(32), Cm(2.0))
bullets(tf, [
    (0, "남은 주의 1건: 특이점 통과 순간 최대 16.2°/frame 스윙 잔존 → 실기는 감속 실행 + 안전 필터(|Δjoint| 한계) 전제",
     {"size": 13, "color": GREY}),
])
foot(s, "검증 리포트: v6 meta/conversion_report.json · knn_check.py")

# ═══════════ 11. 모델 v6 결과 ═══════════
s = slide()
bar(s, "모델 v6 결과 — v5 급 정확도, IK 호출 0회의 실행", "10 v6 결과")
rows = [
    ["항목", "결과"],
    ["학습", "ACT 100,000 스텝 · 2.9시간 · 최종 loss 0.052"],
    ["관절 오차 (open-loop, 3개 ep)", ("평균 0.93~1.55° · 중앙값 0.4~0.6°", {"bold": True})],
    ["그리퍼 개폐 일치율", "96.8~99.2%"],
    ["출력 연속성", "프레임 간 평균 0.32° (청크 경계 점프 1건 → 안전 필터로 처리)"],
    ["RViz 재생", ("예측 joint 직접 발행 — IK 호출 0회, 20초 전 구간 매끄러움", {"bold": True, "color": GREEN})],
]
table(s, Cm(0.9), Cm(2.6), Cm(19.2), rows, widths=[6.4, 12.8], size=12.5, rh=Cm(1.25))
tf = textbox(s, Cm(0.9), Cm(11.0), Cm(19.2), Cm(6.6))
bullets(tf, [
    (0, "영상 비교 포인트: v5 에서 계속 헛돌던 4번 손목 축이 v6 에선 돌지 않는다", {"bold": True, "size": 15}),
    (1, "특이점에서 먼 작업 위치 + 변환 규약 + 학습의 평활 — 3겹으로 원인 제거", {"size": 13.5}),
    (1, "손끝 궤적은 두 방식이 사실상 동일 — 다른 것은 손목이 헛도느냐뿐", {"size": 13.5, "color": GREY}),
    (0, "실기 함의: worst-case 지연이 상수(추론 + 뺄셈 필터), \"IK 안 풀리면?\" 원천 봉쇄", {"size": 14, "after": 4}),
])
videobox(s, Cm(21.0), Cm(2.6), Cm(12.0), Cm(6.8), "v6 RViz 시뮬레이션 영상")
s.shapes.add_picture(os.path.join(IMG14, "v6_f2.png"), Cm(21.0), Cm(9.9), height=Cm(7.2))
caption(s, Cm(27.8), Cm(9.9), Cm(5.2), "v6 무IK 재생 정지화면 — 작업대 오프셋이 반영된 위치")
foot(s, "산출: eval_v6.py · replay_v6.py · /root/train_tx90_act_v6_joint")

# ═══════════ 12. v5·v6 데이터 분포 그림 ═══════════
s = slide()
bar(s, "같은 시연, 두 가지 표현 — 데이터셋 v5(EE) · v6(joint) 전수 분석", "11 데이터")
s.shapes.add_picture(os.path.join(WSDIR, "tx90_v5_workspace.png"),
                     Cm(0.9), Cm(2.35), width=Cm(32.0))
s.shapes.add_picture(os.path.join(WSDIR, "tx90_v6_workspace.png"),
                     Cm(0.9), Cm(9.85), width=Cm(32.0))
caption(s, Cm(0.9), Cm(17.35), Cm(32),
        "위 = v5 (TX90 손끝 좌표) · 아래 = v6 (오프셋 반영 실행 위치 · 관절별 분포와 한계 여유 · 연속성) — "
        "grasp(빨강)/release(파랑) 패턴과 그리퍼 쌍봉은 동일, 표현만 다름", 10.5)
foot(s, "생성: plot_v5_v6_workspace.py")

# ═══════════ 13. 현재 위치 — open-loop vs closed-loop ═══════════
s = slide()
bar(s, "현재 위치 — 지금까지는 전부 open-loop, closed-loop 는 루프의 양 끝이 없다", "12 현재 위치")
y1 = Cm(2.7); h = Cm(2.6)
flowbox(s, Cm(0.9), y1, Cm(10.0), h, "① open-loop 재생  ✔ 완료\n녹화 관측 → 추론 → 통째 재생\n(지금의 RViz 검증)", fill=GREEN, size=13)
flowbox(s, Cm(11.9), y1, Cm(10.0), h, "② 재계획형 closed-loop\n청크마다 관측 → 실행 (stop-and-go)\nIsaac Sim → 실기", fill=NAVY, size=13)
flowbox(s, Cm(22.9), y1, Cm(10.1), h, "③ 반응형 closed-loop\n실행 중 재추론·청크 이어붙이기 (30Hz)\n= \"실시간 action chunk\"", fill=GREY, size=13)
arrow(s, Cm(10.95), y1 + Cm(1.05)); arrow(s, Cm(21.95), y1 + Cm(1.05))
tf = textbox(s, Cm(0.9), Cm(6.0), Cm(32), Cm(11.6))
bullets(tf, [
    (0, "지금의 재생이 open-loop 인 이유", {"bold": True, "size": 16}),
    (1, "관측(카메라 영상)이 녹화본 — 실행 중 주사위를 옮겨도 정책은 모른다", {"size": 14}),
    (0, "closed-loop 로 못 넘어가는 이유는 모델이 아니라 인프라", {"bold": True, "size": 16, "after": 4}),
    (1, "살아있는 관측이 없다 — TX90 에 카메라 미설치 + 학습 영상이 OMX 장면 (시각 도메인 갭)", {"size": 14}),
    (1, "명령 받을 로봇이 없다 — RViz 는 물리 없는 표시 전용, 실기는 외부 제어 미개통", {"size": 14}),
    (0, "청크 출력 자체는 이미 작동 중 — 추론 1회 = 100스텝(3.3초), 반응형 전환은 설정 변경이지 재학습이 아님",
     {"bold": True, "size": 15, "color": NAVY, "after": 4}),
    (1, "n_action_steps 축소(청크 일부만 쓰고 재추론) + 비동기 추론(실행과 겹침) + temporal ensembling(청크 이어붙임) — 전부 ACT 내장·검증된 표준 배포법",
     {"size": 13.5}),
    (1, "반응 지연 ≈ 재추론 주기 + 추론 시간 ≈ 0.1초 이내 — \"주사위를 옮기면 다음 청크가 따라온다\"의 실체", {"size": 13.5}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 14. 실시간 청크 준비물 ═══════════
s = slide()
bar(s, "실시간 청크(반응형)까지 준비물 — 관문은 하나, 나머지는 병렬 진행", "13 준비물")
rows = [
    ["구분", "항목", "상태"],
    [("A. 로봇 명령 통로", {"bold": True}), "VAL3 드라이버 (공식 ROS, 무료) — 재계획형까지 가능", "조사 완료, 설치 대기"],
    ["", "CS9 velocity/motion 확장 — 4ms 연속 추종, 반응형의 열쇠",
     ("최대 관문 — 계정·견적 확인 필요 (사람 몫)", {"color": RED, "bold": True})],
    [("B. 관측 (정책의 눈)", {"bold": True}), "TX90 에 카메라 2대 고정 마운트 + 그리퍼 장착", "미착수"],
    ["", "실기 시연 30~50ep 수집 → 전이학습 v7 (시각 갭 + state 실측화)", "외부 제어 개통 후"],
    [("C. 실행기 (코드)", {"bold": True}), "비동기 추론 러너 + n_action_steps/앙상블 설정 + 안전 필터",
     ("Isaac Sim 에서 실기 없이 전부 완성 가능", {"color": GREEN})],
    [("D. 배치", {"bold": True}), "테이프 사각형을 (−20 cm, +15 cm) 위치로 이동", "v6 데이터 좌표 확정됨"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[3.8, 10.6, 5.6], size=12.5, rh=Cm(1.5))
tf = textbox(s, Cm(0.9), Cm(13.6), Cm(32), Cm(4.2))
bullets(tf, [
    (0, "핵심: 사람이 뚫어야 하는 관문은 A 의 velocity 확장 하나 — B·C·D 는 그와 무관하게 병렬 진행", {"bold": True, "size": 15.5}),
    (1, "v7 전이학습: 사람이 조종하는 게 아니라, 검증된 v6 궤적을 로봇이 실행하는 동안 TX90 카메라+엔코더를 녹화 → 이어학습 (30~50개면 충분 — 새로 배우는 건 화면 해석뿐)",
     {"size": 13.5}),
    (1, "velocity 확장이 불가로 판명되면: 실기는 재계획형(VAL3 만으로 가능)으로 마무리, 반응형 연구는 Isaac Sim 에서 완결 — 플랜 B",
     {"size": 13.5, "color": GREY}),
])
foot(s, "참고: staubli_val3_driver · adaptive_motion_control (FAU-FAPS) · LeRobot 비동기 스택")

# ═══════════ 15. Isaac Sim ═══════════
s = slide()
bar(s, "Isaac Sim — 실기 개통을 기다리는 동안 closed-loop 를 미리 완성하는 무대", "14 Isaac Sim")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(4.4))
bullets(tf, [
    (0, "왜 (한 문장): 실기 closed-loop 의 두 관문 — 명령 통로(A)와 살아있는 관측(B) — 이 Isaac Sim 안에서는 둘 다 공짜다",
     {"bold": True, "size": 16}),
    (1, "시뮬 로봇은 매 스텝 관절 명령을 그냥 받고, 렌더 카메라는 원하는 위치·화각으로 즉시 설치 가능", {"size": 14}),
    (1, "그리고 RViz 는 물리가 없는 애니메이션 — \"그리퍼가 주사위를 실제로 집는가\" 는 Isaac Sim 이 처음 답할 수 있다",
     {"size": 14}),
])
rows = [
    ["용도", "내용", "시점"],
    ["① 물리 검증", "v6 관절 궤적을 물리 엔진 위에서 재생 — 실제로 집히는지·충돌 없는지 첫 증명", ("이번 주 목표", {"bold": True})],
    ["② closed-loop 첫 가동", "렌더 카메라(관측) + 물리 로봇(실행)으로 루프 최초 가동 — 재계획형·반응형 실험, 실행기(C) 튜닝을 실기 전에 완료",
     "그 다음"],
    ["③ 데이터 합성 (선택)", "TX90 장면 렌더로 시연 자동 생성 + 도메인 랜덤화 — 실기 수집 부담 축소", "필요시"],
]
table(s, Cm(0.9), Cm(6.7), Cm(32.1), rows, widths=[3.6, 11.6, 2.8], size=12.5, rh=Cm(1.7))
tf = textbox(s, Cm(0.9), Cm(14.1), Cm(32), Cm(4.0))
bullets(tf, [
    (0, "전략적 의미: Isaac Sim 트랙은 Stäubli 견적 결과와 무관하게 진행된다 — 최악의 경우에도 반응형 연구는 sim 에서 완결",
     {"bold": True, "size": 15}),
    (1, "바로 시작 가능한 자산: tx2_90.urdf (검증된 URDF) + 테이블·주사위 좌표 (v6 데이터 기준) + ROS2 브리지", {"size": 14}),
    (1, "성공 판정 데모: \"실행 중에 주사위를 옮기면 다음 청크가 따라오는가\"", {"size": 14, "color": NAVY}),
])
foot(s, "NVIDIA Isaac Sim · URDF 임포트 → 장면 구성 → v6 재생 → closed-loop")

# ═══════════ 16. 로드맵 ═══════════
s = slide()
bar(s, "로드맵 — 완료 · 진행 · 예정", "15 로드맵")
rows = [
    ["단계", "상태", "비고"],
    ["OMX 분석 → 좌표 이식 → v5 학습·RViz 검증", ("완료", {"color": GREEN, "bold": True}), "8/12~13"],
    ["v6 joint 데이터셋 · 모델 · 무IK 재생", ("완료", {"color": GREEN, "bold": True}), "8/14 — v5(EE)·v6(joint) 두 노선 확보"],
    ["Stäubli 계정 · velocity 확장 견적 확인", ("진행 필요", {"color": RED, "bold": True}), "반응형의 상한을 결정하는 관문"],
    ["Isaac Sim: URDF 임포트 → v6 물리 재생", ("이번 주", {"color": RED, "bold": True}), "closed-loop 의 첫 무대"],
    ["Isaac Sim: closed-loop + 실행기 완성", ("다음", {"color": AMBER}), "재계획형 → 반응형 실험"],
    ["실기: VAL3 open-loop 재현 → 카메라 설치 → v7 전이학습", ("예정", {"color": GREY}), "외부 제어 개통 후"],
    ["실기 closed-loop 데모", ("목표", {"color": GREY}), "\"주사위를 옮기면 따라온다\""],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[9.0, 2.4, 6.6], size=12.5, rh=Cm(1.35))
tf = textbox(s, Cm(0.9), Cm(13.8), Cm(32), Cm(4.2))
bullets(tf, [
    (0, "리스크 유지 목록", {"bold": True, "size": 15}),
    (1, "시각 도메인 갭 (학습 영상이 OMX 장면) — v7 전이학습으로 해소 예정, 실기 closed-loop 의 전제", {"size": 13.5}),
    (1, "청크 경계 관절 점프 1건 실측 — 실기 전 안전 필터(|Δjoint| 한계) 구현 필수", {"size": 13.5}),
    (1, "실기 테이블 앞모서리 충돌 가능 구간(접근 통로) — 배치 시 실측 확인 or MoveIt 충돌 객체", {"size": 13.5}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18")

# ═══════════ 17. 논의 사항 ═══════════
s = slide()
bar(s, "논의드리고 싶은 것", "16 논의")
tf = textbox(s, Cm(0.9), Cm(2.8), Cm(32), Cm(14))
bullets(tf, [
    (0, "① 작업대 재배치 승인 — 테이프 사각형을 로봇 기준 몸쪽 20 cm + 왼쪽 15 cm 위치로",
     {"bold": True, "size": 17}),
    (1, "159개 시연 전체가 손목 특이점을 회피하는 위치 (전수 스윕 결과) · v6 데이터에 이미 반영됨", {"size": 14, "after": 14}),
    (0, "② Stäubli 외부 제어 확인 진행 — CS9 계정 · velocity/motion 확장 견적",
     {"bold": True, "size": 17}),
    (1, "결과에 따라 실기의 상한이 재계획형(무료 VAL3 만)이냐 반응형(확장 필요)이냐 결정 — 40분짜리 확인 작업", {"size": 14, "after": 14}),
    (0, "③ Isaac Sim 진행 우선순위 — 이번 주 물리 검증부터",
     {"bold": True, "size": 17}),
    (1, "실기 개통과 무관하게 closed-loop 실험·실행기 코드를 미리 완성해 두는 트랙", {"size": 14, "after": 14}),
    (0, "④ 실기 수집(v7) 시점 조율 — 카메라 2대 · 그리퍼 조달 포함",
     {"bold": True, "size": 17}),
    (1, "외부 제어 개통 후 반나절 수집 + 2~3시간 학습이면 시각 도메인 갭까지 해소", {"size": 14}),
])
foot(s, "OMX→TX2-90 전이 · 랩미팅 2026-08-18 · 이채원")

out = os.path.join(OUTDIR, "랩미팅_20260818.pptx")
prs.save(out)
print("저장:", out)
