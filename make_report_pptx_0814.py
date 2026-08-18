#!/usr/bin/env python3
"""2026-08-14 작업 정리 PPT (20260813.pptx 이후 → v6 joint 모델 구축 전체).

  python3 make_report_pptx_0814.py
  → /home/kim/tx90/대화록 및 PPT/20260814.pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

NAVY = RGBColor(0x2F, 0x3D, 0x9E)
INK = RGBColor(0x15, 0x18, 0x1D)
GREY = RGBColor(0x62, 0x6B, 0x78)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
RED = RGBColor(0xA9, 0x33, 0x1D)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
BGSOFT = RGBColor(0xED, 0xEF, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LILAC = RGBColor(0xB9, 0xC2, 0xF0)

FONT = "맑은 고딕"
MONO = "Consolas"
SW, SH = Cm(33.867), Cm(19.05)
BASE = "/home/kim/tx90"
OUTDIR = os.path.join(BASE, "대화록 및 PPT")
WSDIR = os.path.join(BASE, "omx_workspace")
IMG13 = os.path.join(OUTDIR, "img_0813")
IMG14 = os.path.join(OUTDIR, "img_0814")

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size=14, bold=False, color=INK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide():
    return prs.slides.add_slide(BLANK)


def bar(s, title, step=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Cm(2.0))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    tf = r.text_frame
    tf.margin_left = Cm(0.9); tf.margin_top = Cm(0.28)
    p = tf.paragraphs[0]
    if step:
        run = p.add_run(); run.text = step + "  "
        _set(run, 15, True, LILAC)
    run = p.add_run(); run.text = title
    _set(run, 21, True, WHITE)
    return r


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullets(tf, items, size=14):
    first = True
    for it in items:
        lv, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lv
        p.space_after = Pt(opt.get("after", 5))
        marks = {0: "▪  ", 1: "–  ", 2: "·  "}
        run = p.add_run()
        run.text = ("" if opt.get("nomark") else marks.get(lv, "")) + txt
        _set(run, opt.get("size", size - lv), opt.get("bold", False),
             opt.get("color", INK), MONO if opt.get("mono") else FONT)
        if opt.get("link"):
            run.hyperlink.address = opt["link"]
            run.font.color.rgb = NAVY


def table(s, x, y, w, rows, widths=None, size=11.5, header=True, rh=Cm(0.72)):
    shp = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, rh * len(rows))
    t = shp.table
    if widths:
        total = sum(widths)
        for i, cw in enumerate(widths):
            t.columns[i].width = int(w * cw / total)
    for ri, row in enumerate(rows):
        t.rows[ri].height = rh
        for ci, cell in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Cm(0.15); c.margin_right = Cm(0.15)
            c.margin_top = Cm(0.04); c.margin_bottom = Cm(0.04)
            txt = cell if isinstance(cell, str) else cell[0]
            opt = {} if isinstance(cell, str) else cell[1]
            c.fill.solid()
            c.fill.fore_color.rgb = (NAVY if (header and ri == 0)
                                     else opt.get("bg", WHITE if ri % 2 else BGSOFT))
            p = c.text_frame.paragraphs[0]
            run = p.add_run(); run.text = txt
            _set(run, opt.get("size", size),
                 opt.get("bold", header and ri == 0),
                 WHITE if (header and ri == 0) else opt.get("color", INK),
                 MONO if opt.get("mono") else FONT)
    return t


def foot(s, txt):
    tf = textbox(s, Cm(0.9), SH - Cm(1.0), SW - Cm(1.8), Cm(0.7))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, 10, False, GREY)


def flowbox(s, x, y, w, h, txt, fill=NAVY, fg=WHITE, size=12, bold=True):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = LINE; b.line.width = Pt(0.75)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.06); tf.margin_bottom = Cm(0.06)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = line
        _set(run, size if i == 0 else size - 2, bold if i == 0 else False, fg)
    return b


def arrow(s, x, y, w=Cm(0.7)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Cm(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = GREY
    a.line.fill.background()
    return a


# ═══════════ 1. 표지 ═══════════
s = slide()
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
r.fill.solid(); r.fill.fore_color.rgb = NAVY; r.line.fill.background()
tf = textbox(s, Cm(2.2), Cm(5.0), Cm(29.5), Cm(7))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "v6: joint 출력 모델 구축"
_set(run, 36, True, WHITE)
p = tf.add_paragraph(); p.space_before = Pt(6)
run = p.add_run()
run.text = "2026-08-14 작업 정리 — 오프라인 IK 변환 · 검증 4종 · 재학습 · RViz 무IK 재생"
_set(run, 19, False, LILAC)
p = tf.add_paragraph(); p.space_before = Pt(28)
run = p.add_run()
run.text = "변환 159/159   |   관절오차 평균 ~1.1°   |   실행 IK 호출 0회   |   작업대 위치 갱신 (−20cm, +15cm)"
_set(run, 15, True, WHITE)
tf2 = textbox(s, Cm(2.2), SH - Cm(1.6), Cm(20), Cm(0.8))
run = tf2.paragraphs[0].add_run(); run.text = "이채원 · 2026-08-14"
_set(run, 12, False, LILAC)

# ═══════════ 2. 파이프라인: 데이터 변환 계보 ═══════════
s = slide()
bar(s, "파이프라인 — 데이터 변환 계보 (오늘 = [4] 단계 완성)", "PIPELINE")
y1 = Cm(2.7); h = Cm(2.15); w = Cm(5.65); gap = Cm(1.0)
steps = [
    ("[0] OMX 관절 기록\n시연 162ep, 4축+그리퍼", GREEN),
    ("[1] OMX EE pose\nFK (fk_omx.py)", GREEN),
    ("[2] TX90 EE (v4)\nUmeyama+툴 보정", GREEN),
    ("[3] 데이터셋 v5\nwrap 수정 → 모델 v5", GREEN),
    ("[4] 데이터셋 v6\n오프셋+IK → 모델 v6 (오늘)", RED),
]
x = Cm(0.9)
for i, (txt, c) in enumerate(steps):
    flowbox(s, x, y1, w, h, txt, fill=c, size=12)
    if i < len(steps) - 1:
        arrow(s, x + w + Cm(0.12), y1 + h / 2 - Cm(0.25), Cm(0.72))
    x += w + gap
tf = textbox(s, Cm(0.9), Cm(5.3), Cm(32), Cm(12.5))
bullets(tf, [
    (0, "각 단계의 분석 그림 (같은 시연의 세 가지 표현 — 패턴 동일, 좌표계·표현만 다름)", {"bold": True, "size": 15}),
    (1, "omx_workspace.png — [1] OMX 좌표 (x 0.05~0.35m, 기준면 28mm)"),
    (1, "tx90_v5_workspace.png — [3] TX90 EE 좌표 (x 0.60~0.85m, 기준면 334mm) ← 오늘 생성"),
    (1, "tx90_v6_workspace.png — [4] TX90 관절 공간 (관절 분포·연속성 패널) ← 오늘 생성"),
    (0, "[4] 단계에 들어간 변환 (오늘 확정한 규약)", {"bold": True, "size": 15, "after": 6}),
    (1, "① 작업대 오프셋 (−0.20, +0.15)m — 159개 전체가 특이점을 벗어나는 위치 (28종 스윕으로 탐색)"),
    (1, "② 궤적 평활 — 7프레임 이동평균 (OMX 시연 지터 저역 필터)"),
    (1, "③ 앵커 IK — 6Hz 에서 체인 IK, 사이 30Hz 는 관절 선형 보간 (RViz replay 검증 방식)"),
    (1, "④ state = action[t-1] — OMX 측정 노이즈 배제, 완전 추종 가정"),
    (1, "⑤ DP 전역 탐색 폴백 — 그리디 체인이 막히는 에피소드용 (ep143 등)"),
    (0, "결과: EE pose 시연 159개 → 단일 관절 배치의 매끄러운 joint 궤적 159개", {"bold": True, "size": 15, "after": 6}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-14")

# ═══════════ 2b. 데이터 분포 그림 (v5 · v6) ═══════════
s = slide()
bar(s, "같은 시연, 두 가지 표현 — 데이터셋 v5(EE) · v6(joint) 전수 분석", "PIPELINE")
s.shapes.add_picture(os.path.join(WSDIR, "tx90_v5_workspace.png"),
                     Cm(0.9), Cm(2.4), width=Cm(32.0))
s.shapes.add_picture(os.path.join(WSDIR, "tx90_v6_workspace.png"),
                     Cm(0.9), Cm(10.4), width=Cm(32.0))
tf = textbox(s, Cm(0.9), Cm(17.9), Cm(32), Cm(1.0))
bullets(tf, [
    (0, "위: v5 (TX90 EE, 기준면 334mm) — 아래: v6 (관절 분포·연속성, 오프셋 반영 실행 위치). "
        "grasp/release 패턴·그리퍼 쌍봉은 동일 — 좌표계와 표현만 다름", {"size": 11.5, "color": GREY}),
])
foot(s, "생성: plot_v5_v6_workspace.py · OMX 원본 분석(omx_workspace.png)과 3부작")

# ═══════════ 3. 오늘 작업 타임라인 ═══════════
s = slide()
bar(s, "오늘 작업 흐름 (시간순)", "PIPELINE")
rows = [
    ["단계", "작업", "결과"],
    ["오전", "v6 착수 — 데이터셋 v5 구조 조사, MoveIt 재기동 (컨테이너 재시작으로 다운)", "159ep × 600프레임 × state/action"],
    ["오전", "1차 시도: 30Hz 전 프레임 IK 체인 → ep0 에서 353/1200 실패", "트러블슈팅 ① 시작"],
    ["오전", "원인 3종 규명: 손목 self-motion 드리프트 · KDL 특이점 스톨 · state 지터", "알고리즘 재설계"],
    ["오전", "앵커 6Hz+보간 · 평활 · state 재정의 · 오프셋 28종 스윕 · DP 폴백", "(−0.20,+0.15) 에서 159/159"],
    ["점심", "검증 4종 (점프·FK·브랜치·kNN) → 데이터셋 v6 조립", "전부 통과"],
    ["10:35~13:28", "ACT v6 학습 100,000 스텝 (2.9시간)", "loss 0.052, train_tx90_act_v6_joint"],
    ["오후", "평가: ep0/80/140 open-loop 관절오차 + 연속성 실측", "평균 0.93~1.55°, 그리퍼 96.8~99.2%"],
    ["오후", "RViz 재생 — 예측 joint 직접 발행", "IK 호출 0회, v5 의 손목 공회전 소멸"],
    ["오후", "v5/v6 워크스페이스 분석 그림 생성 + 변환 계보 정리", "tx90_v5/v6_workspace.png"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[2.4, 10.0, 5.8], size=12, rh=Cm(1.48))
foot(s, "OMX→TX2-90 전이 · 2026-08-14")

# ═══════════ 4. v5 vs v6 모델 차이 (구체) ═══════════
s = slide()
bar(s, "모델 v5 vs 모델 v6 — 구체적으로 뭐가 다른가", "CORE")
rows = [
    ["", "모델 v5 (EE pose 출력)", "모델 v6 (joint 출력)"],
    ["디코더 출력 7D 의 의미", "손끝 위치 3 + 자세 3 + 그리퍼 1", "관절각 6 (j1~j6) + 그리퍼 1"],
    ["학습 데이터 출처", "데이터셋 v5 (시연의 EE 표현)", "같은 시연 — v5 를 오프라인 IK 로 번역한 v6"],
    ["state 입력", "현재 EE pose + 그리퍼", "현재 관절각 + 그리퍼 (로봇 고유감각)"],
    ["실행 시 IK", ("필요 — 매 청크 IK 번역 (0.1~0.5초, 실패·표류 가능)", {"color": RED}),
     ("없음 — 출력이 곧 로봇 명령", {"color": GREEN, "bold": True})],
    ["연속성 보장", "런타임 IK 체인 알고리즘이 강제", "데이터 규약 + 학습 + 안전 필터(|Δj| 검사)"],
    ["작업 위치", "실행 시 오프셋 인자로 지정", "데이터에 구움 — (−0.20,+0.15) 고정"],
    ["OMX 데이터 호환", ("유지 (EE 는 로봇 중립) — 전이학습 노선", {"color": GREEN}), "단절 — TX90 전용"],
    ["모델 구조", "ACT (동일)", "ACT (동일 — 데이터만 다름, 7D 차원도 우연히 동일)"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[4.2, 7.0, 6.8], size=11.5, rh=Cm(1.35))
tf = textbox(s, Cm(0.9), Cm(15.4), Cm(32), Cm(3.2))
bullets(tf, [
    (0, "관계 주의: v6 는 \"모델 v5 의 출력\"이 아니라 \"원본 시연\"에서 나옴 — 두 모델은 같은 스승의 형제", {"bold": True, "size": 14}),
    (1, "모델 v5 출력으로 만들면 예측 오차까지 정답으로 배우는 복사본의 복사본 — 그렇게 하지 않음"),
    (1, "IK 는 사라진 게 아니라 학습 전(오프라인)으로 이동 — 시간 제약 없는 시점에 전역 탐색·검증까지 완료", {"color": GREY}),
])
foot(s, "박사님 질문 회신의 실물: \"디코더에서 joint 출력\" 완성")

# ═══════════ 4b. 데이터셋 해부 ① — v5 의 한 프레임 ═══════════
s = slide()
bar(s, "데이터셋 해부 ① — v5 의 한 프레임 (에피소드 0, t=227 · 집는 순간)", "DATA")
rows = [
    ["필드", "형태", "내용"],
    ["observation.images.camera1 / 2", "480×640 RGB (mp4)", "OMX 수집 당시 두 시점 영상 — 아래 실물 프레임"],
    ["observation.state", "float 7", "그 순간의 EE pose 6D + 그리퍼 1D (action 과 같은 규약)"],
    ["action", "float 7", "명령 EE pose — [x, y, z, rx, ry, rz, grip]"],
    ["timestamp / frame·episode·index", "float / int", "30Hz 타임스탬프, 프레임·에피소드 번호"],
]
table(s, Cm(0.9), Cm(2.5), Cm(21.6), rows, widths=[6.2, 3.6, 9.4], size=11, rh=Cm(1.1))
tf = textbox(s, Cm(0.9), Cm(8.6), Cm(21.6), Cm(9.4))
bullets(tf, [
    (0, "실제 값 — action[227] (parquet 에서 그대로 읽은 수치)", {"bold": True, "size": 14.5}),
    (1, "[0.7771, −0.1251, 0.3862,  3.3821, 0.8397, 3.0977,  0.451]", {"mono": True, "size": 13, "bold": True}),
    (2, "x=0.777m (로봇 앞) · y=−0.125m (오른쪽) · z=0.386m (기준면 334mm 위 5cm)"),
    (2, "rx,ry,rz = 3.38, 0.84, 3.10 rad = 193.8°, 48.1°, 177.5° (extrinsic xyz 오일러) — 그리퍼가 아래를 봄"),
    (2, "grip = 0.451 rad < 임계 0.459 → 닫힘 = 지금 집는 중", {"color": RED}),
    (0, "이 수치의 유래 (여기가 핵심)", {"bold": True, "size": 14.5, "after": 5}),
    (1, "OMX 시연 관절 기록 → FK → OMX EE pose → Umeyama 유사변환 (s·R·p + t, 4코너 대응으로 산출)"),
    (1, "+ 툴 프레임 보정 (Ry+90°) → TX90 base_link 기준 tool0 pose 로 이식된 값"),
    (1, "즉 \"OMX 로봇이 한 동작을 TX90 몸으로 옮기면 손끝이 있어야 할 곳\"", {"color": GREY}),
    (0, "규모: 159 에피소드 · 76,345 프레임 · 30Hz · 에피소드당 ~20초", {"bold": True, "size": 13.5, "after": 4}),
])
s.shapes.add_picture(os.path.join(IMG14, "cam1_grasp.png"), Cm(23.2), Cm(2.6), width=Cm(9.6))
s.shapes.add_picture(os.path.join(IMG14, "cam2_grasp.png"), Cm(23.2), Cm(10.2), width=Cm(9.6))
tf = textbox(s, Cm(23.2), Cm(17.6), Cm(9.6), Cm(1.0))
bullets(tf, [(0, "t=227 의 camera1(위)·camera2(아래) — OMX 장면임에 주의(시각 도메인 갭의 실체)", {"size": 10, "color": GREY})])
foot(s, "같은 프레임이 v6 에서 어떻게 바뀌는지 → 다음 장")

# ═══════════ 4c. 데이터셋 해부 ② — v6 의 같은 프레임 ═══════════
s = slide()
bar(s, "데이터셋 해부 ② — v6 의 '같은' 프레임 (표현만 관절로 번역됨)", "DATA")
rows = [
    ["", "v5 action[227]", "v6 action[227]"],
    ["값 (저장 단위)", ("[0.7771, −0.1251, 0.3862, 3.3821, 0.8397, 3.0977, 0.451]", {"mono": True, "size": 10.5}),
     ("[0.002, 0.627, 2.009, 0.737, −0.421, −0.654, 0.451]", {"mono": True, "size": 10.5})],
    ["의미", "손끝이 있어야 할 위치·자세 (m, rad)", "관절이 가야 할 각도 (rad) = 0.1°, 35.9°, 115.1°, 42.2°, −24.1°, −37.5°"],
    ["그리퍼 (7번째)", "0.451 rad — 동일 (번역 대상 아님)", "0.451 rad — 그대로 복사"],
    ["카메라 영상", "동일 (mp4 링크 공유)", "동일 — 픽셀 단위로 같은 파일"],
    ["차이의 유래", "—", "v5 값에 오프셋(x−0.20, y+0.15) 적용 후 IK 로 관절 번역"],
]
table(s, Cm(0.9), Cm(2.5), Cm(32.1), rows, widths=[3.4, 7.2, 7.4], size=11, rh=Cm(1.45))
tf = textbox(s, Cm(0.9), Cm(11.6), Cm(32), Cm(6.4))
bullets(tf, [
    (0, "state 의 정의가 v6 에서 바뀐 것도 실측 값으로 확인 가능", {"bold": True, "size": 14.5}),
    (1, "v6 state[t] 의 관절 6개 = action[t−1] 의 관절 6개 (완전 추종 가정) — ep0 전 프레임에서 일치 검증됨", {"mono": False}),
    (1, "이유: v5 의 state 는 OMX '측정' FK 라 백래시 지터 포함 — TX90 에 무의미한 노이즈라 배제"),
    (1, "그리퍼만 원본 측정값 유지 (state[0] grip=0.690 vs action[0] grip=0.687 처럼 미세하게 다름)"),
    (0, "정리: v5 와 v6 는 프레임 수·에피소드 수·영상·그리퍼·타임스탬프까지 전부 같고, state/action 의 6개 숫자의 '언어'만 다르다",
     {"bold": True, "size": 14, "after": 4}),
    (1, "그래서 모델 구조 변경이 0 — 7차원 회귀 헤드는 그대로, 정규화 통계(mean/std)만 v6 것으로 재계산", {"color": GREY}),
])
foot(s, "v6 meta/conversion_meta.json 에 변환 규약 전체 기록")

# ═══════════ 4d. 모델 입출력 나란히 (구체 벡터) ═══════════
s = slide()
bar(s, "모델 구조는 동일 — 입출력 벡터의 '의미'만 교체", "DATA")
tf = textbox(s, Cm(0.9), Cm(2.3), Cm(32), Cm(0.8))
bullets(tf, [(0, "모델 v5 — EE pose 를 읽고 EE pose 를 낸다", {"bold": True, "size": 14, "color": NAVY})])
flowbox(s, Cm(0.9), Cm(3.15), Cm(6.6), Cm(2.4), "카메라 2대 480×640\n+ state 7D\n[x,y,z,rx,ry,rz,grip]", fill=BGSOFT, fg=INK, size=11)
flowbox(s, Cm(8.3), Cm(3.15), Cm(9.2), Cm(2.4),
        "ACT (공통 구조)\nResNet18×2 → 토큰 602개\n인코더 4층·512D → 디코더 쿼리 100", size=11)
flowbox(s, Cm(18.3), Cm(3.15), Cm(14.4), Cm(2.4),
        "출력: EE pose 청크 100×7\n예: [0.7771, −0.1251, 0.3862, 3.38, 0.84, 3.10, 0.451]\n→ 실행하려면 IK 번역 필요", fill=RED, size=11)
arrow(s, Cm(7.55), Cm(4.1)); arrow(s, Cm(17.55), Cm(4.1))
tf = textbox(s, Cm(0.9), Cm(6.0), Cm(32), Cm(0.8))
bullets(tf, [(0, "모델 v6 — joint 를 읽고 joint 를 낸다", {"bold": True, "size": 14, "color": GREEN})])
flowbox(s, Cm(0.9), Cm(6.85), Cm(6.6), Cm(2.4), "카메라 2대 480×640 (동일)\n+ state 7D\n[j1..j6, grip]", fill=BGSOFT, fg=INK, size=11)
flowbox(s, Cm(8.3), Cm(6.85), Cm(9.2), Cm(2.4),
        "ACT (완전 동일 구조)\n가중치만 v6 데이터로 재학습\n(100k 스텝, 2.9h)", size=11)
flowbox(s, Cm(18.3), Cm(6.85), Cm(14.4), Cm(2.4),
        "출력: joint 청크 100×7\n예: [0.002, 0.627, 2.009, 0.737, −0.421, −0.654, 0.451]\n→ 그대로 로봇 명령 (IK 불필요)", fill=GREEN, size=11)
arrow(s, Cm(7.55), Cm(7.8)); arrow(s, Cm(17.55), Cm(7.8))
tf = textbox(s, Cm(0.9), Cm(9.8), Cm(32), Cm(8.0))
bullets(tf, [
    (0, "두 예시 벡터는 같은 순간(ep0, t=227 집는 순간)의 실제 데이터 — 같은 물리적 사건의 두 언어", {"bold": True, "size": 14.5}),
    (1, "v5 벡터를 읽는 법: \"손끝을 (0.777, −0.125, 0.386)m 로, 자세는 (194°, 48°, 178°), 그리퍼 닫아라\""),
    (1, "v6 벡터를 읽는 법: \"관절을 (0.1°, 35.9°, 115.1°, 42.2°, −24.1°, −37.5°) 로, 그리퍼 닫아라\""),
    (1, "v6 의 위치가 v5 와 다른 이유: 작업대 오프셋 (x−0.20, y+0.15) 이 관절값 안에 이미 반영되어 있음"),
    (0, "공통으로 안 바뀐 것: 카메라 영상(OMX 장면) · 그리퍼 값 · chunk 100 · CVAE z · 학습 손실(L1+KL)", {"bold": True, "size": 13.5, "after": 4}),
    (1, "그래서 \"어떤 모델이냐\"의 차이가 아니라 \"무엇을 가르쳤냐\"의 차이 — 구조 도식은 한 장으로 충분", {"color": GREY}),
])
foot(s, "실측 값 출처: 각 데이터셋 episode_000000.parquet · 카메라 프레임은 앞 장 참조")

# ═══════════ 4e. 입력 joint vs 출력 joint — 시제가 다르다 ═══════════
s = slide()
bar(s, "입력 joint 와 출력 joint 는 뭐가 다른가 — \"시제\"가 다르다", "MODEL")
flowbox(s, Cm(0.9), Cm(2.6), Cm(8.0), Cm(2.6),
        "state (입력)\n관절 벡터 1개\n= 지금 몸이 어디 있나", fill=BGSOFT, fg=INK, size=13)
flowbox(s, Cm(9.9), Cm(2.6), Cm(5.6), Cm(2.6), "ACT", size=14)
flowbox(s, Cm(16.5), Cm(2.6), Cm(16.4), Cm(2.6),
        "action (출력)\n관절 벡터 100개 (미래 3.3초)\n= 앞으로 관절을 이렇게 움직여라", fill=GREEN, size=13)
arrow(s, Cm(9.05), Cm(3.65)); arrow(s, Cm(15.65), Cm(3.65))
tf = textbox(s, Cm(0.9), Cm(5.7), Cm(32), Cm(12.2))
bullets(tf, [
    (0, "모델이 푸는 문제를 한국어로: \"이 장면이 보이고(주사위 위치), 관절이 지금 여기 있을 때 → 100스텝 동안 이렇게 움직여라\"", {"bold": True, "size": 14.5}),
    (0, "실제 값으로 보면 (같은 언어, 다른 시제)", {"bold": True, "size": 14, "after": 4}),
    (1, "state      = [0.1°, 35.9°, 115.1°, 42.2°, −24.1°, −37.5°]   ← 지금", {"mono": True, "size": 12}),
    (1, "action[0]  = [0.1°, 36.0°, 115.0°, 42.3°, −24.2°, −37.4°]   ← 0.03초 뒤 명령 (거의 같음)", {"mono": True, "size": 12}),
    (1, "action[50] = [−1.2°, 38.4°, 112.3°, ...]                    ← 1.7초 뒤 (꽤 다름)", {"mono": True, "size": 12}),
    (1, "action[99] = [−2.8°, 40.1°, 109.8°, ...]                    ← 3.3초 뒤", {"mono": True, "size": 12}),
    (1, "출력은 \"현재의 복사본\"이 아니라 현재에서 출발해 뻗어나가는 궤적 — 뒤로 갈수록 state 에서 멀어짐"),
    (0, "비유: 내비게이션", {"bold": True, "size": 14, "after": 4}),
    (1, "입력 = 내 차의 현재 위치 (GPS 점 하나) · 출력 = 앞으로 갈 경로 (3.3초 분량의 길) — 둘 다 좌표라는 같은 언어, 역할이 다름"),
    (0, "데이터셋 안의 사슬 구조: state[t] = action[t−1] (직전 명령이 현재 상태가 되고, 거기서 다음 명령들이 이어짐)", {"bold": True, "size": 14, "after": 4}),
])
foot(s, "v5 도 구조는 동일 — 단어(EE pose)만 다름")

# ═══════════ 4f. 왜 v6 는 다른 브랜치로 안 꺾이나 ═══════════
s = slide()
bar(s, "왜 v6 는 다른 관절 배치(브랜치)로 안 꺾이나 — 시드 코드 없이도", "MODEL")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(15.8))
bullets(tf, [
    (0, "결론 먼저: v6 실행부에는 \"직전 해를 시드로\" 같은 코드가 없다 — 그리고 없어도 되는 이유가 v6 의 본질", {"bold": True, "size": 15}),
    (0, "브랜치 선택은 IK 에만 존재하는 사건이다", {"bold": True, "size": 14, "after": 4}),
    (1, "\"한 pose 에 여러 관절 정답 → 그중 하나를 고름\"이 IK. v6 는 실행 때 IK 를 안 풀므로 고를 후보 목록 자체가 없음 — 꺾일 무대가 없다"),
    (0, "그래도 항상 같은 배치의 값이 나오는 3겹의 이유", {"bold": True, "size": 14, "after": 4}),
    (1, "① 다른 브랜치를 본 적이 없다 — 변환 때 체인+DP 로 159개 전부를 단일 배치로 구움 (kNN 0/4770 검증). "
        "회귀 모델은 배운 패턴만 재현 — 팔꿈치 뒤집힌 배치는 출력할 줄 모름"),
    (2, "비유: 통역사(IK)는 매번 여러 번역 후보 중 고르다 실수할 수 있지만, 한 가지 말투만 듣고 자란 원어민(v6)은 다른 말투가 입에서 안 나온다", {"color": GREY}),
    (1, "② state 가 시드 역할을 '암묵적으로' 대신 — \"state 가 j 근처면 action 은 j 에서 이어진다\"는 패턴을 학습. "
        "명시적 코드가 아니라 학습된 조건부 습관"),
    (1, "③ 청크 통짜 생성 — 100스텝을 한 번의 계산으로 출력. 스텝을 따로 풀어 이어붙이다 어긋나는 구조 자체가 없음 (실측: 프레임당 평균 0.32°)"),
    (0, "정직한 각주: 이것은 알고리즘적 '보장'이 아니라 통계적 '습관'", {"bold": True, "size": 14, "color": RED, "after": 4}),
    (1, "분포 밖 입력(못 본 화면·이상한 state)에서는 엉뚱한 값이 나올 수 있음 → 실행단에 안전 필터(|Δjoint| 한계 검사)를 둠"),
    (1, "안전 필터 = 옛 시드 체인의 \"90° 거부\"가 남긴 유일한 후손 — 단 IK 가 아니라 뺄셈 한 번이라 실시간성 무해", {"color": GREY}),
])
foot(s, "시드 체인의 기능 이관: 브랜치 선택→데이터에 소멸 · 연속성→state 조건부 예측 · 점프 보장→안전 필터")

# ═══════════ 4g. state 깊이 알기 — 합성값과 실측값 ═══════════
s = slide()
bar(s, "state 깊이 알기 — 지금은 '만든 값', 실기에선 '잰 값'", "MODEL")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(15.8))
bullets(tf, [
    (0, "v6 의 state 는 합성값이다 — state[t] = action[t−1] (\"로봇이 직전 명령을 완벽히 따라갔다고 치자\")", {"bold": True, "size": 15}),
    (1, "이유: v6 를 만들 때 TX90 은 한 번도 안 움직였음 (전부 책상 위 변환) → 측정값이 존재할 수 없어서 가정으로 대체"),
    (1, "학습 때 본 세상:  명령 35.90° → state 35.90° (오차 0 인 이상적 세계)", {"mono": True, "size": 12}),
    (1, "실기의 세상:      명령 35.90° → 모터 실제 35.87° → state 35.87° (추종 오차 ~0.1°)", {"mono": True, "size": 12}),
    (1, "차이는 먼지 수준(관절 가동 범위 수십°대비)이라 작은 문제 — 실기 수집 때 진짜 측정값을 기록하면 자동 해소"),
    (0, "그럼 실기의 \"진짜 측정값\"은 joint 인가 EE pose 인가 → joint 다", {"bold": True, "size": 15, "after": 6}),
    (1, "로봇에 달린 센서는 관절 각도계(엔코더) 6개뿐 — 손끝 위치를 재는 센서는 없음"),
    (1, "사람 비유: 눈 감고도 손이 어디 있는지 아는 건 손에 GPS 가 있어서가 아니라, 어깨·팔꿈치 굽은 각도를 알고 뇌가 계산하기 때문"),
    (1, "펜던트의 \"X=611mm\" 표시도 실은 엔코더 → FK 계산 결과 — EE pose 는 언제나 측정이 아닌 계산값"),
    (0, "그래서 실기 수집의 철칙: joint 를 원본으로 저장", {"bold": True, "size": 15, "after": 6}),
    (1, "joint → EE (FK): 답이 하나, 언제든 재계산 가능  /  EE → joint (IK): 답이 여럿 — 이번에 고생한 그 문제"),
    (1, "한 번의 수집으로 두 노선 커버: v6 학습엔 joint 그대로, v5 노선엔 FK 돌린 EE 로", {"color": GREY}),
])
foot(s, "v6(joint 모델)가 실기 데이터와 궁합이 가장 좋은 이유 — 측정한 그대로가 학습 입력")

# ═══════════ 4h. 기준 좌표계 — 숫자는 전부 TX90 ═══════════
s = slide()
bar(s, "기준 정리 — 모든 숫자는 TX90 기준, OMX 는 어디에 남아 있나", "MODEL")
rows = [
    ["값", "기준", "읽는 법"],
    ["v5 의 EE pose", ("TX90", {"bold": True}), "TX90 base_link 좌표계에서 본 tool0(손끝). x=0.777m = \"TX90 몸통 앞 77.7cm\""],
    ["v6 의 joint", ("TX90", {"bold": True}), "TX90 의 관절 6개. OMX 는 4축이라 6개 값이 나올 수조차 없음"],
    ["카메라 영상", ("OMX", {"color": RED, "bold": True}), "화면에 찍힌 팔·배경은 OMX 그대로 — 유일하게 남은 OMX 흔적 = 시각 도메인 갭"],
]
table(s, Cm(0.9), Cm(2.5), Cm(32.1), rows, widths=[3.6, 2.0, 12.4], size=12, rh=Cm(1.5))
tf = textbox(s, Cm(0.9), Cm(9.2), Cm(32), Cm(8.6))
bullets(tf, [
    (0, "혼동 포인트: 시연의 '출처'는 OMX 지만, 값의 '기준'은 TX90", {"bold": True, "size": 15}),
    (1, "OMX 가 시연 (OMX 좌표의 움직임)"),
    (1, "→ Umeyama 변환: \"TX90 몸에 옮기면 손끝이 있어야 할 곳\" 으로 번역  = v5 (TX90 EE)"),
    (1, "→ IK: \"그 손끝을 만들려면 TX90 관절이 취해야 할 각도\" 로 번역  = v6 (TX90 joint)"),
    (1, "v5 시점부터 숫자는 전부 TX90 세계 — 그래서 두 모델의 출력을 RViz/실기에 바로 먹일 수 있음"),
    (0, "카메라 영상의 OMX 흔적은 실기 전이학습(v7)에서 TX90 장면으로 교체됨 — 그때 OMX 는 완전히 은퇴", {"bold": True, "size": 14, "after": 4}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-14")

# ═══════════ 5. 트러블슈팅 ①: 30Hz 전 프레임 IK 의 실패 ═══════════
s = slide()
bar(s, "트러블슈팅 ① — 1차 시도(30Hz 전 프레임 IK)가 무너진 3가지 이유", "TODAY")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(15.8))
bullets(tf, [
    (0, "증상: ep0 에서 1200회 중 353회 IK 실패, 에피소드의 3분의 1이 구멍", {"bold": True, "size": 15, "color": RED}),
    (0, "원인 A — 손목 self-motion 드리프트 (실측)", {"bold": True, "size": 14, "after": 4}),
    (1, "j5≈0° 특이점 근처에서는 j4/j6 를 반대로 돌려도 손끝이 그대로 → 해가 '헐거워짐'"),
    (1, "30Hz 미세 스텝을 따라가던 KDL 해가 그 자유도를 타고 표류 → j4=−265°, j6=+261° 로 한계(±270°) 림에 고착"),
    (1, "림에 닿으면 2π 등가각 접기가 한계 밖으로 나가 불능 → 이후 전부 연쇄 실패"),
    (0, "원인 B — KDL 특이점 스톨", {"bold": True, "size": 14, "after": 4}),
    (1, "특이 자세를 시드로 주면 야코비안이 병들어 수치 반복이 수렴 실패 — 해가 존재해도 못 찾음 (ep7 실측)"),
    (0, "원인 C — state 지터", {"bold": True, "size": 14, "after": 4}),
    (1, "observation.state 는 OMX '측정' 관절의 FK — 백래시 노이즈가 커서 특이점 근처에서 풀 수 없는 미세 자세 요구 생성"),
    (1, "진단 결정타: 같은 시점에서 state 는 실패, action 은 같은 시드로 즉시 성공 → state 만의 문제로 분리"),
    (0, "교훈: 예측 궤적(어제, 다운샘플)에서 통했다고 원본 30Hz 데이터에서 통하는 게 아니다", {"color": GREY, "size": 13}),
    (1, "원본에는 시연 지터·측정 노이즈·전 구간 커버리지가 있음 — 파이프라인 검증은 반드시 실데이터 전체로", {"color": GREY}),
])
foot(s, "진단 스크립트: diag_ep8.py · diag_action.py (컨테이너 /root)")

# ═══════════ 6. 알고리즘 시도 비교 (하이라이트) ═══════════
s = slide()
bar(s, "시도한 알고리즘 6종과 채택/기각 이유", "TODAY")
rows = [
    ["시도", "결과", "판정 / 이유"],
    ["① 30Hz 전 프레임 IK 체인 (어제 방식 그대로)", "ep0 353/1200 실패",
     ("기각 — 특이점 self-motion 드리프트로 림 고착 (원인 A)", {"color": RED})],
    ["② 강한 시드 바이어스(j4/j6→0) + 점프 한계 20° 강화", "오히려 악화 (ep120 589 실패)",
     ("기각 — KDL 은 같은 수렴 분지로 돌아감, 특이점의 관절속도 발산은 물리 법칙이라 못 막음", {"color": RED})],
    ["③ 앵커 6Hz IK + 30Hz 관절 보간 + 평활(7프레임)", "실패 대폭 감소, FK 오차 정량화 가능",
     ("채택 — 어제 RViz replay 로 검증된 방식의 일반화. 보간이 특이점을 관절 공간에서 무리 없이 통과", {"color": GREEN, "bold": True})],
    ["④ state 도 IK 로 풀기 (이중 체인)", "state 만 집중 실패 (원인 C)",
     ("기각 → state=action[t-1] 재정의 채택 — OMX 추종 오차는 TX90 에 무의미, 연속성 자동 보장", {"color": GREEN})],
    ["⑤ 오프셋 (−0.10,0) 유지 vs 전역 스윕 28종", "(−0.10,0): 131/159 → (−0.20,+0.15): 158/159",
     ("스윕 결과 채택 — 어제의 −10cm 는 예측궤적 1개 기준이었음. 전체 데이터 기준으로 갱신", {"color": GREEN, "bold": True})],
    ["⑥ 그리디 단독 vs DP 전역 탐색 폴백", "그리디 158 + DP 1 (ep143) = 159/159",
     ("하이브리드 채택 — DP 는 팔 가족 프라이어(j1~j3 HOME 계열)로 브랜치 일관성 유지", {"color": GREEN})],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[6.4, 4.4, 7.8], size=11, rh=Cm(1.95))
tf = textbox(s, Cm(0.9), Cm(15.0), Cm(32), Cm(3.2))
bullets(tf, [
    (0, "설계 원칙: 실시간이 아니라 오프라인이므로 \"그 자리에서 제일 좋은 해\"(그리디)가 아니라 \"에피소드 전체를 보고 최적\"(DP)까지 쓸 수 있다", {"bold": True, "size": 13.5}),
    (1, "이것이 v6 방식의 본질적 이점 — 런타임 IK 는 그리디밖에 못 함 (미래를 모르니까)", {"color": GREY}),
])
foot(s, "산출물: convert_v6.py (규약 전체 구현) · scan_greedy.py (스윕)")

# ═══════════ 7. 오프셋 스윕 + 작업대 결정 갱신 ═══════════
s = slide()
bar(s, "⚠ 작업대 위치 결정 갱신 — (−10cm, 0) → (−20cm, +15cm)", "DECISION")
rows = [
    ["오프셋 (dx, dy)", "그리디 성공 / 159", "비고"],
    ["(−0.10, 0)  ← 어제 결정", "131", ("예측궤적 1개 기준의 한계 노출", {"color": RED})],
    ["(−0.10, +0.05)", "131~134", "평활·앵커 적용 후에도 부족"],
    ["(−0.15, +0.10)", "153", "급격히 개선 — 방향 확인"],
    ["(−0.20, +0.10) / (−0.18,+0.12) / (−0.15,+0.15)", "156", "고원 지대 진입"],
    ["(−0.20, +0.15)  ← 최종 채택", ("158 (+DP 1 = 159)", {"bold": True, "color": GREEN}), ("전 에피소드 통과, 점프 최대 16.2° 로 최소", {"color": GREEN})],
    ["(−0.22, +0.17) 등 주변", "158", "고원 — ±2cm 여유 확인됨"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[6.5, 3.6, 7.0], size=11.5, rh=Cm(1.3))
tf = textbox(s, Cm(0.9), Cm(12.3), Cm(32), Cm(5.8))
bullets(tf, [
    (0, "의미: 159개 시연 전체가 손목 특이점 능선을 벗어나는 위치는 로봇 기준 몸쪽 20cm + 옆 15cm", {"bold": True, "size": 14.5}),
    (1, "ep143 하나만 어느 오프셋에서도 그리디 불가(시연 자체가 특이) → DP 폴백으로 해결"),
    (1, "실물 적용: 테이프 사각형을 이 위치로 이동해야 함 — tx90_v6_workspace.png 의 빨간 사각형이 실측 좌표", {"bold": True, "color": RED}),
    (1, "v5 를 replay 할 때도 --offset=-0.20,0.15,0 으로 맞춰야 동일 조건 비교 가능"),
    (1, "어제 자료(20260813.pptx 6·9장)의 −10cm 표기는 이 슬라이드로 대체됨", {"color": GREY}),
])
foot(s, "스윕 28종 × 159 에피소드 전수 — 변환이 에피소드당 ~0.2초라 가능했던 탐색")

# ═══════════ 8. 검증 4종 ═══════════
s = slide()
bar(s, "데이터셋 v6 검증 4종 — 전부 통과", "VERIFY")
rows = [
    ["검증", "방법", "결과", "판정"],
    ["① 변환 완전성", "159 에피소드 × 앵커 체인 (+DP 폴백)", "159/159, 실패 앵커 보간 간격 ≤2앵커",
     ("통과", {"color": GREEN, "bold": True})],
    ["② 연속성", "30Hz 프레임 간 max|Δjoint| 전수", "최대 16.2° / 평균 4.9° (임계 30°)",
     ("통과", {"color": GREEN, "bold": True})],
    ["③ FK 역검증", "/compute_fk 로 에피소드당 15프레임 왕복 비교", "위치 ≤5.1mm · 자세 ≤1.0°",
     ("통과", {"color": GREEN, "bold": True})],
    ["④ 다봉성 (kNN)", "정책 입력(state joint) 이웃 → action joint 거리", "4770 샘플 중 다봉 0건",
     ("통과", {"color": GREEN, "bold": True})],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[3.2, 7.2, 6.2, 1.6], size=11.5, rh=Cm(1.5))
tf = textbox(s, Cm(0.9), Cm(10.6), Cm(32), Cm(7.2))
bullets(tf, [
    (0, "④번의 교훈 — 검사 기준을 잘못 잡으면 가짜 경보가 뜬다", {"bold": True, "size": 14.5}),
    (1, "처음에 'EE pose 가 비슷한데 joint 가 다른' 쌍을 찾자 41%가 걸림 — 최악은 366° 차이"),
    (1, "실체: ±360° 감김 표현 차이 + 감김 상태가 다른 것 — 그런데 v6 의 정책 입력(state)은 joint 라서 감김 상태가 입력에 포함됨"),
    (1, "→ \"같은 입력, 다른 출력\"이어야 다봉. 정책 입력 공간에서 재검사하니 0건 — 데이터는 정책 관점에서 완전한 단봉", {"bold": True}),
    (0, "잔여 리스크 (정직하게)", {"bold": True, "size": 14.5, "after": 5}),
    (1, "특이점 통과 순간 최대 16.2°/frame(≈490°/s) 스윙이 데이터에 잔존 — 실기는 감속 실행 + 안전 필터 전제"),
    (1, "FK 5.1mm 는 보간·평활의 대가 — 주사위 크기 대비 무해하나 정밀 태스크라면 앵커 간격 축소 필요"),
])
foot(s, "검증 스크립트: convert_v6.py 내장 + knn_check.py · 리포트: v6 meta/conversion_report.json")

# ═══════════ 9. 학습·평가 결과 ═══════════
s = slide()
bar(s, "모델 v6 학습·평가 — v5 와 같은 급의 정확도, IK 없는 실행", "RESULT")
rows = [
    ["항목", "값"],
    ["학습", "ACT 100,000 스텝 · batch 8 · 2.9시간 (10:35~13:28) · 최종 loss 0.052"],
    ["관절 오차 (open-loop, ep0/80/140)", "평균 1.09° / 0.93° / 1.55° · 중앙값 0.43~0.59°"],
    ["관절별 오차", "j1~j3: 0.2~0.7° (매우 정확) · j4/j6: 2~3.7° (특이점 스윙 구간에 집중)"],
    ["그리퍼 개폐 일치율", "99.2% / 96.8% / 99.2% (임계 0.459)"],
    ["출력 연속성", "예측 프레임 간 평균 0.32° · 최대 30.5° (청크 경계 1건 — 안전 필터 필요성 재확인)"],
    ["RViz 재생", "예측 joint 직접 발행 — IK 호출 0회, 20초 전 구간 매끄러움"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[5.2, 12.8], size=11.5, rh=Cm(1.3))
s.shapes.add_picture(os.path.join(IMG14, "v6_f2.png"),
                     Cm(24.6), Cm(12.3), height=Cm(5.6))
s.shapes.add_picture(os.path.join(IMG14, "v6_f4.png"),
                     Cm(29.3), Cm(12.3), height=Cm(5.6))
tf = textbox(s, Cm(0.9), Cm(12.4), Cm(23.2), Cm(5.6))
bullets(tf, [
    (0, "관찰: RViz 에서 v5 는 4번 축이 계속 헛돌았는데 v6 는 안 돈다 — 이유가 곧 이번 작업의 요약", {"bold": True, "size": 14.5}),
    (1, "v5 공회전 = 런타임 IK 가 특이점 근처의 '헐거운' j4/j6 자유도 안에서 표류한 것 (태스크와 무관한 모터 낭비)"),
    (1, "v6 는 ① 특이점에서 먼 작업 위치 ② 변환 규약의 표류 억제 ③ 학습의 평활 — 3겹으로 원인 제거"),
    (1, "실기 함의: v5 방식이면 공회전이 실제 모터 동작 (마모·케이블·한계 정지) — v6 는 구조적으로 재발 불가"),
    (0, "손끝 궤적은 두 방식이 사실상 동일 — 다른 것은 손목이 헛도느냐뿐. 랩미팅 데모 포인트", {"color": GREY, "size": 13}),
])
foot(s, "평가: eval_v6.py · 재생: replay_v6.py")

# ═══════════ 10. open-loop 청크 vs 실시간 청크 ═══════════
s = slide()
bar(s, "현재 위치 — 지금 RViz 재생은 '실시간 청크'가 아니라 'open-loop 청크'다", "STATUS")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(6.6))
bullets(tf, [
    (0, "청크 출력 자체는 이미 작동 중 — 추론 1회 = 100스텝(3.3초) 청크, 600프레임 추론이 6초에 끝난 이유", {"bold": True, "size": 14.5}),
    (0, "그러나 지금까지의 모든 실행은 open-loop:", {"bold": True, "size": 14.5, "after": 4}),
    (1, "[녹화된 관측 → 추론 → 청크] 를 전부 미리 계산해 저장 → 완성본을 통째로 재생", {"mono": False}),
    (1, "실행 중 카메라가 행동에 영향을 주지 않음 — 주사위를 옮겨도 모름"),
    (0, "닫힌 루프(관측→청크→실행→재관측)가 안 되는 이유는 모델이 아니라 루프의 양 끝 부재:", {"bold": True, "size": 14.5, "after": 4}),
    (1, "살아있는 관측 없음 (TX90 카메라 미설치 + 시각 도메인 갭 미해결) · 명령 받을 로봇 없음 (RViz 는 표시 전용)"),
])
rows = [
    ["해결 단계", "무엇이 생기나", "필요한 것", "시점"],
    ["① Isaac Sim (다음 단계)", "렌더 카메라 + 물리 로봇 → 닫힌 루프 최초 가동, 재계획형·반응형 모두 실험",
     "Isaac Sim 설치 + URDF 임포트 + 장면", "이번 주~"],
    ["② 실기 재계획형", "청크마다 관측→실행 (stop-and-go)", "VAL3 드라이버(무료) + TX90 카메라 + 전이학습", "외부제어 개통 후"],
    ["③ 실기 반응형", "실행 중 재추론·청크 이어붙이기 (30Hz)", "CS9 velocity 확장 (견적 확인 중)", "확장 확보 시"],
]
table(s, Cm(0.9), Cm(9.6), Cm(32.1), rows, widths=[4.0, 8.0, 6.6, 2.6], size=11.5, rh=Cm(1.65))
tf = textbox(s, Cm(0.9), Cm(16.6), Cm(32), Cm(1.8))
bullets(tf, [
    (0, "즉 다음 마일스톤 = Isaac Sim 에서 v6 청크의 닫힌 루프 첫 가동 — \"주사위를 옮기면 다음 청크가 따라오는가\"", {"bold": True, "size": 14}),
])
foot(s, "open-loop(검증, 완료) → 재계획형 → 반응형 사다리의 현재 위치: 첫 칸 완료")

# ═══════════ 10b. 실시간 청크의 3개의 열쇠 ═══════════
s = slide()
bar(s, "실시간 청크는 어떻게 가능한가 — 3개의 열쇠 (ACT/ALOHA 검증된 방식)", "REALTIME")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(10.2))
bullets(tf, [
    (0, "열쇠 ① 청크를 다 쓸 의무가 없다 — n_action_steps 하나로 조절 (재학습 불필요)", {"bold": True, "size": 14.5}),
    (1, "open-loop(지금): 100스텝 전부 실행 후 재추론 (3.3초마다 관측)  /  재계획형: 30스텝만 쓰고 재추론 (1초)  /  반응형: 1~2스텝만 쓰고 재추론 (0.03~0.07초)"),
    (1, "뒤의 스텝들은 \"세워두고 즉시 폐기하는 예비 계획\" — 낭비 같지만 이게 정상 사용법"),
    (0, "열쇠 ② 추론을 실행과 겹친다 (비동기)", {"bold": True, "size": 14.5, "after": 4}),
    (1, "로봇: [청크A 실행중......][청크B 실행중......]   /   추론: 실행과 동시에 백그라운드에서 다음 청크 계산 (수십 ms)", {"mono": True, "size": 11.5}),
    (1, "로봇은 항상 버퍼의 관절값을 30Hz 로 꺼내 씀 — 추론 때문에 멈추는 순간이 없음"),
    (0, "열쇠 ③ 청크 이어붙이기 (temporal ensembling)", {"bold": True, "size": 14.5, "after": 4}),
    (1, "청크들이 겹치게 생성됨 → t=50 시점 명령 = 청크A 의 50번째 + 청크B 의 20번째 + 청크C 의 5번째의 가중 평균 (최신에 무게)"),
    (1, "갈아타는 '순간'이 사라지고 부드럽게 미끄러짐 — ACT 내장 기능, 지금은 꺼져 있고 켜면 됨"),
])
tf = textbox(s, Cm(0.9), Cm(12.6), Cm(32), Cm(5.4))
bullets(tf, [
    (0, "주사위를 옮기면 무슨 일이 일어나나 (반응형 기준)", {"bold": True, "size": 14.5}),
    (1, "t=0.0s 관측: 주사위 A → 청크들이 A 로 향함   |   t=0.5s 사람이 B 로 옮김", {"mono": True, "size": 12}),
    (1, "t=0.53s 다음 재추론이 새 화면을 봄 → 새 청크는 B 로   |   t=0.6s~ 앙상블이 A→B 로 부드럽게 휨", {"mono": True, "size": 12}),
    (1, "반응 지연 = 재추론 주기 + 추론 시간 ≈ 0.1초 이내 — 이게 \"실시간 action chunk\"의 실체"),
    (1, "모델·청크 메커니즘은 ACT 원 논문(ALOHA, 50Hz 실기)이 검증한 표준 배포법 — 우리 리스크 목록에서 제외 가능", {"color": GREEN}),
])
foot(s, "남는 관문은 모델이 아니라 인프라 — 다음 장")

# ═══════════ 10c. 실시간 청크 준비물 총정리 ═══════════
s = slide()
bar(s, "실시간 청크(반응형 closed-loop)까지 준비물 총정리", "REALTIME")
rows = [
    ["구분", "#", "항목", "상태 / 비고"],
    [("A. 로봇\n(명령 통로)", {"bold": True}), "1", "VAL3 드라이버 (ros_server + TCP 소켓 4개) — 무료", "조사 완료, 설치 대기"],
    ["", "2", "CS9 velocity/motion 확장 — 4ms 연속 추종, '실행 중 청크 갈아타기'의 열쇠",
     ("최대 관문 — 계정·견적·SRS 에뮬레이터로 판가름 (사람 몫)", {"color": RED, "bold": True})],
    [("B. 관측\n(정책의 눈)", {"bold": True}), "3", "카메라 2대 고정 마운트 (수집 위치 = 실행 위치)", "미착수"],
    ["", "4", "그리퍼 장착 + 제어 연결", "미착수"],
    ["", "5", "스크립트 시연 수집 30~50ep → 전이학습 v7 (시각 갭 + state 실측화)", "외부 제어 개통 후"],
    [("C. 실행기\n(코드만)", {"bold": True}), "6", "비동기 추론 러너 (LeRobot 기성 스택)", "Isaac Sim 에서 먼저 완성"],
    ["", "7", "설정: n_action_steps 축소 + temporal ensembling 켜기 (재학습 불필요)", "〃"],
    ["", "8", "안전 필터 (|Δjoint| 한계 + 속도 제한 + 소프트 리밋) / 지연 보상 튜닝", "〃"],
    [("D. 배치", {"bold": True}), "9", "테이프 사각형을 (x−20cm, y+15cm) 위치로 이동", "tx90_v6_workspace.png 좌표"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[2.6, 0.9, 9.6, 4.9], size=11.5, rh=Cm(1.42))
tf = textbox(s, Cm(0.9), Cm(16.4), Cm(32), Cm(2.0))
bullets(tf, [
    (0, "핵심: A-2 만 사람이 뚫어야 하는 관문이고, B·C 는 병렬 진행 가능 — C 는 Isaac Sim 에서 실기 없이 전부 완성해 둘 수 있다", {"bold": True, "size": 13.5}),
])
foot(s, "A-2 실패 시 플랜 B: 실기=재계획형(1번만으로 가능) + 반응형=Isaac Sim 정량 검증")

# ═══════════ 10d. Isaac Sim — 어디에 왜 쓰나 ═══════════
s = slide()
bar(s, "Isaac Sim — 어디에, 왜 쓰나", "REALTIME")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(4.2))
bullets(tf, [
    (0, "왜 (한 문장): 실기 closed-loop 의 두 관문(명령 통로 A · 관측 B)이 준비되기 전에는 한 발짝도 못 도는데, "
        "Isaac Sim 안에서는 둘 다 공짜다", {"bold": True, "size": 15}),
    (1, "시뮬 로봇은 매 스텝 명령을 그냥 받음 (A 불필요) · 렌더 카메라는 학습 화면과 일치시킬 수 있음 (B 불필요)"),
    (1, "그리고 RViz 는 물리가 없는 애니메이션 — \"실제로 집히는가\"는 Isaac Sim 이 처음 답할 수 있음"),
])
rows = [
    ["용도", "내용", "시점"],
    ["① 물리 검증", "v6 관절 궤적을 물리 엔진 위에 재생 — 그리퍼가 주사위를 실제로 집는지, 충돌은 없는지 첫 증명",
     ("이번 주 목표", {"bold": True})],
    ["② closed-loop 첫 가동", "렌더 카메라(관측) + 물리 로봇(실행)으로 관측→청크→실행→재관측 루프 최초 가동. "
     "재계획형·반응형 실험, 실행기(C-6~8) 튜닝을 실기 전에 완료. \"주사위를 옮기면 따라오는가\"를 처음 보는 곳", "다음"],
    ["③ 데이터 합성 (선택)", "TX90 장면 렌더로 시연 수백 개 자동 생성 + 조명·질감 랜덤화 — 실기 수집 부담 축소, 시각 갭 보조", "필요시"],
]
table(s, Cm(0.9), Cm(7.0), Cm(32.1), rows, widths=[3.4, 11.6, 2.6], size=11.5, rh=Cm(1.9))
tf = textbox(s, Cm(0.9), Cm(13.6), Cm(32), Cm(4.6))
bullets(tf, [
    (0, "전략적 핵심: Isaac Sim 트랙은 CS9 견적 결과와 무관하게 진행된다", {"bold": True, "size": 14.5}),
    (1, "최악의 경우(velocity 확장 불가)에도 반응형 연구는 sim 에서 완결, 실기는 재계획형으로 마무리 — 플랜 B 의 실체"),
    (0, "순서:  지금 [A-2 판가름 착수 ∥ Isaac ① 물리검증] → [Isaac ② closed-loop + 실행기 완성 ∥ 카메라·그리퍼 마운트] "
        "→ 수집·v7 → A-2 결과 따라 실기 반응형 or 재계획형", {"size": 13, "after": 4}),
])
foot(s, "우리 자산으로 바로 시작 가능: tx2_90.urdf 임포트 + 테이블·주사위·카메라 장면 + ROS2 브리지")

# ═══════════ 10e. 전이학습 v7 절차 ═══════════
s = slide()
bar(s, "전이학습 v7 — 실기 closed-loop 의 마지막 조각", "REALTIME")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(32), Cm(6.2))
bullets(tf, [
    (0, "한 줄: TX90 을 실제로 움직여, 카메라 영상 + 관절 측정값을 30~50 에피소드 녹화하고, v6 체크포인트에서 이어 학습", {"bold": True, "size": 15}),
    (0, "수집 (반나절) — 사람이 조종하지 않는다", {"bold": True, "size": 14, "after": 4}),
    (1, "사람: 주사위를 격자 위치에 놓음 (매번 다르게) → 스크립트: v6 예측 궤적을 TX90 에서 실행 (이미 검증된 동작)"),
    (1, "30Hz 기록: 카메라 2대 (드디어 TX90 장면!) + 엔코더 joint 6 (state 의 진짜 측정값) + 그리퍼 → LeRobot 포맷"),
    (1, "좌표 변환·Umeyama·IK 전부 불필요 — 로봇이 직접 움직인 기록이라 처음부터 TX90 값. 변환 파이프라인은 이 시점에 은퇴", {"color": GREY}),
])
tf = textbox(s, Cm(0.9), Cm(9.0), Cm(32), Cm(8.8))
bullets(tf, [
    (0, "학습 (2~3시간)", {"bold": True, "size": 14}),
    (1, "v6 체크포인트(OMX 시연 159개의 동작 지식) + TX90 데이터 30~50개 이어 학습 → v7"),
    (1, "새로 배우는 건 사실상 \"TX90 카메라 화면 해석\" 하나 (동작은 상속) — 그래서 159개가 아니라 30~50개로 충분"),
    (1, "모델 구조·차원 변경 0 — 정규화 통계만 갱신"),
    (0, "v7 이 나오면 열리는 것", {"bold": True, "size": 14, "after": 4}),
    (1, "시각 도메인 갭 해소 (실기 카메라 직결) · state 문제 자동 해소 (실측값 학습) · 실기 closed-loop 첫 데모 (\"주사위 옮기면 따라감\")"),
    (0, "계보 완성:  OMX 159 시연 → v5(EE) → v6(joint) → v7(joint + TX90 시각) — 각 단계가 이전의 지식을 상속", {"bold": True, "size": 14, "after": 4}),
])
foot(s, "전제: 외부 제어 개통 (A-1) — 현재 병목. Stäubli 계정/견적/에뮬레이터가 최우선인 이유")

# ═══════════ 11. v5→v6 좋아진 것 / 감수한 것 ═══════════
s = slide()
bar(s, "v5 → v6 — 좋아진 것과 감수한 것 (정직한 결산)", "RESULT")
rows = [
    ["", "내용"],
    [("좋아진 것", {"color": GREEN, "bold": True}), "실행 루프에서 IK 소멸 — worst-case 지연이 상수(추론 ms + 필터 µs), \"IK 안 풀리면?\"이 원천 봉쇄"],
    ["", "특이점·브랜치 실패 소멸 — 44/52 구멍, 재시도, 표류가 데이터 빌드 시점에 해소됨"],
    ["", "손목 공회전 제거 — 태스크 무관 모터 동작 없음 (마모·케이블·한계 정지 리스크 제거)"],
    ["", "재생 무결성 100% — 159개 전 시연이 관절 궤적으로 완전 재현 (v5 런타임 IK 는 에피소드별 성공률 가변)"],
    ["", "실기 인터페이스 직결 — CS9 드라이버가 받는 것이 곧 모델 출력 (중간 계층 0)"],
    [("감수한 것", {"color": RED, "bold": True}), "TX90 전용 — OMX EE 호환 단절, 다른 로봇이면 변환부터 다시 (v5 노선을 병행 유지로 헤지)"],
    ["", "작업 위치가 데이터에 고정 — 테이프를 (−20,+15cm) 위치로 옮겨야 하며, 바꾸려면 재변환+재학습"],
    ["", "원 궤적과 미세 편차 — 평활+보간의 대가 (FK 위치 ≤5.1mm, 자세 ≤1.0°)"],
    ["", "연속성이 '보장'에서 '학습'으로 — 분포 밖 상황 대비 안전 필터(|Δj| 한계)가 실기 필수품이 됨"],
    ["", "state 가 합성값 — 완전 추종 가정(action[t-1]). 실기 전이학습 때는 진짜 측정 joint 로 교체됨"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[2.6, 15.4], size=11.5, rh=Cm(1.32))
foot(s, "결론: 실행 안정성·실시간성을 얻고, 범용성·위치 자유도를 지불 — 실기 배포 목적에는 남는 장사")

# ═══════════ 12. 생성 파일 ═══════════
s = slide()
bar(s, "오늘 생성 파일 전체 — 뭐가 뭔지", "FILES")
rows = [
    ["파일 (위치)", "설명"],
    ["convert_v6.py (컨테이너 /root)", "v6 변환기 본체 — 오프셋·평활·앵커 체인·DP 폴백·검증 3종 내장. --offset/--anchor/--smooth 인자"],
    ["scan_greedy.py / diag_ep8.py / diag_action.py", "오프셋 스윕 도구 / state·action 실패 해부 진단 스크립트"],
    ["knn_check.py / normalize_wind.py", "다봉성 kNN 검사 (정책 입력 공간) / 감김 정규화 시도(불필요 판명, 분석용 보존)"],
    ["assemble_v6.py / eval_v6.py / replay_v6.py", "데이터셋 조립(stats 재계산+영상 링크) / 학습 후 평가 / RViz 무IK 재생"],
    ["데이터셋: dlcodnjs/tx90_act_pick_and_place_v6_joint", "joint 표현 159ep. meta 에 conversion_meta.json(규약)·conversion_report.json(에피소드별 결과) 포함"],
    ["모델: /root/train_tx90_act_v6_joint/checkpoints/", "100k 학습 산출물, last = 최종. 설정: train_v6_config.json, 로그: train_v6.log"],
    ["예측: policy_rollouts/ep{000,080,140}_pred_v6joint.npz", "open-loop 예측 joint 궤적 (replay_v6.py 입력)"],
    ["plot_v5_v6_workspace.py (호스트 ~/tx90)", "v5/v6 전수 분석 그림 생성기"],
    ["omx_workspace/tx90_v5_workspace.png · tx90_v6_workspace.png", "데이터셋 v5(EE)·v6(joint) 4패널 분석 그림 — 계보 3부작 완성"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[7.2, 10.8], size=11, rh=Cm(1.52))
tf = textbox(s, Cm(0.9), Cm(17.0), Cm(32), Cm(1.4))
bullets(tf, [
    (0, "재현: python3 convert_v6.py --offset=-0.20,0.15,0 --anchor 5 --smooth 7 → assemble_v6.py → 학습 → eval_v6.py → replay_v6.py", {"mono": True, "size": 10.5}),
])
foot(s, "기존 산출물(데이터셋 v5·모델 v5·train_tx90_act_v5)은 일절 무변경")

# ═══════════ 13. 진행상황 + 다음 단계 ═══════════
s = slide()
bar(s, "진행상황과 다음 단계", "NEXT")
rows = [
    ["단계", "상태", "비고"],
    ["OMX 수집 → 변환 → v5 학습 → RViz 검증", ("완료", {"color": GREEN, "bold": True}), "~8/13"],
    ["v6 joint 데이터셋 + 모델 + 무IK 재생", ("완료", {"color": GREEN, "bold": True}), "8/14 — B안(v5+IK)·C안(v6) 모두 확보"],
    ["Stäubli 계정·velocity 확장 확인 + 대리점 견적", ("진행 필요", {"color": RED, "bold": True}), "사람 몫 — 오늘/내일 40분"],
    ["SRS CS9 에뮬레이터 검증", ("다음", {"color": RED}), "실기 없이 드라이버·스트리밍 한계 실측"],
    ["Isaac Sim: URDF 임포트 → v6 궤적 물리 재생", ("다음", {"color": RED}), "이번 주 목표 — 닫힌 루프의 첫 무대"],
    ["실기: VAL3 드라이버 open-loop 재현", ("예정", {"color": GREY}), "외부 제어 개통 후 — v6 가 바로 입력이 됨"],
    ["TX90 카메라 + 스크립트 수집 + 전이학습", ("예정", {"color": GREY}), "시각 도메인 갭 해결"],
    ["SmolVLA (재계획형 → 반응형)", ("예정", {"color": GREY}), "velocity 확장 결과에 따라 상한 결정"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[8.6, 2.6, 6.8], size=12, rh=Cm(1.35))
tf = textbox(s, Cm(0.9), Cm(14.6), Cm(32), Cm(3.4))
bullets(tf, [
    (0, "고려사항 (리스크 유지 목록)", {"bold": True, "size": 14}),
    (1, "시각 도메인 갭(OMX 영상 학습) — 실기 closed-loop 전 최우선 · CS9 스트리밍 미지수 — 에뮬레이터+견적으로 판가름"),
    (1, "청크 경계 30.5° 점프 1건 — 실기 전 안전 필터 구현 필수 · 테이프 이동(−20,+15) 실측 반영"),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-14")

# ═══════════ 14. 참고 링크 ═══════════
s = slide()
bar(s, "참고 링크", "REF")
tf = textbox(s, Cm(0.9), Cm(2.6), Cm(32), Cm(15))
bullets(tf, [
    (0, "외부 제어 (8/13 조사분 — 유효)", {"bold": True, "size": 15}),
    (1, "staubli_val3_driver (CS8/CS9 공식 ROS 드라이버)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver"}),
    (1, "Staubli_ROS2 (ROS2 포팅 + adaptive motion middleware)",
     {"link": "https://github.com/IvoD1998/Staubli_ROS2"}),
    (1, "adaptive_motion_control (velocity/pose tracking)",
     {"link": "https://github.com/FAU-FAPS/adaptive_motion_control"}),
    (1, "TX2-90 실시간 서보 논의 (issue #32)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver/issues/32"}),
    (0, "오늘 작업 관련", {"bold": True, "size": 15, "after": 6}),
    (1, "LeRobot — ACT 학습·데이터셋 포맷·비동기 추론 스택",
     {"link": "https://github.com/huggingface/lerobot"}),
    (1, "ACT 논문 (Zhao et al., 2023) — action chunking + CVAE",
     {"link": "https://arxiv.org/abs/2304.13705"}),
    (1, "MoveIt2 — IK(/compute_ik)·FK(/compute_fk) 서비스",
     {"link": "https://moveit.picknik.ai/"}),
    (1, "NVIDIA Isaac Sim — 다음 단계 물리 검증·닫힌 루프 실험",
     {"link": "https://developer.nvidia.com/isaac-sim"}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-14")

prs.save(os.path.join(OUTDIR, "20260814.pptx"))
print("저장:", os.path.join(OUTDIR, "20260814.pptx"))
