#!/usr/bin/env python3
"""2026-08-12 OMX → TX2-90 전이 작업 정리 PPT 생성.

  python3 make_report_pptx.py
  → /home/kim/tx90/OMX_TX90_전이_작업정리_20260812.pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ── 색 ──
NAVY = RGBColor(0x2F, 0x3D, 0x9E)
INK = RGBColor(0x15, 0x18, 0x1D)
GREY = RGBColor(0x62, 0x6B, 0x78)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
RED = RGBColor(0xA9, 0x33, 0x1D)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
BGSOFT = RGBColor(0xED, 0xEF, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "맑은 고딕"
MONO = "Consolas"

SW, SH = Cm(33.867), Cm(19.05)          # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size=14, bold=False, color=INK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide():
    return prs.slides.add_slide(BLANK)


def bar(s, title, step=None):
    """상단 타이틀 바."""
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Cm(2.0))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    tf = r.text_frame
    tf.margin_left = Cm(0.9); tf.margin_top = Cm(0.28)
    p = tf.paragraphs[0]
    if step:
        run = p.add_run(); run.text = step + "  "
        _set(run, 15, True, RGBColor(0xB9, 0xC2, 0xF0))
    run = p.add_run(); run.text = title
    _set(run, 21, True, WHITE)
    return r


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullets(tf, items, size=14):
    """items: (level, text, opts) — opts: bold/color/mono"""
    first = True
    for it in items:
        lv, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lv
        p.space_after = Pt(opt.get("after", 5))
        marks = {0: "▪  ", 1: "–  ", 2: "·  "}
        run = p.add_run()
        run.text = ("" if opt.get("nomark") else marks.get(lv, "")) + txt
        _set(run, opt.get("size", size - lv), opt.get("bold", False),
             opt.get("color", INK), MONO if opt.get("mono") else FONT)


def table(s, x, y, w, rows, widths=None, size=11.5, header=True, rh=Cm(0.72)):
    shp = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, rh * len(rows))
    t = shp.table
    if widths:
        total = sum(widths)
        for i, cw in enumerate(widths):
            t.columns[i].width = int(w * cw / total)
    for ri, row in enumerate(rows):
        t.rows[ri].height = rh
        for ci, cell in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Cm(0.15); c.margin_right = Cm(0.15)
            c.margin_top = Cm(0.04); c.margin_bottom = Cm(0.04)
            txt = cell if isinstance(cell, str) else cell[0]
            opt = {} if isinstance(cell, str) else cell[1]
            c.fill.solid()
            c.fill.fore_color.rgb = (NAVY if (header and ri == 0)
                                     else opt.get("bg", WHITE if ri % 2 else BGSOFT))
            p = c.text_frame.paragraphs[0]
            run = p.add_run(); run.text = txt
            _set(run, opt.get("size", size),
                 opt.get("bold", header and ri == 0),
                 WHITE if (header and ri == 0) else opt.get("color", INK),
                 MONO if opt.get("mono") else FONT)
    return t


def foot(s, txt):
    tf = textbox(s, Cm(0.9), SH - Cm(1.0), SW - Cm(1.8), Cm(0.7))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, 10, False, GREY)


# ═════════════════ 1. 표지 ═════════════════
s = slide()
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
r.fill.solid(); r.fill.fore_color.rgb = NAVY; r.line.fill.background()
tf = textbox(s, Cm(2.6), Cm(5.4), Cm(28), Cm(9))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "OMX → TX2-90 EE Pose 전이 파이프라인"
_set(run, 40, True, WHITE)
p = tf.add_paragraph(); p.space_before = Pt(14)
run = p.add_run(); run.text = "작업 정리 · 트러블슈팅 · 진행 상황"
_set(run, 22, False, RGBColor(0xB9, 0xC2, 0xF0))
p = tf.add_paragraph(); p.space_before = Pt(40)
run = p.add_run()
run.text = "2026. 08. 12  ·  한국생산기술연구원 현장실습  ·  이채원"
_set(run, 16, False, RGBColor(0xB9, 0xC2, 0xF0))
p = tf.add_paragraph(); p.space_before = Pt(6)
run = p.add_run()
run.text = "Umeyama similarity transform + MoveIt2 Cartesian + ACT 전이학습"
_set(run, 14, False, RGBColor(0x8A, 0x94, 0xC8))

# ═════════════════ 2. 왜 하는가 ═════════════════
s = slide(); bar(s, "이 작업을 왜 하는가", "배경")
tf = textbox(s, Cm(0.9), Cm(2.5), Cm(32), Cm(6))
bullets(tf, [
    (0, "목표: OpenManipulator-X(OMX)로 수집한 pick & place 시연 163개를 Staubli TX2-90 이 따라 하도록 전이", {"bold": True, "size": 15}),
    (1, "OMX(소형, 리치 0.31 m) 데이터 → TX2-90(산업용, 리치 1 m) 좌표계로 변환 후 전이학습(ACT)에 사용"),
    (0, "8/7~8/8 에 막혀 있던 문제", {"bold": True, "size": 15}),
    (1, "RViz Cartesian 재현 fraction 이 76.3% 에서 정체 — 경로의 24% 계획 실패"),
    (1, "움직이긴 하는데 pick & place 처럼 안 보임"),
    (1, "RViz 시작 시 모든 joint 가 0 → 지금 뭘 하는지 판단 불가"),
    (0, "기존 코드의 원인 진단 (실행가이드 PDF)", {"bold": True, "size": 15}),
    (1, "position: 근거 불명확한 손튜닝 affine (SCALE_X=1.0292 등)"),
    (1, "orientation: OMX 값을 좌표계 보정 없이 날것으로 사용 → 둘이 따로 놀아 재현 실패"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.9), Cm(13.7), Cm(32), Cm(3.6))
r.fill.solid(); r.fill.fore_color.rgb = BGSOFT; r.line.color.rgb = NAVY; r.line.width = Pt(1.5)
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.3); tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run(); run.text = "핵심 아이디어"
_set(run, 14, True, NAVY)
p = tf.add_paragraph()
run = p.add_run()
run.text = ("position 과 orientation 을 하나의 동일한 similarity transform T(s·R·t) 로 일관되게 변환한다. "
            "T 는 TX90 실물에서 측정한 4개 코너로 딱 한 번 계산하고 163개 에피소드 전체에 재사용한다. "
            "스케일 s 는 position 에만 적용하고 orientation 에는 곱하지 않는다.")
_set(run, 13)
foot(s, "결과 미리보기: 오늘 fraction 76.3% → 96.3%, all-zero 문제 해결, 학습 데이터셋 완성, 학습 파이프라인 검증 완료")

# ═════════════════ 3. 파이프라인 4개 층 ═════════════════
s = slide(); bar(s, "전체 파이프라인 — 4개 층", "구조")
layers = [
    ("층 1  캘리브레이션 (딱 한 번)", "OMX 4코너(데이터에서 추출) + TX90 4코너(펜던트 실측) → Umeyama → T(s, R, t) 저장", GREEN, "완료"),
    ("층 2  검증 (대표 에피소드)", "OMX EE → T 적용 → TX90 EE → MoveIt2 Cartesian → RViz 로 pick & place 재현 확인", GREEN, "완료"),
    ("층 3  배치 변환 (전체 자동)", "162개 에피소드 전체를 TX90 EE pose 로 변환해 저장 — IK 는 풀지 않음", GREEN, "완료"),
    ("층 4  전이학습", "TX90 EE pose 데이터셋으로 ACT 정책 학습 → 실행 시 MoveIt2 가 IK 를 그때그때 풂", RGBColor(0xC8, 0x8A, 0x00), "스모크 테스트 통과, 본 학습은 내일"),
]
y = Cm(2.7)
for title, desc, color, status in layers:
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.9), y, Cm(26), Cm(3.1))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.color.rgb = LINE; r.line.width = Pt(1.25)
    tf = r.text_frame; tf.margin_left = Cm(0.45); tf.margin_top = Cm(0.25); tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = title
    _set(run, 15, True, NAVY)
    p = tf.add_paragraph()
    run = p.add_run(); run.text = desc
    _set(run, 12.5)
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(27.3), y + Cm(0.75), Cm(5.7), Cm(1.6))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    tp = b.text_frame; tp.word_wrap = True
    tp.margin_top = Cm(0.05); tp.margin_left = Cm(0.1); tp.margin_right = Cm(0.1)
    pp = tp.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    run = pp.add_run(); run.text = status
    _set(run, 11.5, True, WHITE)
    y += Cm(3.55)
foot(s, "중요: RViz Cartesian 은 '변환이 맞는지 눈으로 보는 검증'이지 학습 데이터 생성이 아니다 — IK 를 미리 풀면 branch 문제가 163배가 된다")

# ═════════════════ 4. STEP 1~7 진행 현황 ═════════════════
s = slide(); bar(s, "STEP 1 ~ 7 진행 현황", "진행")
rows = [
    ["STEP", "내용", "상태", "핵심 결과"],
    ["1", "EE 데이터 구조 확인", "완료", "163 에피소드 · 7차원 [x,y,z,rx,ry,rz,grip] · FK 교차검증 오차 2×10⁻⁸ m"],
    ["2", "OMX 작업영역·4코너 추출", "완료", "241×361 mm (1:1.499) · 기준면 z=28.2 mm · 불량 ep62 발견"],
    ["3", "TX2-90 펜던트 실측", "완료", "4코너+HOME 측정 · 변 길이 오차 0.1 mm · z 편차 0.1 mm"],
    ["4", "Umeyama 변환식 T 계산", "완료", "s=0.999838 · R≈I · residual 0.02 mm (판정 기준 5 mm)"],
    ["5", "변환 코드 v2 작성", "완료", "T 적용 + tool frame 90° 보정 + HOME 반영 + 진단 기능"],
    ["6", "RViz 검증", "완료", "fraction 96.3% · all-zero 해결 · pick&place 패턴 그래프로 확인"],
    ["7", "배치 변환 + 학습 데이터셋", "완료", "전수·영상 검사로 불량 4건 제외 → 159 에피소드 76,345 프레임 · lerobot 로드 확인"],
    ["8", "ACT 전이학습", "내일", "스모크 테스트 통과 (loss 60.8→23.9) · 본 학습 약 3시간 예상"],
    ["9", "정책 실행 배선", "예정", "정책 출력(EE pose) → MoveIt2 IK → TX90 — v2 코드 재사용 예정"],
]
t = table(s, Cm(0.9), Cm(2.5), Cm(32), rows, widths=[1, 3.4, 1.1, 8.2], size=11.5)
for ri in (8, 9):
    for ci in range(4):
        c = t.cell(ri, ci)
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFA, 0xF0, 0xDC)
foot(s, "STEP 8 부터 내일 진행 — 명령어와 절차는 마지막 장에")

# ═════════════════ 5. 데이터 흐름 ═════════════════
s = slide(); bar(s, "데이터 흐름 한눈에", "구조")
flow = [
    ("OMX 원본 데이터셋", "joint 6차원 + 카메라 2대 영상\n(LeRobot v2.1, 163 ep, 1.2 GB)", False),
    ("FK (omx_f.urdf)", "joint → EE pose 7차원\n[x,y,z,rx,ry,rz,gripper]", False),
    ("Umeyama T(s,R,t)", "OMX 좌표계 → TX90 base_link\n+ tool frame Ry(+90°) 보정", True),
    ("TX90 EE 데이터셋", "162 ep · 77,608 프레임\n영상은 원본에 심볼릭 링크", False),
    ("ACT 학습 → 실행", "정책이 EE pose 출력\n→ MoveIt2 가 IK → TX90", False),
]
x = Cm(0.7)
for i, (t1, t2, hl) in enumerate(flow):
    w = Cm(6.0)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(4.2), w, Cm(4.4))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY if hl else WHITE
    r.line.color.rgb = NAVY; r.line.width = Pt(1.5)
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.25); tf.margin_right = Cm(0.25); tf.margin_top = Cm(0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = t1
    _set(run, 13, True, WHITE if hl else NAVY)
    for line in t2.split("\n"):
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(4)
        run = p.add_run(); run.text = line
        _set(run, 10.5, False, RGBColor(0xD5, 0xDA, 0xF5) if hl else INK)
    if i < 4:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Cm(0.02), Cm(5.9), Cm(0.62), Cm(0.9))
        a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background()
    x += w + Cm(0.66)
tf = textbox(s, Cm(0.9), Cm(9.6), Cm(32), Cm(7))
bullets(tf, [
    (0, "설계 원칙", {"bold": True, "size": 15}),
    (1, "T 는 실측 4점으로 한 번만 계산하고 전 에피소드에 재사용 — 에피소드별 튜닝 없음"),
    (1, "IK 는 데이터에 저장하지 않는다. joint 해는 실행 시점에 MoveIt2 가 현재 상태 기준으로 푼다"),
    (2, "미리 풀면 163개 에피소드마다 branch(팔꿈치 위/아래 등 다중해) 문제가 생김"),
    (1, "그리퍼 값(7번째 채널)은 좌표가 아니므로 변환 없이 그대로 통과"),
    (0, "오늘 사용한 도구", {"bold": True, "size": 15}),
    (1, "host: Python (numpy/scipy/pandas/matplotlib) — 데이터 처리·검증·그래프"),
    (1, "Docker physical_ai_server: ROS2 Jazzy + MoveIt2 + RViz + lerobot 0.2.0 (RTX 5070 Ti)"),
])

# ═════════════════ 6. 알고리즘 ① FK ═════════════════
s = slide(); bar(s, "정기구학 (Forward Kinematics)", "알고리즘 ①")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(14))
bullets(tf, [
    (0, "무엇: joint 각도 6개 → EE(그리퍼 끝) 위치·자세", {"bold": True}),
    (1, "URDF 의 관절 체인을 따라 4×4 동차변환 행렬을 곱해 나감"),
    (1, "T = Π ( T_origin(고정 오프셋) · T_joint(회전) )"),
    (0, "어디에 썼나", {"bold": True}),
    (1, "OMX: joint parquet → EE pose 재계산, EE 데이터셋과 대조"),
    (2, "전 프레임 오차 2×10⁻⁸ m — float32 한계 → 데이터 계보 확정"),
    (1, "TX90: 펜던트 joint 값 → FK → 펜던트 좌표와 대조"),
    (2, "→ 좌표계 원점 차이 478 mm 발견 (슬라이드 13)"),
    (0, "URDF 선택", {"bold": True}),
    (1, "omx_f.urdf 사용 — xacro 2개는 ROS 패키지 참조라 전개 불가"),
    (1, "EE 체인: world→link0→…→link5→end_effector_link"),
    (1, "gripper_joint_1 은 형제 브랜치 → EE pose 에 영향 없음"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.6), Cm(2.4), Cm(15.3), Cm(9.2))
r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.4); tf.word_wrap = True
code = [
    ("# fk_omx.py 핵심 (검증: 영자세)", GREY),
    ("chain = load_chain(urdf, 'world',", WHITE),
    ("                   'end_effector_link')", WHITE),
    ("T = fk(chain, [0,0,0,0,0])", WHITE),
    ("", WHITE),
    ("FK      = [0.31288, -0.0016, 0.21065]", RGBColor(0x7F, 0xE0, 0xA8)),
    ("URDF 합 = [0.31288, -0.0016, 0.21065]", RGBColor(0x7F, 0xE0, 0xA8)),
    ("# 소수점 5자리까지 일치", GREY),
    ("", WHITE),
    ("# EE 데이터셋 전 프레임 대조", GREY),
    ("pos 오차 2.0e-08 m  rpy 1.1e-06 rad", RGBColor(0x7F, 0xE0, 0xA8)),
]
first = True
for txt, col in code:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, 12.5, False, col, MONO)
tf = textbox(s, Cm(17.6), Cm(12.1), Cm(15.3), Cm(4.6))
bullets(tf, [
    (0, "왜 중요한가", {"bold": True, "color": NAVY}),
    (1, "EE 데이터셋이 어떤 URDF·어떤 규약(extrinsic xyz 오일러)으로 만들어졌는지 역으로 확정"),
    (1, "이 규약 확인이 없으면 이후 모든 회전 변환이 틀어짐"),
], size=13)

# ═════════════════ 7. 알고리즘 ② Umeyama ═════════════════
s = slide(); bar(s, "Umeyama Similarity Transform", "알고리즘 ②")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(14.2))
bullets(tf, [
    (0, "무엇: 대응점들로 두 좌표계 사이 s(스케일)·R(회전)·t(평행이동) 을 최소제곱 추정", {"bold": True}),
    (1, "pos_tx90 = s · (R @ pos_omx) + t", {"mono": True}),
    (1, "Kabsch(s=1 고정)와 달리 스케일까지 추정 — 두 로봇 작업영역 크기가 다를 가능성에 대비"),
    (0, "구현 (compute_transform.py)", {"bold": True}),
    (1, "교차공분산 C = X_dst^T X_src / n 을 SVD → R = U·S·V^T"),
    (1, "반사(reflection) 방지: det<0 이면 S 의 마지막 원소를 -1 로", {"color": RED}),
    (2, "이게 없으면 거울상 R 이 나와도 residual 0 — 로봇이 좌우 반대로 움직임"),
    (1, "s 는 position 에만 적용, orientation 은 R 만 (회전에 크기 개념 없음)"),
    (0, "합성 데이터 검증", {"bold": True}),
    (1, "노이즈 0: s·R·t 복원 오차 10⁻¹⁶ / 노이즈 1mm: residual 0.87mm"),
    (1, "몬테카를로 400회: 4점(동일평면)으로도 궤적 최대오차 1.7 mm → 5번째 점 불필요 판정"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.6), Cm(2.4), Cm(15.3), Cm(8.2))
r.fill.solid(); r.fill.fore_color.rgb = BGSOFT; r.line.color.rgb = NAVY; r.line.width = Pt(1.5)
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.35); tf.word_wrap = True
res = [
    ("실측 결과 (STEP 4)", NAVY, True, 14),
    ("s = 0.999838          (예상대로 ≈1)", INK, False, 13),
    ("R ≈ 단위행렬 (Z축 회전 0.01°)", INK, False, 13),
    ("t = [+0.4999, -0.0142, +0.3059] m", INK, False, 13),
    ("residual mean 0.02 mm / max 0.03 mm", GREEN, True, 13),
    ("", INK, False, 8),
    ("→ 사실상 순수 평행이동", NAVY, True, 13),
    ("두 사각형을 같은 크기로 잡았고 축도 나란히", GREY, False, 11.5),
]
first = True
for txt, col, b, sz in res:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, sz, b, col, MONO if txt.strip().startswith(("s =", "R ≈", "t =", "resid")) else FONT)
tf = textbox(s, Cm(17.6), Cm(11.0), Cm(15.3), Cm(5.6))
bullets(tf, [
    (0, "residual 의 올바른 해석", {"bold": True, "color": RED}),
    (1, "4점이 동일 평면 + 좌표 조그로 정확히 만든 사각형 → residual 은 '기하가 안 어긋났다'는 뜻이지 전이 정확도가 아님"),
    (1, "단위 착오(mm↔m)·코너 순서 오류 검출용 sanity check"),
    (1, "전이가 맞는지의 실제 검증은 STEP 6 RViz + 그래프"),
], size=12.5)

# ═════════════════ 8. 알고리즘 ③ tool frame ═════════════════
s = slide(); bar(s, "Tool Frame 90° 보정 — 그리퍼가 옆을 보던 진짜 이유", "알고리즘 ③")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(13.6))
bullets(tf, [
    (0, "두 로봇의 EE 접근축 규약이 다르다", {"bold": True}),
    (1, "OMX end_effector_link: 접근축 = 로컬 +X"),
    (2, "URDF: link5 에서 x 로 0.09193 m 나간 지점"),
    (1, "TX90 tool0: 접근축 = 로컬 +Z"),
    (2, "flange 의 Ry(-90°) 와 tool0 의 Ry(+90°) 가 상쇄 → tool0 = link_6 프레임, 접근축은 joint_6 축(+Z)"),
    (0, "보정 없이 쿼터니언을 그대로 넣으면", {"bold": True}),
    (1, "Umeyama R 이 완벽해도 그리퍼가 90° 돌아가 옆을 봄", {"color": RED}),
    (1, "기존 코드(v1)가 정확히 이 상태 — '그리퍼가 옆을 봄' 증상의 원인"),
    (0, "수정 (v2)", {"bold": True}),
    (1, "R_ee = R_umeyama @ R_omx @ Ry(+90°)", {"mono": True}),
    (1, "s 는 회전에 곱하지 않음"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.6), Cm(2.4), Cm(15.3), Cm(8.6))
r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.4); tf.word_wrap = True
code = [
    ("# ep0 첫 프레임으로 수치 확인", GREY),
    ("OMX 접근축(+X)을 월드로:", WHITE),
    ("  [ 0.083, -0.002, -0.997]  # 수직 아래", RGBColor(0x7F, 0xE0, 0xA8)),
    ("같은 값을 tool0 에 그대로:", WHITE),
    ("  [ 0.987,  0.138,  0.082]  # 수평 옆!", RGBColor(0xF0, 0x9A, 0x8A)),
    ("Ry(+90°) 보정 후:", WHITE),
    ("  [ 0.083, -0.002, -0.997]  # 일치", RGBColor(0x7F, 0xE0, 0xA8)),
    ("", WHITE),
    ("# dry run 전체 프레임 통계", GREY),
    ("보정 O: 아래를 보는 비율 97.9%", RGBColor(0x7F, 0xE0, 0xA8)),
    ("보정 X: 아래를 보는 비율  0.0%", RGBColor(0xF0, 0x9A, 0x8A)),
]
first = True
for txt, col in code:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, 12, False, col, MONO)
tf = textbox(s, Cm(17.6), Cm(11.4), Cm(15.3), Cm(5))
bullets(tf, [
    (0, "교훈", {"bold": True, "color": NAVY}),
    (1, "PDF 원안의 R_ee = R @ R_omx 에는 이 보정이 없었음 — 실행 전에 URDF 를 읽고 잡아낸 문제"),
    (1, "v2 에 --no_tool_fix 옵션을 둬서 언제든 전후 비교 가능"),
], size=13)

# ═════════════════ 9. 알고리즘 ④ 그리퍼 이벤트 ═════════════════
s = slide(); bar(s, "그리퍼 개폐 검출 — 데이터가 말하게 하기", "알고리즘 ④")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(13.6))
bullets(tf, [
    (0, "그리퍼 값(7번째 채널)의 분포가 깨끗한 이봉", {"bold": True}),
    (1, "닫힘 ≈ 0.23  /  열림 ≈ 0.69 → 임계값 0.459 (10/90 퍼센타일 중간)"),
    (0, "개폐 전이 = 이벤트", {"bold": True}),
    (1, "열림→닫힘 전이 = grasp(집기), 닫힘→열림 = release(놓기)"),
    (1, "163개 중 162개에서 grasp/release 검출 성공"),
    (0, "이 이벤트로 해결한 것들", {"bold": True}),
    (1, "실제 작업영역: grasp/release 지점만으로 사각형 → 241×361 mm"),
    (2, "전 프레임 percentile 방식은 HOME 뭉치 때문에 277 mm 로 왜곡"),
    (1, "pick & place 패턴 정량 검증 (슬라이드 16)"),
    (1, "불량 에피소드 검출: ep62 는 401 프레임 내내 그리퍼가 안 닫힘", {"color": RED}),
    (2, "물건을 안 집은 실패 시연 → 학습 데이터에서 제외"),
])
rows = [
    ["기준면(테이블) 후보", "z", "문제점"],
    ["전 프레임 1% percentile", "21.9 mm", "HOME 대기 자세가 하위 구간 독점 — 테이블 아님"],
    ["grasp/release 순간 평균", "57.4 mm", "물체를 '집은' 높이 — 궤적 17.7%가 기준면 아래로 최대 29mm 침투"],
    [("사각형 안 최저 도달 z (채택)", {"bold": True, "color": GREEN}), ("28.2 mm", {"bold": True, "color": GREEN}), ("기준면 아래 프레임 0개 — 테이블 표면에 가장 가까움", {"color": GREEN})],
]
table(s, Cm(17.6), Cm(2.6), Cm(15.3), rows, widths=[2.6, 1.1, 4.3], size=10.5, rh=Cm(1.35))
tf = textbox(s, Cm(17.6), Cm(9.0), Cm(15.3), Cm(7))
bullets(tf, [
    (0, "기준면을 3번 고쳐 잡은 과정", {"bold": True, "color": NAVY}),
    (1, "percentile → grasp 평균 → 최저 도달 z 순으로 개선"),
    (1, "각 단계마다 '이 값이면 실제 로봇이 어디를 짚는가'를 수치로 확인"),
    (1, "잘못된 기준면은 4코너가 같은 평면이라 t 에 흡수 — residual 로는 절대 안 잡히는 계통 오차", {"color": RED}),
], size=12.5)

# ═════════════════ 10. STEP 1 ═════════════════
s = slide(); bar(s, "EE 데이터 구조 확인", "STEP 1")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.4), Cm(13.5))
bullets(tf, [
    (0, "163개 parquet 전수 스캔 (inspect_ee_dataset.py)", {"bold": True}),
    (1, "컬럼 구성 고유 1가지 — 전 파일 동일 · NaN/Inf 없음"),
    (1, "observation.state / action = 7개 배열 (동일 구조)"),
    (1, "총 78,009 프레임 (에피소드당 282~601)"),
    (0, "인덱스 순서 — FK 교차검증으로 확정", {"bold": True}),
    (1, "[x, y, z, rx, ry, rz, gripper] · 단위 m / rad", {"mono": True}),
    (1, "오일러 = extrinsic xyz (scipy 'xyz') — 기존 코드와 동일 규약"),
    (0, "발견 ① 에피소드는 162개가 아니라 163개", {"bold": True, "color": RED}),
    (1, "episode_index 0~162 결번 없음 — PDF의 'for ep in 0..161' 그대로 쓰면 1개 누락"),
    (0, "발견 ② meta/info.json 이 낡음", {"bold": True, "color": RED}),
    (1, "아직 shape=[6], joint 이름 — parquet 실값과 불일치 → info.json 을 믿지 말고 실값 기준"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.9), Cm(2.4), Cm(15), Cm(9))
r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.4); tf.word_wrap = True
code = [
    ("# 데이터 계보 확인", GREY),
    ("원본 joint parquet (tx90 폴더)", WHITE),
    ("  └ FK(omx_f.urdf) 재계산", WHITE),
    ("     └ EE 데이터셋과 대조:", WHITE),
    ("", WHITE),
    ("pos  오차 max 2.0e-08 m", RGBColor(0x7F, 0xE0, 0xA8)),
    ("rpy  오차 max 1.1e-06 rad", RGBColor(0x7F, 0xE0, 0xA8)),
    ("grip 오차 0", RGBColor(0x7F, 0xE0, 0xA8)),
    ("", WHITE),
    ("# float32 저장 한계 = 완전 일치", GREY),
    ("# → _162_ee 는 omx_f.urdf FK 산물로 확정", GREY),
]
first = True
for txt, col in code:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, 12.5, False, col, MONO)
foot(s, "생성: inspect_ee_dataset.py, fk_omx.py — 이후 모든 단계의 규약(컬럼·오일러)이 여기서 확정됨")

# ═════════════════ 11. STEP 2 그림 ═════════════════
s = slide(); bar(s, "OMX 작업영역 · 4코너 추출", "STEP 2")
img = os.path.expanduser("~/omx_workspace/omx_workspace.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Cm(0.6), Cm(2.5), width=Cm(32.6))
tf = textbox(s, Cm(0.9), Cm(11.6), Cm(32), Cm(6.5))
bullets(tf, [
    (0, "왼쪽부터: XY 평면(사각형·grasp/release 점) · XZ 측면 · z 분포(기준면 확인) · 그리퍼 분포(개폐 판정)", {"size": 13}),
    (0, "결과: 사각형 241×361 mm (1:1.499) — grasp(빨강)는 y<0, release(파랑)는 y>0 로 명확히 두 구역", {"size": 13}),
    (1, "PDF 원안(전 프레임 1~99% percentile)은 241×277 (1:1.152) — HOME 뭉치가 왜곡, 세로 84mm 손실", {"color": RED, "size": 12.5}),
    (1, "종횡비가 다르면 TX90 사각형과 닮은꼴이 깨져 residual 40mm대로 상승 — 실측 전에 잡아서 재측정 회피", {"size": 12.5}),
    (0, "OMX 4코너 확정: P1[0.1119,-0.1596] P2[0.3527,-0.1596] P3[0.3527,0.2014] P4[0.1119,0.2014], z=0.0282 (m)", {"mono": True, "size": 11.5}),
])
foot(s, "생성: extract_omx_workspace.py → ~/omx_workspace/omx_workspace.png, omx_corners.txt")

# ═════════════════ 12. STEP 3 실측 ═════════════════
s = slide(); bar(s, "TX2-90 펜던트 실측", "STEP 3")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(13.8))
bullets(tf, [
    (0, "사전 결정 사항", {"bold": True}),
    (1, "TX90 에 그리퍼 없음 → flange(tool 없음) 기준 측정, link_name='tool0' 유지"),
    (1, "5번째 상단 점: 몬테카를로 결과 이득 미미(1.7→1.4mm) → 4점으로 확정"),
    (0, "측정 방법 — 테이프에 눈으로 맞추지 않기", {"bold": True}),
    (1, "P1 만 물리적으로 정하고, P2~P4 는 펜던트 좌표를 보며 World 축으로 정확히 241.0 / 361.0 mm 조그"),
    (1, "눈대중 정렬 오차 원천 제거 · z 는 자동으로 동일"),
    (0, "측정 품질", {"bold": True}),
    (1, "변 길이: 목표 [240.8, 361.0] → 실측 [240.7, 360.9] — 오차 0.1 mm", {"color": GREEN}),
    (1, "4코너 z 편차 0.1 mm — 테이블 수평 양호", {"color": GREEN}),
    (0, "HOME 자세", {"bold": True}),
    (1, "사각형 중앙, 테이블 높이 · joint [-3.31, 56.92, 73.58, -1.04, 48.85, -47.38]°"),
    (1, "tool0 접근축 [0.011, -0.014, -1.000] — 정확히 수직 아래 → 첫 구간 이동 156mm 로 짧음"),
])
rows = [
    ["점", "X (mm)", "Y (mm)", "Z (mm)"],
    ["P1", "611.84", "-173.72", "-143.89"],
    ["P2", "852.56", "-173.72", "-143.89"],
    ["P3", "852.53", "187.17", "-143.97"],
    ["P4", "611.75", "187.17", "-143.97"],
    ["HOME", "732.03", "6.42", "-144.03"],
]
table(s, Cm(17.9), Cm(2.6), Cm(15), rows, widths=[1, 1.3, 1.3, 1.3], size=12)
tf = textbox(s, Cm(17.9), Cm(8.0), Cm(15), Cm(8))
bullets(tf, [
    (0, "실측 카드(Artifact) 제공", {"bold": True, "color": NAVY}),
    (1, "배치도·체크리스트·기록지를 웹 페이지로 만들어 현장 휴대"),
    (0, "함정 예방 항목", {"bold": True, "color": NAVY}),
    (1, "World/Frame 모드 (Tool 모드 금지 — 과거 실수 이력)"),
    (1, "단위 명시 (mm/도) — 단위 착오는 residual 수십 cm 의 주범"),
    (1, "짧은 변 241 이 P1→P2 방향인지 대각선 434mm 로 확인"),
], size=12.5)
foot(s, "펜던트 측정값은 compute_transform.py 상단에 기록되어 있음")

# ═════════════════ 13. 발견: 478mm ═════════════════
s = slide(); bar(s, "펜던트 좌표계는 base_link 가 아니었다 (+478 mm)", "핵심 발견")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(13.8))
bullets(tf, [
    (0, "검증 방법: 측정한 5개 자세의 joint 값을 URDF 규약대로 FK → 펜던트 좌표와 대조", {"bold": True}),
    (0, "결과", {"bold": True}),
    (1, "X, Y: 5개 전부 0.1 mm 이내 일치"),
    (1, "Z: 5개 전부 정확히 478.0 mm 차이", {"color": RED, "bold": True}),
    (0, "원인 — URDF 에 답이 있었다", {"bold": True}),
    (1, "<joint name=\"base_link-base\"> origin xyz=\"0 0 0.478\"", {"mono": True}),
    (1, "펜던트 World 원점 = ROS-Industrial 'base' 링크(어깨 높이)"),
    (1, "MoveIt 요청 프레임 = base_link(바닥) → 478 mm 차이"),
    (0, "조치", {"bold": True}),
    (1, "compute_transform.py 에 PENDANT_Z_OFFSET = 0.478 적용"),
    (0, "부수 소득", {"bold": True, "color": GREEN}),
    (1, "X·Y 일치 = 펜던트 joint 규약이 URDF 와 동일 — 부호·영점 보정 불필요, deg→rad 만"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.9), Cm(2.4), Cm(15), Cm(8.6))
r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.4); tf.word_wrap = True
code = [
    ("      FK(base_link)   펜던트      차이", GREY),
    ("P1    z=+334.2       z=-143.89   478.0", WHITE),
    ("P2    z=+334.1       z=-143.89   478.0", WHITE),
    ("P3    z=+334.0       z=-143.97   478.0", WHITE),
    ("P4    z=+334.0       z=-143.97   478.0", WHITE),
    ("HOME  z=+334.0       z=-144.03   478.0", WHITE),
    ("", WHITE),
    ("# 랜덤 오차가 아니라 상수 → 좌표계 원점 차이", RGBColor(0xF0, 0x9A, 0x8A)),
]
first = True
for txt, col in code:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, 12, False, col, MONO)
tf = textbox(s, Cm(17.9), Cm(11.4), Cm(15), Cm(5))
bullets(tf, [
    (0, "놓쳤다면?", {"bold": True, "color": RED}),
    (1, "궤적 전체가 478 mm 아래 → 전부 도달 불가 → fraction 폭락"),
    (1, "joint 값을 같이 기록해온 덕에 실행 전에 검증으로 잡음 — 측정 시 여분 정보를 남기는 것의 가치"),
], size=13)

# ═════════════════ 14. STEP 5·6 트러블슈팅 ═════════════════
s = slide(); bar(s, "RViz 검증 트러블슈팅 — 4연쇄", "STEP 5·6")
rows = [
    ["#", "증상", "원인", "해결"],
    ["1", "fraction 84.3% 에서 '끊김' 경고",
     "에피소드 끝 정지 구간(같은 좌표 반복) — 경로 길이 0이면 planner 가 진행 못함",
     "정지 중복 waypoint 제거(dedup) → 47→38개, fraction 96.3%"],
    ["2", "로봇이 '고개 까닥'하고 끝",
     "compute_cartesian_path 는 시간 정보를 안 채움(전부 0초) → 164포인트가 3ms 에 재생",
     "원본 에피소드 길이(15.4초)로 시간 균등 배분"],
    ["3", "재생이 뚝뚝 끊김",
     "158포인트/14.8초 = 10.7 Hz 발행 — 영화 24fps 의 절반도 안 됨",
     "실시간 시계 기준 joint 선형보간, 50 Hz 발행"],
    ["4", "동작 중 all-zero 로 자꾸 튐",
     "demo.launch 의 joint_state_publisher 와 /joint_states 이중 발행 충돌",
     "source_list 토픽(fake_controller_joint_states)으로 발행 → 발행자 단일화"],
]
table(s, Cm(0.9), Cm(2.5), Cm(32), rows, widths=[0.5, 2.6, 4.6, 4.2], size=11, rh=Cm(1.9))
tf = textbox(s, Cm(0.9), Cm(11.2), Cm(32), Cm(6))
bullets(tf, [
    (0, "의미: 8/7~8/8 의 미해결 증상들이 전부 여기서 설명됨", {"bold": True, "size": 14}),
    (1, "'pick&place 처럼 안 보임' = #2 시간 미기입 (v1 도 동일 코드) — 변환 문제가 아니었음"),
    (1, "'RViz 시작 시 all-zero' = #4 joint_state_publisher 가 소스 없을 때 0 을 발행 — branch switching 아님"),
    (1, "'fraction 76.3% 정체' = 좌표계 문제(tool 90°, 손튜닝 affine) + #1 정지 구간의 복합"),
    (0, "최종: fraction 96.3%, 나머지 3.7% 는 시연 끝 잔여 정지 구간 — 동작은 사실상 100% 계획됨", {"bold": True, "size": 14, "color": GREEN}),
])
foot(s, "각 수정은 ee_to_tx90_cartesian_v2.py 에 반영 — --no_dedup, --rate, --joint_topic 옵션으로 전후 비교 가능")

# ═════════════════ 15. pick&place 검증 ═════════════════
s = slide(); bar(s, "pick & place 패턴 정량 검증", "STEP 6")
img = os.path.expanduser("~/omx_workspace/pickplace_check.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Cm(3.2), Cm(2.3), height=Cm(11.2))
rows = [
    ["ep", "집기", "놓기", "들어올림", "이동", "그리퍼↓"],
    ["0", "9.6s", "12.3s", "125 mm", "197 mm", "97.2%"],
    ["80", "2회", "2회", "87 mm", "181 mm", "78.8%"],
    ["161", "7.8s", "12.8s", "60 mm", "214 mm", "100%"],
]
table(s, Cm(0.9), Cm(14.0), Cm(19), rows, widths=[0.7, 1, 1, 1.2, 1, 1], size=11, rh=Cm(0.75))
tf = textbox(s, Cm(20.6), Cm(13.9), Cm(12.3), Cm(4.6))
bullets(tf, [
    (0, "세 에피소드 모두 내려감→집기→올림→놓기 구조 확인", {"bold": True, "size": 12.5, "color": GREEN}),
    (1, "RViz 에서 안 보인 이유: 들어올림이 60~125mm 로 작고 그리퍼가 없어 시각 변화 없음", {"size": 11.5}),
    (1, "ep80 은 재시도 시연(집기 2회) — 시연 자체 특성", {"size": 11.5}),
], size=12)
foot(s, "생성: check_pickplace.py → ~/omx_workspace/pickplace_check.png — 임의 에피소드 검증 가능 (--episodes 5 40 120)")

# ═════════════════ 16. STEP 7 + 데이터셋 조립 ═════════════════
s = slide(); bar(s, "배치 변환 + 학습용 데이터셋 조립", "STEP 7")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(14))
bullets(tf, [
    (0, "1차: EE pose 만 변환 (batch_transform.py)", {"bold": True}),
    (1, "163개 → ~/tx90_ee_dataset · sanity check 통과 (리치 초과 0, 바닥 아래 0)"),
    (1, "검증: RViz 에서 본 v2 변환과 비트 단위 일치 (오차 0.0)"),
    (0, "문제: 학습에는 카메라 영상이 필수", {"bold": True, "color": RED}),
    (1, "EE 발췌본에는 영상 없음 — 영상 없이 학습하면 평균 궤적만 외움"),
    (1, "진짜 원본(1.2 GB, mp4 326개)은 컨테이너 HF 캐시에서 발견"),
    (0, "2차: 완전한 LeRobot 데이터셋 조립 (build_tx90_dataset.py)", {"bold": True}),
    (1, "불량 에피소드 제외 후 0..N 재번호 — lerobot 은 결번을 허용 안 함"),
    (1, "영상 324개는 심볼릭 링크 — 내용 불변이므로 1.2 GB 복사 불필요"),
    (1, "episodes_stats 재계산 — joint(rad) 통계로 EE(m) 를 정규화하면 학습 붕괴", {"color": RED}),
    (1, "scipy 없는 컨테이너 대응: 오일러↔행렬 numpy 구현 (scipy 대조 10⁻¹⁶)"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.9), Cm(2.4), Cm(15), Cm(9.4))
r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.4); tf.word_wrap = True
code = [
    ("# lerobot 실제 로드 검증", GREY),
    ("LeRobotDataset(", WHITE),
    (" 'dlcodnjs/tx90_act_pick_and_place", WHITE),
    ("  _v4_162_ee')", WHITE),
    ("", WHITE),
    ("최종: 에피소드 159, 프레임 76,345", RGBColor(0x7F, 0xE0, 0xA8)),
    ("state/action: (7,) [x,y,z,rx,ry,rz,g]", RGBColor(0x7F, 0xE0, 0xA8)),
    ("camera1/2: (3,480,640) 디코딩 정상", RGBColor(0x7F, 0xE0, 0xA8)),
    ("", WHITE),
    ("# 심볼릭 링크 영상도 문제없이 로드", GREY),
]
first = True
for txt, col in code:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, 12.5, False, col, MONO)
tf = textbox(s, Cm(17.9), Cm(12.2), Cm(15), Cm(4.6))
bullets(tf, [
    (0, "알아둘 것: 시각 도메인 갭", {"bold": True, "color": RED}),
    (1, "영상 속 팔은 OMX, 배포 시 카메라가 볼 팔은 TX90"),
    (1, "실기 성능 저하 시 TX90 시연 몇 개로 파인튜닝이 정석 — 재촬영 없이는 회피 불가"),
], size=12.5)

# ═════════════════ 17. 학습 스모크 테스트 ═════════════════
s = slide(); bar(s, "ACT 학습 파이프라인 검증 (스모크 테스트)", "STEP 8 준비")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(16.2), Cm(13.8))
bullets(tf, [
    (0, "환경", {"bold": True}),
    (1, "lerobot 0.2.0 (physical_ai_tools 내장) · torch 2.7.0+cu128"),
    (1, "GPU: RTX 5070 Ti · CUDA 정상 인식"),
    (0, "ACT 정책 (기본 설정 유지)", {"bold": True}),
    (1, "chunk_size=100 (30fps 기준 3.3초치 동작을 한 번에 예측)"),
    (1, "backbone resnet18 · dim_model 512 · lr 1e-5 · VAE 사용"),
    (1, "학습 파라미터 52M"),
    (0, "10스텝 스모크 테스트 결과", {"bold": True}),
    (1, "loss 60.8 (step5) → 23.9 (step10) — 정상 하강", {"color": GREEN}),
    (1, "데이터 로드·정규화 통계·심볼릭 링크 영상·체크포인트 저장 전부 정상", {"color": GREEN}),
    (1, "스텝당 0.107초 → 10만 스텝 약 3시간 예상"),
])
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.9), Cm(2.4), Cm(15), Cm(10.6))
r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.4); tf.word_wrap = True
code = [
    ("# 내일 실행할 본 학습 명령 (컨테이너 안)", GREY),
    ("cd /root/ros2_ws/src/physical_ai_tools/\\", WHITE),
    ("   lerobot/src", WHITE),
    ("python3 -m lerobot.scripts.train \\", WHITE),
    ("  --dataset.repo_id=dlcodnjs/\\", WHITE),
    ("    tx90_act_pick_and_place_v4_162_ee \\", WHITE),
    ("  --policy.type=act \\", WHITE),
    ("  --policy.device=cuda \\", WHITE),
    ("  --policy.push_to_hub=false \\", WHITE),
    ("  --batch_size=8 --steps=100000 \\", WHITE),
    ("  --log_freq=200 --save_freq=10000 \\", WHITE),
    ("  --output_dir=/root/train_tx90_act \\", WHITE),
    ("  --job_name=tx90_act_ee", WHITE),
]
first = True
for txt, col in code:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run(); run.text = txt
    _set(run, 11.5, False, col, MONO)
tf = textbox(s, Cm(17.9), Cm(13.4), Cm(15), Cm(3.4))
bullets(tf, [
    (1, "GPU 메모리 여유 시 --batch_size=16 상향 가능"),
    (1, "loss 하강·gradient norm 안정 여부 관찰 / 체크포인트는 checkpoints/ 에 1만 스텝마다"),
], size=12.5)

# ═════════════════ 18. 생성 파일 ① ═════════════════
s = slide(); bar(s, "생성 파일 안내 ① — 스크립트 (/home/kim/tx90/)", "파일")
rows = [
    ["파일", "역할", "비고"],
    ["fk_omx.py", "OMX URDF 파싱 + FK. joint parquet → EE pose CSV 생성", "영자세 검증 포함"],
    ["inspect_ee_dataset.py", "STEP 1. 163개 parquet 전수 스캔 — 구조·범위·NaN 검사", "info.json 불신, 실값 기준"],
    ["extract_omx_workspace.py", "STEP 2. 작업영역 사각형·4코너·기준면 추출 + 그림", "grasp 이벤트 기반"],
    ["compute_transform.py", "STEP 4. Umeyama T 계산 + 사전점검(z편차·변길이·s범위)", "실측값·478mm 보정 포함"],
    ["ee_to_tx90_cartesian_v2.py", "STEP 5·6. 변환→MoveIt2 Cartesian→RViz 재생 (v1 대체)", "컨테이너 /root/ 에도 복사됨"],
    ["check_pickplace.py", "pick&place 패턴 정량 검증 그래프 (임의 에피소드)", "ROS 불필요, host 실행"],
    ["summarize_all_episodes.py", "162개 전수 검사 — 품질 분포·이상치 그래프", "품질 지도 포함"],
    ["batch_transform.py", "STEP 7. EE pose 만 일괄 변환 (영상 없는 버전)", "~/tx90_ee_dataset 생성"],
    ["build_tx90_dataset.py", "학습용 완전체 조립 — 영상 링크·재번호·통계 재계산", "컨테이너 안에서 실행"],
    ["export_transcript.py", "Claude Code 대화 기록 → 마크다운 전문 추출", "claude.ai 이관용"],
    ["make_report_pptx.py", "이 PPT 를 생성한 스크립트", "수정 후 재생성 가능"],
]
table(s, Cm(0.9), Cm(2.5), Cm(32), rows, widths=[2.5, 6.2, 2.6], size=10.5, rh=Cm(1.15))
foot(s, "모든 스크립트는 인자 없이 기본값으로 동작 · 상단 docstring 에 사용법 기재")

# ═════════════════ 19. 생성 파일 ② ═════════════════
s = slide(); bar(s, "생성 파일 안내 ② — 데이터·결과물", "파일")
rows = [
    ["위치", "내용", "설명"],
    ["~/omx_workspace/", "omx_corners.txt", "OMX 4코너 좌표 (compute_transform 에 붙여넣는 형식)"],
    ["", "omx_workspace.png", "작업영역 4패널 그림 (STEP 2 결과)"],
    ["", "transform.npy", "T(s,R,t) + residual 기록 — 전 단계가 참조하는 핵심 산출물"],
    ["", "pickplace_check.png", "ep 0/80/161 pick&place 검증 3×3 그래프"],
    ["", "all_episodes_summary.png", "162개 전수 검사 6패널 그래프"],
    ["~/tx90_ee_dataset/", "parquet 163개 + meta/", "EE pose 만 변환한 중간 산출물 (영상 없음)"],
    ["~/tx90/", "대화전문_OMX_TX90.md", "오늘 작업 대화 전문 (근거·수치 포함, claude.ai 업로드용)"],
    ["", "OMX_TX90_전이_작업정리_....pptx", "이 발표자료"],
    ["컨테이너 /root/", "transform.npy · v2 · build_...", "host 와 동일본 (docker cp 로 동기화됨)"],
    ["컨테이너 HF 캐시", "tx90_act_pick_and_place_v4_162_ee", "학습용 최종 데이터셋 — 159 ep, 영상 심볼릭 링크"],
    ["HF 캐시/frames_check/", "영상 프레임 시트 PNG 13장", "영상 재생 불가 환경용 — ep111/115 실패·ep0 카메라 증거"],
]
table(s, Cm(0.9), Cm(2.5), Cm(32), rows, widths=[2.6, 3.6, 6.2], size=10.5, rh=Cm(1.1))
tf = textbox(s, Cm(0.9), Cm(15.9), Cm(32), Cm(2.4))
bullets(tf, [
    (1, "웹 실측 카드(Artifact): https://claude.ai/code/artifact/8f08108c-46c0-43cb-b661-165f83ce12f9 — 재측정 시 재사용", {"size": 12}),
    (1, "주의: 학습 데이터셋의 영상은 원본(omx_act_pick_and_place_v4_162)에 링크 — 원본 캐시를 지우면 영상이 깨짐", {"size": 12, "color": RED}),
])

# ═════════════════ 20. 현재 위치 / 내일 ═════════════════
s = slide(); bar(s, "현재 위치와 내일 할 일", "정리")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(15.8), Cm(13.5))
bullets(tf, [
    (0, "오늘 완료 (층 1~3 + 층 4 준비)", {"bold": True, "size": 15, "color": GREEN}),
    (1, "캘리브레이션 T 확정 (residual 0.02 mm)"),
    (1, "RViz 검증 통과 — fraction 96.3%, 미해결 증상 4건 전부 원인 규명"),
    (1, "학습용 데이터셋 완성 + lerobot 로드 확인"),
    (1, "학습 파이프라인 10스텝 검증 (loss 정상 하강)"),
    (0, "예방한 사고들", {"bold": True, "size": 15}),
    (1, "펜던트 z +478 mm (전 궤적 도달 불가가 될 뻔)"),
    (1, "tool frame 90° (그리퍼가 옆을 볼 뻔)"),
    (1, "기준면 29 mm 오차 (테이블 침투할 뻔)"),
    (1, "통계 미갱신 (정규화 붕괴로 학습이 안 될 뻔)"),
    (1, "불량 4건 검출·제외 — ep0 카메라 뒤바뀜, ep62 그리퍼 미작동, ep111·115 운반 실패(영상 확정)"),
])
tf = textbox(s, Cm(17.3), Cm(2.4), Cm(15.7), Cm(13.5))
bullets(tf, [
    (0, "내일 ① — 본 학습 (약 3시간, GPU 방치 가능)", {"bold": True, "size": 15, "color": NAVY}),
    (1, "슬라이드 17 의 명령 실행 → loss 곡선 확인"),
    (0, "내일 ② — 정책 실행 배선", {"bold": True, "size": 15, "color": NAVY}),
    (1, "카메라 2대 + 현재 EE pose → ACT → 목표 EE pose → MoveIt2 IK → TX90"),
    (1, "v2 코드의 'EE pose → MoveIt' 부분 재사용 — parquet 입력을 정책 출력으로 교체"),
    (0, "그다음 (실기 단계)", {"bold": True, "size": 15}),
    (1, "TX90 에 그리퍼 장착 → URDF tool link 추가 → link_name 만 변경 (T 재계산 불필요)"),
    (1, "테이프로 표시해 둔 사각형 자리에 작업판 배치"),
    (1, "실기 성능 미흡 시: TX90 시연 소량 수집 → 파인튜닝 (도메인 갭 대응)"),
    (0, "김박사님께 보고할 것", {"bold": True, "size": 15}),
    (1, "residual 0.02 mm · fraction 96.3% · 데이터셋 159 ep 확정(불량 4건 영상 근거로 제외)"),
    (1, "펜던트 좌표계 +478 mm 발견 건 공유"),
])
foot(s, "이어서 할 때: Claude Code 를 /home/kim/tx90 에서 열면 오늘 결정사항이 메모리에 저장되어 있어 맥락이 복원됨")


# ═════════════════ A. 162개 전수 검사 ═════════════════
s = slide(); bar(s, "162개 전수 검사 — 품질 분포와 이상치", "검증 ①")
img = os.path.expanduser("~/omx_workspace/all_episodes_summary.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Cm(0.6), Cm(2.4), width=Cm(23.8))
tf = textbox(s, Cm(24.9), Cm(2.5), Cm(8.4), Cm(14.5))
bullets(tf, [
    (0, "판정", {"bold": True, "color": NAVY}),
    (1, "깨끗한 1집기-1놓기 149개 (91%)"),
    (1, "재시도 13개 — 회복 행동 학습에 유익, 유지"),
    (1, "의심 2개 (ep111·115) — 품질 지도에서 유일하게 고립 → 영상 검사로 이관"),
    (0, "전체 평균 (± 표준편차)", {"bold": True, "color": NAVY}),
    (1, "집기 높이 29±15 mm (테이블 위)"),
    (1, "놓기 높이 30±15 mm — 공중 투하 아님"),
    (1, "들어올림 111±35 mm"),
    (1, "운반 거리 197±45 mm"),
    (1, "집기 시점 36% / 놓기 64% (진행률)"),
    (0, "그래프 6패널", {"bold": True, "color": NAVY}),
    (1, "① z 곡선 겹침(공통 구조) ② 집기·놓기 위치 ③ 높이 분포 ④ 들어올림 ⑤ 운반 ⑥ 품질 지도"),
], size=11)
foot(s, "생성: summarize_all_episodes.py → ~/omx_workspace/all_episodes_summary.png")

# ═════════════════ B1. 영상 검사 — ep111/115 ═════════════════
s = slide(); bar(s, "영상 프레임 검사 ① — ep111·115 는 실패 시연으로 확정", "검증 ②")
tf = textbox(s, Cm(0.9), Cm(2.3), Cm(32), Cm(4.0))
bullets(tf, [
    (0, "방법 — host 에서 영상이 안 열려(ffmpeg 부재·opencv 손상) 컨테이너의 pyav 로 균등 간격 10프레임을 추출해 PNG 시트로 저장"),
    (1, "frames_check/ 폴더 — PNG 라 코덱 없이 어디서든 열람 가능. 정상 기준(ep0)과 나란히 비교"),
    (0, "판정: 두 에피소드 모두 파란 주사위를 집으려다 실패 — 주사위가 끝까지 처음 자리에 그대로", {"bold": True, "color": RED}),
    (1, "통계 이상과 정확히 일치: 집기높이 63mm(정상 29) · 수평이동 15mm(정상 197) · 놓기높이 101/123mm(정상 30)"),
])
img = os.path.expanduser("~/physical_ai_tools/docker/huggingface/frames_check/ep111_camera2.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Cm(0.9), Cm(6.8), width=Cm(32))
tf = textbox(s, Cm(0.9), Cm(16.6), Cm(32), Cm(1.8))
bullets(tf, [
    (1, "ep111 측면 카메라 10프레임 (좌→우, 상→하): 팔이 왼쪽 파란 주사위로 접근·집기 시도하지만 주사위는 제자리. ep115 도 동일 패턴", {"size": 12}),
])

# ═════════════════ B2. 영상 검사 — ep0 카메라 ═════════════════
s = slide(); bar(s, "영상 프레임 검사 ② — ep0 만 카메라가 뒤바뀜", "검증 ③")
img = os.path.expanduser("~/physical_ai_tools/docker/huggingface/frames_check/ep012_check.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Cm(0.7), Cm(2.6), height=Cm(13.4))
tf = textbox(s, Cm(15.2), Cm(2.5), Cm(17.7), Cm(14.5))
bullets(tf, [
    (0, "발견 — 11개 에피소드 카메라 조사 중 ep0 만 다름", {"bold": True, "color": RED}),
    (1, "ep0: camera1 = 측면 전경, camera2 = 회색 블러(고장·가려진 손목캠)"),
    (1, "ep1~162: camera1 = 손목 근접(그리퍼 시점), camera2 = 측면 전경 — 전부 일관"),
    (1, "판별 근거: camera2 중간 프레임 표준편차 ep0=46 vs 정상 75~84"),
    (0, "왜 문제인가", {"bold": True}),
    (1, "ACT 는 \"camera1 은 항상 손목 시점\" 으로 배우는데 ep0 만 반대 → 시각 입력 오염"),
    (1, "블러 손목캠은 대체 불가(영상 자체가 없음) → ep0 제외"),
    (0, "부수 확정 ① 카메라 의미", {"bold": True, "color": GREEN}),
    (1, "camera1 = 손목, camera2 = 측면 — 실기 때 TX90 에 같은 구성으로 설치해야 함"),
    (0, "부수 확정 ② 데이터는 두 배치", {"bold": True, "color": GREEN}),
    (1, "task0(v3_90) = ep0~122, task1(v3_add) = ep123~162 (주사위 색·개수 다름)"),
    (1, "같은 pick&place 과제라 함께 학습해도 무방"),
], size=12.5)
foot(s, "왼쪽 그림: ep0/1/2 의 camera1(좌)·camera2(우) 중간 프레임 — ep0 만 좌우가 반대이고 cam2 가 블러")

# ═════════════════ C. 최종 데이터셋 + 실기 주의 ═════════════════
s = slide(); bar(s, "최종 학습 데이터셋 확정 (162 → 159) + 실기 주의사항", "검증 ④")
rows = [
    ["ep", "제외 사유", "근거"],
    ["0", "카메라 구성 상이", "영상: cam1↔cam2 반대, 손목캠 블러"],
    ["62", "그리퍼 미작동", "401프레임 내내 0.663 고정 (물체 안 집음)"],
    ["111", "운반 실패", "영상: 주사위 제자리 · 이동 15mm"],
    ["115", "운반 실패", "영상: 주사위 제자리 · 이동 15mm"],
]
table(s, Cm(0.9), Cm(2.6), Cm(15.8), rows, widths=[0.8, 2.2, 4.4], size=11, rh=Cm(0.95))
tf = textbox(s, Cm(0.9), Cm(8.2), Cm(15.8), Cm(8.5))
bullets(tf, [
    (0, "재조립", {"bold": True, "color": NAVY}),
    (1, "build_tx90_dataset.py --drop 0 62 111 115", {"mono": True}),
    (1, "결과: 159 에피소드 · 76,345 프레임 · lerobot 로드 재검증 통과"),
    (1, "주의: 재조립 전 출력 폴더 삭제 필수 — 이전 빌드의 높은 번호 파일이 남으면 데이터 오염", {"color": RED}),
    (1, "손실 4/163 = 2.5% — 학습에 영향 없는 수준"),
], size=12.5)
r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17.4), Cm(2.6), Cm(15.5), Cm(13.9))
r.fill.solid(); r.fill.fore_color.rgb = BGSOFT; r.line.color.rgb = RED; r.line.width = Pt(1.5)
tf = r.text_frame; tf.margin_left = Cm(0.5); tf.margin_top = Cm(0.35); tf.word_wrap = True
p = tf.paragraphs[0]; run = p.add_run(); run.text = "실기 주의 — 테이블면 아래 접근 구간"
_set(run, 14, True, RED)
for txt in (
    "전 프레임의 3.5%(2,723개)가 실측 테이블면(334.1mm)보다 낮게 지나감 — 최대 23.4mm 아래",
    "전부 실측 사각형 밖, 로봇~사각형 사이 접근 통로 (x 539~612mm 구간)",
    "OMX 리그에는 그 구간에 테이블이 없었다는 뜻 (HOME 접근 경로)",
    "",
    "→ 실기 전 확인: TX90 실제 테이블 앞 모서리가 테이프 사각형보다 로봇 쪽으로 나와 있으면 flange 가 최대 23mm 파고듦",
    "→ 대응: 테이블 위치 확인 후 필요 시 MoveIt 에 테이블 충돌 객체 추가",
):
    p = tf.add_paragraph(); p.space_after = Pt(6)
    run = p.add_run(); run.text = txt
    _set(run, 12.5)
foot(s, "RViz 검증에는 영향 없음(장면에 테이블 없음) — 실기·실행 단계에서만 해당")

# ═════════════════ D. 기존 모델과의 차이 ═════════════════
s = slide(); bar(s, "기존 OMX 모델과 무엇이 다른가 + 실행 구조", "학습 이해")
rows = [
    ["", "기존 모델 (기학습)", "이번 모델 (내일)"],
    ["학습 데이터", "OMX joint 6차원", "TX90 EE pose 7차원 (오늘 변환본)"],
    ["정책 출력", "\"관절을 이 각도로 꺾어라\"", "\"그리퍼 끝을 이 위치·자세로\""],
    ["실행 가능 로봇", "OMX 뿐 (몸에 종속된 언어)", "TX90 — MoveIt IK 가 관절로 번역"],
    ["카메라·시연", "동일 (같은 영상 163개)", "동일 (불량 4건만 제외)"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32), rows, widths=[1.6, 3.2, 3.6], size=12, rh=Cm(1.05))
flow = [
    ("카메라 2대", "손목 + 측면\n30fps", False),
    ("ACT 정책", "PC·GPU 에서 실행\n(로봇에 안 올라감)", True),
    ("목표 EE pose", "다음 3.3초\n100프레임 chunk", False),
    ("MoveIt2 IK", "실행 순간마다\njoint 해 계산", False),
    ("TX90", "joint 명령만\n받아 실행", False),
]
x = Cm(0.7)
for i, (t1, t2, hl) in enumerate(flow):
    w = Cm(6.0)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(9.6), w, Cm(3.6))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY if hl else WHITE
    r.line.color.rgb = NAVY; r.line.width = Pt(1.5)
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.2); tf.margin_right = Cm(0.2); tf.margin_top = Cm(0.25)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = t1
    _set(run, 13, True, WHITE if hl else NAVY)
    for line in t2.split("\n"):
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(3)
        run = p.add_run(); run.text = line
        _set(run, 10.5, False, RGBColor(0xD5, 0xDA, 0xF5) if hl else INK)
    if i < 4:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Cm(0.02), Cm(11.0), Cm(0.62), Cm(0.9))
        a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background()
    x += w + Cm(0.66)
tf = textbox(s, Cm(0.9), Cm(13.9), Cm(32), Cm(4.2))
bullets(tf, [
    (0, "핵심", {"bold": True, "size": 14}),
    (1, "joint 각도는 로봇 몸에 묶인 언어라 로봇 간 이식 불가 — 두 로봇이 공유하는 언어(EE pose)로 바꾼 것이 오늘 작업의 목적"),
    (1, "모델은 로봇 컨트롤러에 올라가지 않는다. PC 에서 30fps 제어 루프로 돌며 TX90 은 joint 명령만 수신"),
    (1, "IK 를 학습 데이터에 미리 넣지 않는 이유: 163개를 미리 풀면 branch(팔꿈치 다중해) 문제가 163배"),
    (1, "기존 모델의 의미: 같은 시연·카메라로 ACT 가 이 과제를 배울 수 있음을 이미 증명 — 이번엔 행동 표현만 교체", {"color": GREEN}),
])
foot(s, "실기 추가 요건: TX90 그리퍼 장착 + 카메라 2대(손목·측면) 설치 · 영상 속 팔이 OMX 인 도메인 갭은 성능 미흡 시 파인튜닝으로 대응")

# ═════════════════ 순서 재배치 ═════════════════
def move_slide(old, new):
    lst = prs.slides._sldIdLst
    el = list(lst)[old]
    lst.remove(el)
    lst.insert(new, el)

for _old, _new in ((20, 16), (21, 17), (22, 18), (23, 19), (24, 21)):
    move_slide(_old, _new)

out = "/home/kim/tx90/OMX_TX90_전이_작업정리_20260812.pptx"
prs.save(out)
print(f"저장: {out}")
print(f"슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장")
