#!/usr/bin/env python3
"""2026-08-13 작업 정리 PPT 생성 (어제 20260812.pptx 이후 → 오늘 전체).

  python3 make_report_pptx_0813.py
  → /home/kim/tx90/대화록 및 PPT/20260813.pptx

내용: v5 재학습 → RViz 검증 → IK 체인 → 특이점 트러블슈팅(알고리즘 4종 비교)
      → 외부 제어 조사 → 아키텍처 결정 → 생성 파일 → 다음 단계 → 고려사항
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ── 색 (20260812.pptx 와 동일) ──
NAVY = RGBColor(0x2F, 0x3D, 0x9E)
INK = RGBColor(0x15, 0x18, 0x1D)
GREY = RGBColor(0x62, 0x6B, 0x78)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
RED = RGBColor(0xA9, 0x33, 0x1D)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
BGSOFT = RGBColor(0xED, 0xEF, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LILAC = RGBColor(0xB9, 0xC2, 0xF0)

FONT = "맑은 고딕"
MONO = "Consolas"

SW, SH = Cm(33.867), Cm(19.05)
BASE = "/home/kim/tx90"
OUTDIR = os.path.join(BASE, "대화록 및 PPT")
IMG = os.path.join(OUTDIR, "img_0813")

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size=14, bold=False, color=INK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide():
    return prs.slides.add_slide(BLANK)


def bar(s, title, step=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Cm(2.0))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    tf = r.text_frame
    tf.margin_left = Cm(0.9); tf.margin_top = Cm(0.28)
    p = tf.paragraphs[0]
    if step:
        run = p.add_run(); run.text = step + "  "
        _set(run, 15, True, LILAC)
    run = p.add_run(); run.text = title
    _set(run, 21, True, WHITE)
    return r


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullets(tf, items, size=14):
    """items: (level, text, opts) — opts: bold/color/mono/link/after/size/nomark"""
    first = True
    for it in items:
        lv, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lv
        p.space_after = Pt(opt.get("after", 5))
        marks = {0: "▪  ", 1: "–  ", 2: "·  "}
        run = p.add_run()
        run.text = ("" if opt.get("nomark") else marks.get(lv, "")) + txt
        _set(run, opt.get("size", size - lv), opt.get("bold", False),
             opt.get("color", INK), MONO if opt.get("mono") else FONT)
        if opt.get("link"):
            run.hyperlink.address = opt["link"]
            run.font.color.rgb = NAVY


def table(s, x, y, w, rows, widths=None, size=11.5, header=True, rh=Cm(0.72)):
    shp = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, rh * len(rows))
    t = shp.table
    if widths:
        total = sum(widths)
        for i, cw in enumerate(widths):
            t.columns[i].width = int(w * cw / total)
    for ri, row in enumerate(rows):
        t.rows[ri].height = rh
        for ci, cell in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Cm(0.15); c.margin_right = Cm(0.15)
            c.margin_top = Cm(0.04); c.margin_bottom = Cm(0.04)
            txt = cell if isinstance(cell, str) else cell[0]
            opt = {} if isinstance(cell, str) else cell[1]
            c.fill.solid()
            c.fill.fore_color.rgb = (NAVY if (header and ri == 0)
                                     else opt.get("bg", WHITE if ri % 2 else BGSOFT))
            p = c.text_frame.paragraphs[0]
            run = p.add_run(); run.text = txt
            _set(run, opt.get("size", size),
                 opt.get("bold", header and ri == 0),
                 WHITE if (header and ri == 0) else opt.get("color", INK),
                 MONO if opt.get("mono") else FONT)
    return t


def foot(s, txt):
    tf = textbox(s, Cm(0.9), SH - Cm(1.0), SW - Cm(1.8), Cm(0.7))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = txt
    _set(run, 10, False, GREY)


def flowbox(s, x, y, w, h, txt, fill=NAVY, fg=WHITE, size=12, bold=True,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    b = s.shapes.add_shape(shape, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = LINE; b.line.width = Pt(0.75)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.06); tf.margin_bottom = Cm(0.06)
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = line
        _set(run, size if i == 0 else size - 2, bold if i == 0 else False, fg)
    return b


def arrow(s, x, y, w=Cm(0.7)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Cm(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = GREY
    a.line.fill.background()
    return a


# ════════════════════ 1. 표지 ════════════════════
s = slide()
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
r.fill.solid(); r.fill.fore_color.rgb = NAVY; r.line.fill.background()
tf = textbox(s, Cm(2.2), Cm(5.2), Cm(29.5), Cm(6))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "OMX → TX2-90 모방학습 전이"
_set(run, 36, True, WHITE)
p = tf.add_paragraph(); p.space_before = Pt(6)
run = p.add_run(); run.text = "2026-08-13 작업 정리 — v5 재학습 · RViz 검증 · 특이점 해결 · 실기 로드맵"
_set(run, 19, False, LILAC)
p = tf.add_paragraph(); p.space_before = Pt(28)
run = p.add_run()
run.text = "자세오차 28.9° → 1.8°   |   IK 성공률 84.6% → 100%   |   외부 제어 경로 확보"
_set(run, 16, True, WHITE)
foot_tf = textbox(s, Cm(2.2), SH - Cm(1.6), Cm(20), Cm(0.8))
run = foot_tf.paragraphs[0].add_run()
run.text = "이채원 · 2026-08-13"
_set(run, 12, False, LILAC)

# ════════════════════ 2. 전체 파이프라인 (오늘 위치) ════════════════════
s = slide()
bar(s, "전체 파이프라인 — 오늘 작업의 위치", "PIPELINE")
y1 = Cm(3.2); h = Cm(2.1); w = Cm(4.35); gap = Cm(0.72)
steps = [
    ("OMX 시연 수집\n162 에피소드 (완료)", GREEN),
    ("FK + 좌표 변환\nUmeyama+툴 보정 (완료)", GREEN),
    ("TX90 데이터셋\nv4→v5 재조립 (오늘)", RED),
    ("ACT 학습\n100k 스텝 (오늘)", RED),
    ("RViz 검증\nIK 체인 재생 (오늘)", RED),
    ("실기 TX90\n외부제어+수집 (예정)", GREY),
]
x = Cm(0.9)
for i, (txt, c) in enumerate(steps):
    flowbox(s, x, y1, w, h, txt, fill=c)
    if i < len(steps) - 1:
        arrow(s, x + w + Cm(0.03), y1 + h / 2 - Cm(0.25), Cm(0.62))
    x += w + gap
tf = textbox(s, Cm(0.9), Cm(6.1), Cm(32), Cm(1.0))
bullets(tf, [(0, "초록 = 완료(~8/12)   빨강 = 오늘(8/13) 작업   회색 = 예정", {"color": GREY, "size": 13})])
tf = textbox(s, Cm(0.9), Cm(7.2), Cm(32), Cm(10.5))
bullets(tf, [
    (0, "프로젝트 핵심 아이디어", {"bold": True, "size": 16}),
    (1, "OMX(4축)와 TX90(6축)은 관절 공간이 달라 직접 전이 불가 → EE pose(6D)+그리퍼(1D)로 동작 공간 통일"),
    (1, "OMX 시연의 EE 궤적을 Umeyama 변환 + 툴 프레임 보정(Ry+90°)으로 TX90 base_link 좌표계로 이식"),
    (1, "정책(ACT)은 카메라 2대 영상 + state(7D) → action(7D) 청크 100스텝을 출력"),
    (0, "오늘의 질문: \"학습된 정책이 TX90 기구학으로 실제 실행 가능한 궤적을 내는가?\"", {"bold": True, "size": 16, "after": 8}),
    (1, "답: 가능. 단, 손목 특이점 회피를 위해 작업 위치를 로봇 쪽으로 10cm 이동해야 함 (오늘 실증)"),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 3. 오늘 작업 타임라인 ════════════════════
s = slide()
bar(s, "오늘 작업 흐름 (시간순)", "PIPELINE")
rows = [
    ["시각", "작업", "결과"],
    ["~10:04", "v4 데이터셋으로 학습 시작 → 10k 시점 중간 평가", "자세오차 28.9° — 비정상, 원인 추적 시작"],
    ["10:14", "원인 발견: 오일러 각 wrap 결함 → fix_euler_wrap.py로 수정, v5 재조립", "tx90_act_pick_and_place_v5_ee 생성"],
    ["10:33~13:03", "v5로 ACT 재학습 100,000 스텝 (약 2.5시간, 10k마다 체크포인트)", "train_tx90_act_v5/checkpoints/last"],
    ["13:03", "open-loop 추론(predict): ep0 / ep80 / ep140 예측 궤적 생성", "자세오차 1.8° — 정상"],
    ["13:23~", "RViz 재생(replay): waypoint별 IK 체인 계획", "44/52 성공(84.6%) — wp34~39 실패"],
    ["오후", "특이점 트러블슈팅: 알고리즘 4종 실험 (A~D)", "작업대 x−10cm 이동 시 52/52(100%) 확정"],
    ["오후", "run_policy_tx90.py에 --offset 옵션 추가 → RViz 100% 재생 검증", "해결 완료"],
    ["오후", "TX90 외부 제어 조사 + 실행 방식/디코더 출력 아키텍처 논의", "실기 로드맵 확정"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[2.2, 9.5, 6.5], size=12.5, rh=Cm(1.55))
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 4. 트러블슈팅 #1: 오일러 wrap ════════════════════
s = slide()
bar(s, "트러블슈팅 ① — 오일러 각 wrap 결함", "TODAY")
tf = textbox(s, Cm(0.9), Cm(2.5), Cm(32), Cm(15.5))
bullets(tf, [
    (0, "증상", {"bold": True, "size": 16, "color": RED}),
    (1, "v4 데이터셋 학습 10k 시점 평가에서 자세(회전) 오차 28.9° — 위치는 정상인데 자세만 크게 어긋남"),
    (0, "원인", {"bold": True, "size": 16, "color": RED, "after": 6}),
    (1, "EE 자세를 오일러 각(extrinsic xyz)으로 저장할 때 ±180° 경계에서 wrap 발생"),
    (2, "예: rx가 +179° → −179°로 건너뜀 — 물리적으로는 2° 회전인데 수치로는 358° 점프"),
    (1, "회귀 모델은 수치를 그대로 학습 → 경계 근처 프레임에서 평균을 내다 엉뚱한 자세를 출력"),
    (1, "이 작업 자세들의 rx가 하필 ±180° 근처에 몰려 있어 피해가 컸음"),
    (0, "해결 (fix_euler_wrap.py)", {"bold": True, "size": 16, "color": GREEN, "after": 6}),
    (1, "에피소드별로 오일러 각 시퀀스를 시간순으로 unwrap — 프레임 간 차이가 ±180° 넘으면 2π 보정해 연속화"),
    (1, "수정본으로 v5 데이터셋(tx90_act_pick_and_place_v5_ee) 재조립 — 영상은 심볼릭 링크로 v4 재사용(용량 절약)"),
    (0, "결과", {"bold": True, "size": 16, "color": GREEN, "after": 6}),
    (1, "v5 재학습(100k) 후 자세오차 28.9° → 1.8° — 결함 해소 확인", {"bold": True}),
    (0, "교훈: 각도 표현을 학습 데이터로 쓸 때는 연속성(unwrap)을 반드시 전수 검사할 것", {"color": GREY}),
])
foot(s, "산출물: fix_euler_wrap.py · tx90_act_pick_and_place_v5_ee")

# ════════════════════ 5. v5 재학습 상세 ════════════════════
s = slide()
bar(s, "ACT v5 재학습 — 설정과 산출물", "TODAY")
rows = [
    ["항목", "값"],
    ["데이터셋", "dlcodnjs/tx90_act_pick_and_place_v5_ee (162 에피소드, 오일러 wrap 수정본)"],
    ["모델 / 스텝", "ACT · 100,000 스텝 · batch 8 (10:33~13:03, 약 2.5시간, RTX 5070 Ti)"],
    ["입력", "카메라 2대(480×640) + state 7D (EE pose 6D + gripper 1D)"],
    ["출력", "action 7D × chunk 100 스텝 (한 추론 = 3.3초 분량 궤적)"],
    ["체크포인트", "/root/train_tx90_act_v5/checkpoints/{010000..100000, last} — 10k마다 보존"],
    ["최종 모델", "checkpoints/last/pretrained_model (= 100000)"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[4.2, 13.8], size=12.5, rh=Cm(1.1))
tf = textbox(s, Cm(0.9), Cm(11.2), Cm(32), Cm(6.5))
bullets(tf, [
    (0, "검증 1단계 — open-loop 추론 (predict 모드)", {"bold": True, "size": 16}),
    (1, "데이터셋 에피소드의 관측(녹화 영상+pose)을 순서대로 넣어 예측 궤적 생성 → 정답과 비교"),
    (1, "ep0 / ep80 / ep140 세 에피소드 rollout 저장 (policy_rollouts/*.parquet)"),
    (1, "자세오차 1.8° — 학습 자체는 성공. 단 이것은 open-loop(관측이 녹화본)이므로 \"모방 정확도\" 검증임", {"bold": True}),
    (1, "실시간 관측 기반 제어(closed-loop)는 실기/시뮬레이션 단계의 별도 과제", {"color": GREY}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 6. RViz 재생 배선 + 채택 알고리즘 ════════════════════
s = slide()
bar(s, "검증 2단계 — RViz 재생과 IK 체인 알고리즘 (채택안)", "TODAY")
tf = textbox(s, Cm(0.9), Cm(2.4), Cm(19.5), Cm(15.8))
bullets(tf, [
    (0, "배선: run_policy_tx90.py replay 모드", {"bold": True, "size": 15}),
    (1, "예측 parquet → 다운샘플(1/10) → 정지 중복 제거 → 52 waypoint"),
    (1, "waypoint → IK → 관절 궤적 → 50Hz 보간 → fake_controller 토픽 발행"),
    (0, "1차 시도: MoveIt compute_cartesian_path → 기각", {"bold": True, "size": 15, "color": RED, "after": 6}),
    (1, "특이점 위 자세라 mm 차이로 IK 성패가 갈려 1.9%에서 조기 종료"),
    (0, "채택: waypoint별 IK 체인 (직접 구현)", {"bold": True, "size": 15, "color": GREEN, "after": 6}),
    (1, "① 직전 해를 시드로 waypoint마다 IK (GetPositionIK, 충돌 회피, 200ms)"),
    (1, "② 2π 접기: 관절 한계 안에서 시드에 가장 가까운 등가각으로 (+361°→+1°)"),
    (1, "③ 접은 뒤에도 max|Δ|>90°면 브랜치 플립(팔꿈치 뒤집힘)으로 보고 기각"),
    (1, "④ 실패 시 시드 섭동 4종으로 재시도, 그래도 실패면 해당 waypoint 건너뜀"),
    (0, "결과: 44/52 성공(84.6%) — 궤적 재생됨, 단 wp1,3,34~39 실패", {"bold": True, "size": 15, "after": 6}),
    (1, "7.7s 그리퍼 닫힘(집기) · 13.1s 열림(놓기) 재현 — 트러블슈팅 ②로", {"color": GREY}),
])
s.shapes.add_picture(os.path.join(IMG, "replay_t1.png"), Cm(21.0), Cm(2.6), height=Cm(12.5))
tf = textbox(s, Cm(21.0), Cm(15.3), Cm(11.5), Cm(1.2))
bullets(tf, [(0, "RViz 재생 화면 (TX2-90, base_link 고정)", {"color": GREY, "size": 11})])
foot(s, "부수 트러블슈팅: RViz에 map만 보이는 문제 = zenoh 라우터 부재 → rmw_zenohd 선기동으로 해결")

# ════════════════════ 7. 트러블슈팅 #2: 특이점 ════════════════════
s = slide()
bar(s, "트러블슈팅 ② — 손목 특이점 (wp34~39 실패)", "TODAY")
tf = textbox(s, Cm(0.9), Cm(2.5), Cm(32), Cm(15.5))
bullets(tf, [
    (0, "증상", {"bold": True, "size": 16, "color": RED}),
    (1, "놓기 직전 구간 wp34~39(약 2초)가 IK 재시도 4회에도 일관된 해 없음 → 건너뜀 (44/52)"),
    (1, "건너뛴 구간은 관절 보간으로 이어져 화면상 부드러워 보이지만 실제 EE 경로와 다름 — 정직하게 한계로 기록"),
    (0, "원인 (실측으로 확인)", {"bold": True, "size": 16, "color": RED, "after": 6}),
    (1, "이 작업 자세들의 TX90 관절해가 joint_5 ≈ 0° 위에 있음 — 4번·6번 축이 일직선이 되는 손목 특이점"),
    (1, "특이점 위에서는 같은 EE pose에 대한 관절해(브랜치)가 급변 → 연속 궤적이 성립 안 함"),
    (1, "wp34~39의 해는 존재하지만 현재 팔의 관절 배치에서 90° 미만 점프로는 도달 불가한 다른 브랜치에 있음"),
    (0, "해결 전략: 알고리즘 4종을 실험해 정량 비교 후 선택 (다음 장)", {"bold": True, "size": 16, "color": GREEN, "after": 6}),
    (1, "A. 브랜치 전환 허용   B. 작업 영역 평행이동   C. 시작 관절 배치 변경   D. IK 해 그래프 전역 탐색"),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 8. 알고리즘 4종 비교 (하이라이트) ════════════════════
s = slide()
bar(s, "특이점 해결 — 시도한 알고리즘 4종과 선택 이유", "TODAY")
rows = [
    ["안", "방법", "결과", "판정 / 이유"],
    ["A", "브랜치 플립 허용 — 일관된 해가 없으면 뒤집힌 해라도 수용",
     "52/52", ("기각 — 3곳에서 165~359° 관절 점프, 실기서 팔 전체가 휙 도는 동작", {"color": RED})],
    ["B", "작업 영역 평행이동 — 궤적 전체를 dx,dy,dz 이동 후 기존 체인 적용 (그리드 12종 탐색)",
     ("x−10cm: 52/52", {"bold": True}), ("채택 — 코드 한 줄(--offset), 어떤 궤적에도 견고, 실기는 테이프만 이동", {"color": GREEN, "bold": True})],
    ["C", "시작 관절 배치 변경 — wp34~39를 푸는 브랜치의 해를 시작 시드로 사용",
     "31/52", ("기각 — 그 브랜치는 중간 구간(wp14~30)을 못 풀음. 단일 브랜치로 전 구간 커버 불가 입증", {"color": RED})],
    ["D", "IK 해 그래프 탐색 — waypoint마다 다중 시드로 해 7~18개 수집(2π 등가 포함), 점프<90° 연결을 DP로 탐색(건너뛰기 허용)",
     "51/52", ("참고용 — 손목 뒤집힌 배치로 놓기 구간까지 통과. 단 에피소드 전용·최대 점프 83°·실시간 불가", {"color": GREY})],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[1.0, 8.2, 2.6, 7.8], size=11.5, rh=Cm(1.9))
tf = textbox(s, Cm(0.9), Cm(13.0), Cm(32), Cm(5.2))
bullets(tf, [
    (0, "실험 B 상세 — 평행이동 그리드 결과 (성공 waypoint / 52)", {"bold": True, "size": 15}),
    (1, "x−10cm: 52 ✅   x+10: 52   y+5: 52   (x−5,y+5): 52   |   y+10: 49   x−5: 50   z±5: 50~51   y−5: 34 ❌", {"mono": True, "size": 12.5}),
    (1, "y+5cm도 100%지만 y+10에서 다시 실패 → 여유 폭이 좁음. x−10cm은 대칭적으로 안전 → 채택"),
    (1, "의미: 현재 작업 영역이 특이점 능선 바로 위 — 로봇 쪽으로 10cm만 들어오면 능선을 완전히 벗어남"),
    (0, "실험 D의 부수 성과: wp10→11 DP 단절 원인 추적 → 2π 등가각 누락 버그 발견·수정, 전 구간 해 지도 확보", {"color": GREY, "size": 12.5}),
])
foot(s, "산출물: probe_singularity.py(A·B) · probe_branch.py(C) · probe_graph2.py+probe_dp2.py(D) · ik_solutions_ep000.npz")

# ════════════════════ 9. 최종 검증 ════════════════════
s = slide()
bar(s, "최종 검증 — 작업대 −10cm에서 52/52 (100%) 재생", "TODAY")
s.shapes.add_picture(os.path.join(IMG, "f04.png"), Cm(0.9), Cm(2.7), height=Cm(11.0))
s.shapes.add_picture(os.path.join(IMG, "f07.png"), Cm(12.2), Cm(2.7), height=Cm(11.0))
s.shapes.add_picture(os.path.join(IMG, "f10.png"), Cm(23.5), Cm(2.7), height=Cm(11.0))
tf = textbox(s, Cm(0.9), Cm(14.0), Cm(32), Cm(4.5))
bullets(tf, [
    (0, "run_policy_tx90.py replay --offset=-0.10,0,0 → IK 체인 52/52 waypoint 성공 (100%)", {"bold": True, "size": 15, "mono": True}),
    (1, "궤적 범위: x 0.47~0.68m · 어깨거리 0.48~0.70m (최대 리치 0.95m 내) · 18.7초 재생"),
    (1, "이전에 건너뛰던 놓기 직전 구간(wp34~39)까지 정책 예측 경로 그대로 재현 — 화면 전체가 이제 정책 출력과 일치"),
    (1, "실기 적용 시 주의: 실물 테이프 사각형(집기/놓기 위치)도 로봇 쪽으로 10cm 이동 필수", {"bold": True, "color": RED}),
])
foot(s, "캡처: 재생 중 3프레임 (접근 → 집기 하강 → 종료 자세)")

# ════════════════════ 10. 외부 제어 조사 ════════════════════
s = slide()
bar(s, "조사 — TX90(CS9) 외부 제어, 뚫을 수 있는가?", "SURVEY")
rows = [
    ["경로", "내용", "실시간성", "비용/조건"],
    ["① VAL3 드라이버", "ros_server VAL3 앱 + TCP 소켓 4개(11000~11003) → FollowJointTrajectory 궤적 실행 + 상태 피드백",
     "궤적 일괄 실행", ("무료", {"color": GREEN, "bold": True})],
    ["② velocity/motion 확장", "FAU-FAPS 미들웨어의 속도·포즈 추종 모드 — 보간 클럭(~4ms) 주기 연속 제어",
     ("연속 스트리밍", {"bold": True}), "CS9 확장 옵션 (유상 추정, 견적 문의 예정)"],
    ["③ uniVAL Drive", "CS9를 EtherCAT/PROFINET 서보 드라이브처럼 구동 — 필드버스 사이클 위치 명령",
     "하드 실시간", ("유료 + PLC 스택 — 연구용 과함", {"color": GREY})],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[3.4, 9.6, 3.2, 4.4], size=11.5, rh=Cm(1.75))
tf = textbox(s, Cm(0.9), Cm(10.4), Cm(32), Cm(7.8))
bullets(tf, [
    (0, "핵심 판단", {"bold": True, "size": 15}),
    (1, "검증 단계(ACT open-loop 재현)와 스크립트 데이터 수집은 ①만으로 가능 — 리스크 낮음"),
    (1, "반응형(30Hz closed-loop)은 ②가 관건 — Stäubli Robotics Suite의 CS9 에뮬레이터로 실기 없이 선검증 가능"),
    (1, "온라인 스트리밍을 ①로 흉내내는 시도는 공식 이슈에서 불안정 판정 → ② 없이는 재계획형까지만"),
    (0, "참고 링크", {"bold": True, "size": 15, "after": 6}),
    (1, "ros-industrial/staubli_val3_driver (CS8/CS9 공식 드라이버)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver"}),
    (1, "IvoD1998/Staubli_ROS2 (ROS2 포팅 + 미들웨어 포함)",
     {"link": "https://github.com/IvoD1998/Staubli_ROS2"}),
    (1, "FAU-FAPS/adaptive_motion_control (velocity/pose tracking 미들웨어)",
     {"link": "https://github.com/FAU-FAPS/adaptive_motion_control"}),
    (1, "TX2-90 실시간 서보 논의 (staubli_val3_driver issue #32)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver/issues/32"}),
    (1, "Stäubli uniVAL Drive 공식 페이지",
     {"link": "https://www.staubli.com/us/en/robotics/products/robot-software/uniVAL-drive.html"}),
])
foot(s, "우리 로봇: TX2-90 → CS9 컨트롤러 세대 (staubli_tx2_90_description 확인)")

# ════════════════════ 11. 아키텍처 결정 ① 실행 방식 ════════════════════
s = slide()
bar(s, "아키텍처 정리 ① — 실행 방식 스펙트럼", "DESIGN")
y1 = Cm(3.0); h = Cm(3.4); w = Cm(9.6); gap = Cm(1.2)
boxes = [
    ("open-loop (현재)\n관측 없이 좌표 재생\n= 오늘 RViz 검증", GREEN),
    ("재계획형 (실기 1차 목표)\n청크마다 관측→계획→실행\n기본 드라이버로 가능", NAVY),
    ("반응형 (최종 목표)\n30Hz 연속 관측·수정\nvelocity 확장 필요", GREY),
]
x = Cm(0.9)
for i, (txt, c) in enumerate(boxes):
    flowbox(s, x, y1, w, h, txt, fill=c, size=14)
    if i < len(boxes) - 1:
        arrow(s, x + w + Cm(0.2), y1 + h / 2 - Cm(0.25), Cm(0.75))
    x += w + gap
tf = textbox(s, Cm(0.9), Cm(7.2), Cm(32), Cm(10.8))
bullets(tf, [
    (0, "재계획형도 closed-loop이며 VLA로 성립 — 관측이 청크마다 행동을 결정 (RT-2 등도 1~3Hz 저주파 실행)", {"bold": True}),
    (1, "주사위를 옮기면: 실행 중인 청크는 못 보지만 다음 청크(1~3초 내)가 새 위치로 향함"),
    (1, "짧은 청크(1초) + 집기 직전 확인 관측 가드레일로 반응 지연 최소화 가능"),
    (1, "연속 이동 추적만 반응형 전용 영역 — 이 경계가 velocity 확장의 가치를 정함"),
    (0, "플랜: 실기 상한선은 CS9 에뮬레이터+견적으로 판가름, 반응형 연구 본론은 Isaac Sim에서 선행", {"bold": True, "after": 6}),
    (1, "경로 ② 실패 시에도: 실기=재계획형 완성 + 반응형=sim 정량 검증 → \"산업 컨트롤러 제약별 VLA 배포 스펙트럼\" 연구 프레임"),
    (0, "용어 주의: 현 ACT는 언어 입력이 없어 visuomotor policy. VLA 타이틀은 SmolVLA(언어 지시) 단계부터", {"color": GREY}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 12. 아키텍처 결정 ② 디코더 출력 ════════════════════
s = slide()
bar(s, "아키텍처 정리 ② — 디코더 출력: EE pose vs joint", "DESIGN")
rows = [
    ["", "B안: EE pose 유지 + 청크 IK 레이어", "C안: joint 직접 출력 (v6)"],
    ["구조", "디코더(EE 청크 100) → IK 체인 레이어(시드=직전 명령) → joint 버퍼 → 30Hz 스트림",
     "디코더가 joint 청크를 직접 출력 → 안전 필터 → 로봇"],
    ["재학습", ("불필요 — v5 그대로", {"color": GREEN}), "필요 (~2.5h) — v6 데이터셋으로"],
    ["런타임 IK", "루프 밖(청크 단위, 3.3초 여유) — 실패는 실행 전 감지", ("없음 — 특이점 문제 학습 전에 소멸", {"color": GREEN})],
    ["joint 라벨", "—", "v5 EE action을 오프라인 IK 체인으로 일괄 변환 (실기 시연 0개)"],
    ["OMX 전이성", ("유지 (EE 규격 공통)", {"color": GREEN}), "TX90 전용으로 고정"],
    ["연속성 보장", "체인 알고리즘이 보장", "학습으로 획득 + 안전 필터(|Δj| 한계) 검사"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[2.6, 7.9, 7.6], size=11.5, rh=Cm(1.6))
tf = textbox(s, Cm(0.9), Cm(14.2), Cm(32), Cm(4.0))
bullets(tf, [
    (0, "결정: 양자택일하지 않고 둘 다 확보해 비교 (v5는 이미 있고 v6는 변환+재학습만)", {"bold": True, "size": 15}),
    (1, "v6 변환 규약: 전 에피소드 공통 HOME 시드 · −10cm 오프셋 위치에서 변환 · 90° 점프 거부/2π 접기 동일 적용"),
    (1, "v6 검증 3종: 프레임 간 최대 점프 · 에피소드 간 브랜치 분포(단봉성) · FK 역검증(mm 오차)"),
    (1, "MoveIt 체인의 역할 이관: 브랜치 선택→변환 때 소멸 · 시드 연속성→state 조건부 예측 · 점프 보장→출력 안전 필터"),
])
foot(s, "박사님 질문 회신: \"디코더에서 joint 출력 가능 — IK는 런타임이 아닌 데이터셋 변환 시 1회\"")

# ════════════════════ 12a. ACT 모델 해부 (입력·중간·출력) ════════════════════
s = slide()
bar(s, "ACT 모델 해부 — 입력 · 중간 과정 · 출력 (v5 실측 설정)", "MODEL")
# 입력 열
flowbox(s, Cm(0.9), Cm(2.6), Cm(5.4), Cm(1.5), "camera1\n480×640 RGB", fill=BGSOFT, fg=INK, size=12)
flowbox(s, Cm(0.9), Cm(4.3), Cm(5.4), Cm(1.5), "camera2\n480×640 RGB", fill=BGSOFT, fg=INK, size=12)
flowbox(s, Cm(0.9), Cm(6.0), Cm(5.4), Cm(1.5), "state 7D\nEE pose 6 + grip 1", fill=BGSOFT, fg=INK, size=12)
flowbox(s, Cm(0.9), Cm(7.7), Cm(5.4), Cm(1.5), "latent z 32D\n(학습 때만, 추론 시 0)", fill=GREY, size=11)
# 토큰화 열
flowbox(s, Cm(7.6), Cm(3.0), Cm(6.4), Cm(2.6),
        "ResNet18 ×2\n(ImageNet 사전학습)\n→ 카메라당 300 토큰", size=12)
flowbox(s, Cm(7.6), Cm(6.0), Cm(6.4), Cm(1.5), "선형 투영 → 1 토큰", size=12)
flowbox(s, Cm(7.6), Cm(7.7), Cm(6.4), Cm(1.5), "투영 → 1 토큰", fill=GREY, size=12)
for y in (Cm(3.3), Cm(5.0), Cm(6.5), Cm(8.2)):
    arrow(s, Cm(6.5), y, Cm(0.9))
# 인코더/디코더/출력
flowbox(s, Cm(15.2), Cm(3.4), Cm(5.6), Cm(4.8),
        "Transformer 인코더\n4층 · 512D · 8헤드\n\n602 토큰 융합\n(화면·자세·스타일 연결)", size=13)
flowbox(s, Cm(22.0), Cm(3.4), Cm(5.2), Cm(4.8),
        "디코더 1층\n\n학습된 쿼리 100개\n= 청크의 각 스텝\ncross-attention", size=13)
flowbox(s, Cm(28.4), Cm(3.4), Cm(4.6), Cm(4.8),
        "action 청크\n100스텝 × 7D\n\n= 3.3초 궤적\n(30Hz)", fill=GREEN, size=13)
arrow(s, Cm(14.2), Cm(5.5), Cm(0.85))
arrow(s, Cm(20.95), Cm(5.5), Cm(0.9))
arrow(s, Cm(27.35), Cm(5.5), Cm(0.9))
tf = textbox(s, Cm(0.9), Cm(9.8), Cm(32), Cm(8.2))
bullets(tf, [
    (0, "중간 과정 핵심", {"bold": True, "size": 15}),
    (1, "영상: ResNet18이 480×640 → 15×20 특징맵으로 압축, 격자 한 칸 = 토큰 1개 (2D 위치 인코딩으로 \"화면 어디\"를 보존)"),
    (1, "latent z (CVAE): 학습 때는 별도 VAE 인코더(4층)가 정답 청크에서 \"시연 스타일\"을 뽑아 z에 흡수 (KL 가중치 10)"
        " — 시연 간 다양성이 평균으로 뭉개지는 것을 방지. 추론 때는 z=0(평균 스타일), VAE 인코더는 돌지 않음"),
    (1, "손실: 정답 청크와의 L1 + KL — \"다음 한 스텝\"이 아니라 청크 통째 회귀 → 스텝별 오차 누적(compounding error)을 차단"),
    (0, "현재 소비 방식과 조정 노브", {"bold": True, "size": 15, "after": 6}),
    (1, "n_action_steps=100: 청크를 끝까지 실행 후 재추론 (3.3초마다 재계획) · temporal ensembling은 현재 꺼짐(None)"),
    (1, "반응형으로 갈 때: 청크 단축(예: 30스텝=1초) + 앙상블 켜기 → 관측 반영 주기 상승 — 모델 재학습 없이 실행기 설정만 변경"),
])
foot(s, "모델 해부·방식 비교 · 2026-08-14 추가")

# ════════════════════ 12b. B안 vs C안 실행 파이프라인 비교 ════════════════════
s = slide()
bar(s, "B안 vs C안 — \"번역(IK)을 언제 하느냐\"의 차이", "MODEL")
tf = textbox(s, Cm(0.9), Cm(2.25), Cm(32), Cm(0.8))
bullets(tf, [(0, "B안 — 실행할 때마다 번역 (모델 v5 그대로, 재학습 불필요)", {"bold": True, "size": 14, "color": NAVY})])
bx = [
    (Cm(0.9),  Cm(4.7), "카메라 2대\n+ state (EE pose)", BGSOFT, INK),
    (Cm(6.3),  Cm(5.2), "모델 v5\nEE 청크 100스텝", NAVY, WHITE),
    (Cm(12.2), Cm(6.6), "IK 번역기 — 매 청크 가동\n체인 방식 · 0.1~0.5초\n실패 가능(낮지만 0 아님)", RED, WHITE),
    (Cm(19.5), Cm(4.4), "joint 버퍼", BGSOFT, INK),
    (Cm(24.6), Cm(4.2), "로봇 30Hz", GREEN, WHITE),
]
for x, w, txt, c, fg in bx:
    flowbox(s, x, Cm(3.1), w, Cm(2.0), txt, fill=c, fg=fg, size=11.5)
for x in (Cm(5.65), Cm(11.55), Cm(18.85), Cm(23.95)):
    arrow(s, x, Cm(3.85), Cm(0.6))
tf = textbox(s, Cm(0.9), Cm(5.5), Cm(32), Cm(0.8))
bullets(tf, [(0, "C안 — 미리 번역해두고 학습 (모델 v6 새로, 실행 시 IK 없음)", {"bold": True, "size": 14, "color": GREEN})])
flowbox(s, Cm(0.9), Cm(6.35), Cm(31.8), Cm(1.15),
        "사전 작업 딱 1회:  데이터셋 v5(EE 표현) ─ 오프라인 IK 일괄 번역 (−10cm 위치, 공통 HOME 시드) ─▶ 데이터셋 v6(joint 표현) ─▶ 재학습 ~2.5h",
        fill=BGSOFT, fg=GREY, size=11.5)
cx = [
    (Cm(0.9),  Cm(5.6), "카메라 2대\n+ state (joint)", BGSOFT, INK),
    (Cm(7.2),  Cm(6.2), "모델 v6\njoint 청크 100스텝", NAVY, WHITE),
    (Cm(14.2), Cm(7.6), "안전 필터\n|Δjoint| 한계 검사 (뺄셈+비교, 비용≈0)", GREEN, WHITE),
    (Cm(22.6), Cm(4.2), "로봇 30Hz", GREEN, WHITE),
]
for x, w, txt, c, fg in cx:
    flowbox(s, x, Cm(7.85), w, Cm(2.0), txt, fill=c, fg=fg, size=11.5)
for x in (Cm(6.55), Cm(13.45), Cm(21.85)):
    arrow(s, x, Cm(8.6), Cm(0.6))
rows = [
    ["상황", "B안 (EE + 번역기)", "C안 (joint 직접)"],
    ["실행 중 IK 실패", "청크 폐기·재시도 로직 필요", ("그런 상황 자체가 없음", {"color": GREEN})],
    ["특이점", "실행 때마다 회피 (−10cm 전제)", "변환 때 1회 회피로 종결"],
    ["다른 로봇 재사용 / OMX 호환", ("EE는 로봇 중립 — 전이 노선 유지", {"color": GREEN}), "TX90 전용으로 고정"],
    ["오작동 원인 진단", "모델 탓/번역기 탓 분리 가능", "원인이 모델 하나로 좁혀짐"],
    ["재학습", ("불필요 — 지금 바로 사용 가능", {"color": GREEN}), "필요 (~2.5h) + v6 변환 선행"],
]
table(s, Cm(0.9), Cm(10.6), Cm(32.1), rows, widths=[4.6, 6.7, 6.7], size=11.5, rh=Cm(1.15))
foot(s, "비유: B = 실행마다 통역사 대동 / C = 교재를 미리 번역해 원어로 학습 — 결정: 둘 다 확보해 비교")

# ════════════════════ 12c. v5·v6 관계와 연속성 ════════════════════
s = slide()
bar(s, "오해 방지 — v6는 \"모델 v5\"가 아니라 \"원본 시연\"에서 나온다", "MODEL")
flowbox(s, Cm(0.9), Cm(2.7), Cm(7.8), Cm(2.0),
        "OMX 사람 시연 162ep\n(원본 정답 — 좌표 변환 완료)", fill=NAVY, size=12.5)
arrow(s, Cm(8.85), Cm(3.45), Cm(0.7))
flowbox(s, Cm(9.7), Cm(2.7), Cm(6.8), Cm(2.0),
        "데이터셋 v5\naction = EE pose 표현", fill=BGSOFT, fg=INK, size=12.5)
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(16.65), Cm(3.2), Cm(3.3), Cm(1.0))
a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background()
tfa = a.text_frame; p = tfa.paragraphs[0]
run = p.add_run(); run.text = "오프라인 IK 번역 (1회)"
_set(run, 9.5, True, WHITE)
flowbox(s, Cm(20.2), Cm(2.7), Cm(6.8), Cm(2.0),
        "데이터셋 v6\naction = joint 표현", fill=BGSOFT, fg=INK, size=12.5)
for x in (Cm(12.6), Cm(23.1)):
    d = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, Cm(4.85), Cm(0.7), Cm(0.9))
    d.fill.solid(); d.fill.fore_color.rgb = GREY; d.line.fill.background()
flowbox(s, Cm(9.7), Cm(5.9), Cm(6.8), Cm(1.7), "모델 v5 (EE 출력)\n— 학습 완료", fill=GREEN, size=12.5)
flowbox(s, Cm(20.2), Cm(5.9), Cm(6.8), Cm(1.7), "모델 v6 (joint 출력)\n— 새로 학습 예정", fill=GREEN, size=12.5)
tf = textbox(s, Cm(0.9), Cm(8.0), Cm(32), Cm(2.2))
bullets(tf, [
    (0, "모델 v5는 v6 생성 과정 어디에도 관여하지 않음 — 두 모델은 같은 시연을 다른 언어(EE/joint)로 배운 \"형제\"", {"bold": True, "size": 14}),
    (1, "만약 모델 v5의 예측으로 v6를 만들면 예측 오차(1.8°)까지 정답으로 배우는 \"복사본의 복사본\" — 그렇게 하지 않는 이유", {"color": GREY, "size": 12.5}),
])
rows = [
    ["연속성(안 튐)은 누가 보장하나", "v5 + IK 체인 (현행)", "v6 joint 출력 (예정)"],
    ["보장 주체", "실행 중 알고리즘이 강제 (직전 해 시드 + 90° 점프 거부)", "데이터 규약 + 학습 + 안전 필터로 분산"],
    ["청크 내부", "waypoint마다 독립적으로 풀어 이어붙임 → 브랜치 튈 여지", "100스텝을 한 번의 계산으로 출력 → 구조적으로 매끄러움"],
    ["청크 경계", "직전 관절해를 다음 시드로", "state(현재 joint) 조건부 예측이 시드 역할 대체"],
    ["실패 모드", "해 없음 → waypoint 건너뜀 (궤적에 구멍, 44/52 사례)", "분포 밖 관측 → 어긋난 값 가능 → |Δj| 필터가 차단"],
    ["검증 방법", "재생 성공률 (오늘 100% 확인)", "학습 후 전 에피소드 프레임 간 최대 점프 실측 (검증 3종)"],
]
table(s, Cm(0.9), Cm(10.7), Cm(32.1), rows, widths=[4.4, 6.9, 6.9], size=11.5, rh=Cm(1.25))
foot(s, "모델 해부·방식 비교 · 2026-08-14 추가")

# ════════════════════ 13. 생성 파일 ① 호스트 ════════════════════
s = slide()
bar(s, "오늘 생성·수정 파일 ① — 호스트 (~/tx90)", "FILES")
rows = [
    ["파일", "설명"],
    ["fix_euler_wrap.py", "v4 데이터셋의 오일러 각 ±180° wrap을 에피소드별로 unwrap → v5 데이터셋 재조립 스크립트"],
    ["run_policy_tx90.py", "STEP 8 배선 (오늘 완성+수정). predict: 체크포인트→예측 parquet / replay: IK 체인→RViz 재생. 오늘 --offset(평행이동) 옵션 추가"],
    ["대화록 및 PPT/img_0813/*.png", "RViz 검증 캡처 4장 (재생 전 1 + 재생 중 3프레임)"],
    ["make_report_pptx_0813.py", "이 PPT의 생성 스크립트 (수정 후 재실행하면 갱신)"],
    ["~/physical_ai_tools/.../dlcodnjs/tx90_act_pick_and_place_v5_ee", "v5 데이터셋 (162 에피소드, wrap 수정본. 영상은 v4로 심볼릭 링크)"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[6.8, 11.2], size=12, rh=Cm(1.75))
tf = textbox(s, Cm(0.9), Cm(13.6), Cm(32), Cm(4.2))
bullets(tf, [
    (0, "재현 명령 (컨테이너 안)", {"bold": True, "size": 14}),
    (1, "python3 /root/run_policy_tx90.py replay --parquet /root/policy_rollouts/ep000_pred_v5_100k.parquet --offset=-0.10,0,0",
     {"mono": True, "size": 11.5}),
    (1, "전제: demo.launch.py(RViz) + zenoh 라우터 실행 중. --offset 없이 돌리면 기존 44/52 재현", {"color": GREY}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 14. 생성 파일 ② 컨테이너 ════════════════════
s = slide()
bar(s, "오늘 생성 파일 ② — 컨테이너 physical_ai_server (/root)", "FILES")
rows = [
    ["파일", "설명"],
    ["train_tx90_act_v5/checkpoints/", "v5 100k 학습 산출물. last→100000이 최종 모델 (pretrained_model)"],
    ["train_tx90_act/ (10k 중단)", "wrap 수정 전(v4) 학습분 — 결함 진단용으로만 의미, 사용 안 함"],
    ["policy_rollouts/ep{000,80,140}_pred_v5*.parquet", "open-loop 예측 궤적 3개 (RViz 재생 입력)"],
    ["probe_singularity.py", "실험 A(브랜치 플립 허용)·B(평행이동 그리드 12종) 스크립트"],
    ["probe_branch.py", "실험 C — 좋은 브랜치 시드로 strict 체인 재시작 (31/52 → 기각 근거)"],
    ["probe_graph.py / probe_graph2.py", "실험 D — waypoint별 IK 해 다중 수집 + 연속 경로 DP (v2에서 2π 등가 버그 수정)"],
    ["probe_dp2.py / probe_debug.py", "건너뛰기 허용 DP (51/52 도출) / wp10→11 DP 단절 원인 디버그"],
    ["ik_solutions_ep000.npz", "ep0 52 waypoint별 IK 해 집합 (해 7~18개/wp, 2π 등가 포함) — 재분석용 캐시"],
    ["branch_path_ep000.npy (+kept)", "실험 D가 찾은 51/52 연속 관절 궤적 (손목 뒤집힌 배치) — 데모·비교용"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[7.2, 10.8], size=11.5, rh=Cm(1.42))
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 15. 진행상황 맵 ════════════════════
s = slide()
bar(s, "진행상황 — 어디까지 왔나", "STATUS")
rows = [
    ["단계", "상태", "비고"],
    ["OMX 시연 수집 (162ep) · FK · 좌표 변환", ("완료", {"color": GREEN, "bold": True}), "~8/12"],
    ["TX90 EE 데이터셋 (v4→v5 wrap 수정)", ("완료", {"color": GREEN, "bold": True}), "8/13 오전"],
    ["ACT 학습 (100k) + open-loop 검증 (1.8°)", ("완료", {"color": GREEN, "bold": True}), "8/13"],
    ["RViz 기구학 검증 + 특이점 해결 (−10cm, 100%)", ("완료", {"color": GREEN, "bold": True}), "8/13 — 작업대 이동 결정"],
    ["외부 제어 경로 조사", ("완료", {"color": GREEN, "bold": True}), "8/13 — VAL3 무료 경로 확인"],
    ["CS9 에뮬레이터 검증 + velocity 확장 견적", ("다음", {"color": RED, "bold": True}), "내일부터 — 반응형 가능 여부 판가름"],
    ["v6 joint 데이터셋 변환 + 재학습", ("다음", {"color": RED, "bold": True}), "내일 — Claude가 배치 변환"],
    ["Isaac Sim 물리 검증 (v6 궤적 재생)", ("이번 주", {"color": RED}), "그리퍼가 실제로 집는지"],
    ["실기 open-loop 재현 → 스크립트 수집 → 전이학습", ("예정", {"color": GREY}), "외부 제어 확보 후"],
    ["SmolVLA (언어 지시 VLA) + 재계획형/반응형 비교", ("예정", {"color": GREY}), "ACT 실기 성공 후"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[9.5, 2.2, 6.3], size=12, rh=Cm(1.35))
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 16. 다음 단계 (내일) ════════════════════
s = slide()
bar(s, "다음 단계 — 내일 할 일", "NEXT")
rows = [
    ["순서", "작업", "목적"],
    ["오전 1", "Stäubli 계정 생성 → 소프트웨어 센터에서 motion/velocity 확장 확인, SRS 다운로드(Windows)", "경로② 실존 확인"],
    ["오전 2", "대리점 견적 메일 — \"TX2-90/CS9 motion·velocity expansion 설치 비용·절차\" (시리얼 첨부)", "리드타임 먼저 확보"],
    ["오후", "SRS 설치 → CS9 에뮬레이터 기동 → TCP 소켓 4개(11000~11003) 설정 연습", "실기 없는 검증 환경"],
    ["Claude", "v6 joint 데이터셋 변환 (배치 IK 체인 + 검증 리포트 3종) → 통과 시 재학습", "\"v6 변환 시작해줘\""],
    ["병행", "랩미팅 PPT (이 자료 기반, 이번 주 전체로 확장)", "구성안 확정됨"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[1.8, 11.2, 5.0], size=12, rh=Cm(1.6))
tf = textbox(s, Cm(0.9), Cm(11.8), Cm(32), Cm(6.2))
bullets(tf, [
    (0, "이번 주 Isaac Sim 목표 — \"RViz 재생을 물리로\"", {"bold": True, "size": 15}),
    (1, "① 설치 + tx2_90.urdf 임포트  ② 테이블·주사위·카메라 2대 장면 구성 (−10cm 반영)  ③ v6 joint 궤적 물리 재생"),
    (1, "성공 기준: 그리퍼가 주사위를 실제로 집고 옮기는가 (RViz는 물리 없음 — 여기가 처음)"),
    (0, "판단 포인트: 에뮬레이터+견적 결과가 나오면 실행 방식 상한선(재계획형 vs 반응형) 자동 확정", {"bold": True, "size": 15}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 17. 고려사항 ════════════════════
s = slide()
bar(s, "고려사항 — 앞으로 밟을 수 있는 지뢰", "RISK")
rows = [
    ["리스크", "내용", "대응"],
    ["시각 도메인 갭", "정책이 본 영상은 OMX 장면 — TX90 카메라를 달아도 인식 불가 (팔·배경·그리퍼 다름)",
     "실기 closed-loop 전에 TX90 데이터 수집 or Isaac Sim 렌더 합성 + 전이학습"],
    ["v6 브랜치 일관성", "에피소드마다 IK 브랜치가 다르면 학습 데이터가 다봉(multimodal) → 평균내다 엉터리 관절값",
     "공통 HOME 시드 + 오프셋 위치 변환 + 검증 3종(점프/분포/FK)"],
    ["CS9 스트리밍 미지수", "velocity 확장이 비싸거나 불가할 수 있음",
     "플랜B 확정: 실기=재계획형(짧은 청크+확인 관측), 반응형=sim 정량 검증"],
    ["반응형 지연 보상", "비동기 추론 시 새 청크는 반 박자 늦은 관측 기준 — 경계 블렌딩 필요",
     "Isaac Sim에서 지연 주입 실험으로 선행 해결"],
    ["실기 안전", "학습 정책은 분포 밖 상황에서 보장 없음",
     "속도 제한 + |Δjoint| 안전 필터 + 작업 영역 소프트 리밋 + CS9 자체 한계 이중화"],
    ["오프셋 정합", "−10cm는 데이터·실행·실물 배치가 전부 함께 움직여야 함",
     "실물 테이프도 10cm 이동, v6에는 오프셋을 데이터에 구움(일관성)"],
]
table(s, Cm(0.9), Cm(2.6), Cm(32.1), rows, widths=[3.6, 8.0, 6.4], size=11.5, rh=Cm(1.95))
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

# ════════════════════ 18. 참고 링크 모음 ════════════════════
s = slide()
bar(s, "참고 링크 모음", "REF")
tf = textbox(s, Cm(0.9), Cm(2.6), Cm(32), Cm(15))
bullets(tf, [
    (0, "TX90 외부 제어", {"bold": True, "size": 16}),
    (1, "staubli_val3_driver — CS8/CS9 공식 ROS 드라이버 (무료, FollowJointTrajectory)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver"}),
    (1, "Staubli_ROS2 — ROS2 포팅 + adaptive motion middleware 포함",
     {"link": "https://github.com/IvoD1998/Staubli_ROS2"}),
    (1, "adaptive_motion_control — velocity/pose tracking 미들웨어 (FAU-FAPS)",
     {"link": "https://github.com/FAU-FAPS/adaptive_motion_control"}),
    (1, "issue #32 — TX2-90 실시간 서보 가능성 논의 (개발자 답변)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver/issues/32"}),
    (1, "issue #20 — 온라인 관절 제어 시도 (기본 드라이버 스트리밍의 한계)",
     {"link": "https://github.com/ros-industrial/staubli_val3_driver/issues/20"}),
    (1, "Stäubli CS9 컨트롤러 공식 페이지",
     {"link": "https://www.staubli.com/global/en/robotics/products/robot-controllers/cs9-robot-controller.html"}),
    (1, "Stäubli uniVAL Drive — 필드버스 저수준 제어 (경로 ③)",
     {"link": "https://www.staubli.com/us/en/robotics/products/robot-software/uniVAL-drive.html"}),
    (0, "학습 프레임워크·시뮬레이션", {"bold": True, "size": 16, "after": 6}),
    (1, "LeRobot (ACT·SmolVLA 학습/비동기 추론 스택)",
     {"link": "https://github.com/huggingface/lerobot"}),
    (1, "NVIDIA Isaac Sim (물리 검증 + 데이터 합성)",
     {"link": "https://developer.nvidia.com/isaac-sim"}),
    (1, "ROS-Industrial staubli_experimental (현재 사용 중인 TX2-90 description/moveit config 출처)",
     {"link": "https://github.com/ros-industrial/staubli_experimental"}),
])
foot(s, "OMX→TX2-90 전이 · 2026-08-13")

prs.save(os.path.join(OUTDIR, "20260813.pptx"))
print("저장:", os.path.join(OUTDIR, "20260813.pptx"), f"({len(prs.slides.__iter__.__self__._sldIdLst)}슬라이드)")
